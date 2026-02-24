from __future__ import annotations

from pathlib import Path

import pytest

try:
    from PIL import Image
    from pptx import Presentation
except ModuleNotFoundError:  # pragma: no cover - optional dependency in some environments
    pytest.skip(
        "python-pptx required for PNG validation tests",
        allow_module_level=True,
    )

from counter_risk.ppt.pptx_static import PngValidationError, _rebuild_pptx_from_slide_images


def _make_source_pptx(path: Path) -> None:
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(path))


def _write_valid_png(path: Path) -> None:
    Image.new("RGB", (1, 1), color=(255, 255, 255)).save(path, format="PNG")


def test_rebuild_pptx_from_slide_images_fails_for_missing_png_before_write(tmp_path: Path) -> None:
    source = tmp_path / "source.pptx"
    output = tmp_path / "distribution.pptx"
    missing_png = tmp_path / "missing.png"
    _make_source_pptx(source)

    with pytest.raises(PngValidationError) as exc_info:
        _rebuild_pptx_from_slide_images(
            source_pptx=source,
            slide_images=[missing_png],
            output_path=output,
        )

    assert "PNG validation failed" in str(exc_info.value)
    assert str(missing_png) in str(exc_info.value)
    assert not output.exists()


def test_rebuild_pptx_from_slide_images_fails_for_empty_png_before_write(tmp_path: Path) -> None:
    source = tmp_path / "source.pptx"
    output = tmp_path / "distribution.pptx"
    empty_png = tmp_path / "empty.png"
    _make_source_pptx(source)
    empty_png.write_bytes(b"")

    with pytest.raises(PngValidationError) as exc_info:
        _rebuild_pptx_from_slide_images(
            source_pptx=source,
            slide_images=[empty_png],
            output_path=output,
        )

    assert "PNG validation failed" in str(exc_info.value)
    assert str(empty_png) in str(exc_info.value)
    assert not output.exists()


def test_rebuild_pptx_from_slide_images_fails_for_corrupted_png_before_write(
    tmp_path: Path,
) -> None:
    source = tmp_path / "source.pptx"
    output = tmp_path / "distribution.pptx"
    corrupt_png = tmp_path / "corrupt.png"
    _make_source_pptx(source)
    corrupt_png.write_bytes(b"not-a-real-png")

    with pytest.raises(PngValidationError) as exc_info:
        _rebuild_pptx_from_slide_images(
            source_pptx=source,
            slide_images=[corrupt_png],
            output_path=output,
        )

    assert "PNG validation failed" in str(exc_info.value)
    assert str(corrupt_png) in str(exc_info.value)
    assert not output.exists()


def test_rebuild_pptx_from_slide_images_accepts_valid_png(tmp_path: Path) -> None:
    source = tmp_path / "source.pptx"
    output = tmp_path / "distribution.pptx"
    valid_png = tmp_path / "valid.png"
    _make_source_pptx(source)
    _write_valid_png(valid_png)

    _rebuild_pptx_from_slide_images(
        source_pptx=source,
        slide_images=[valid_png],
        output_path=output,
    )

    assert output.exists()
