from __future__ import annotations

import base64
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from counter_risk.writers.pptx_screenshots import replace_screenshot_pictures

_RED_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO6n8e0AAAAASUVORK5CYII="
)
_BLUE_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAwMBAX0X6XQAAAAASUVORK5CYII="
)
_GREEN_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP4DwQACfsD/QX+vJ8AAAAASUVORK5CYII="
)

_TARGET_PICTURE_GEOMETRY = {
    0: (0, 811705, 9144000, 5817695),
    5: (0, 1304636, 9144000, 5172364),
    15: (0, 2031298, 9144000, 2795403),
}


def _write_png(path: Path, payload: bytes) -> None:
    path.write_bytes(payload)


def _make_sample_pptx(path: Path, base_image: Path) -> None:
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    for idx in range(23):
        slide = prs.slides.add_slide(blank_layout)
        title = f"Section {idx + 1}"
        if idx == 0:
            title = "All Programs"
        elif idx == 5:
            title = "Ex Trend"
        elif idx == 15:
            title = "Trend"

        textbox = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.25),
            Inches(8.5),
            Inches(0.6),
        )
        textbox.text_frame.text = title

        if idx in _TARGET_PICTURE_GEOMETRY:
            left, top, width, height = _TARGET_PICTURE_GEOMETRY[idx]
            slide.shapes.add_picture(
                str(base_image),
                left,
                top,
                width=width,
                height=height,
            )

    prs.save(str(path))


def test_replace_screenshot_pictures_replaces_without_appending(tmp_path: Path) -> None:
    source = tmp_path / "input.pptx"
    output = tmp_path / "output.pptx"

    base_image = tmp_path / "base.png"
    all_programs_image = tmp_path / "all_programs.png"
    ex_trend_image = tmp_path / "ex_trend.png"
    trend_image = tmp_path / "trend.png"

    _write_png(base_image, _RED_PNG)
    _write_png(all_programs_image, _BLUE_PNG)
    _write_png(ex_trend_image, _GREEN_PNG)
    _write_png(trend_image, _BLUE_PNG)

    _make_sample_pptx(source, base_image)

    before = Presentation(str(source))
    assert len(before.slides) == 23
    before_counts = [sum(1 for s in slide.shapes if s.shape_type == 13) for slide in before.slides]

    replace_screenshot_pictures(
        source,
        {
            "All Programs": all_programs_image,
            "Ex Trend": ex_trend_image,
            "Trend": trend_image,
        },
        output,
    )

    after = Presentation(str(output))
    assert len(after.slides) == 23

    after_counts = [sum(1 for s in slide.shapes if s.shape_type == 13) for slide in after.slides]
    assert before_counts == after_counts

    target_map = {
        0: all_programs_image.read_bytes(),
        5: ex_trend_image.read_bytes(),
        15: trend_image.read_bytes(),
    }

    for index, expected_blob in target_map.items():
        pictures = [s for s in after.slides[index].shapes if s.shape_type == 13]
        assert len(pictures) == 1
        assert pictures[0].image.blob == expected_blob


def test_replace_screenshot_pictures_raises_when_section_not_found(tmp_path: Path) -> None:
    source = tmp_path / "input.pptx"
    output = tmp_path / "output.pptx"
    base_image = tmp_path / "base.png"
    replacement = tmp_path / "replacement.png"

    _write_png(base_image, _RED_PNG)
    _write_png(replacement, _BLUE_PNG)
    _make_sample_pptx(source, base_image)

    with pytest.raises(ValueError, match="no matching slide title found"):
        replace_screenshot_pictures(
            source,
            {"Missing Section": replacement},
            output,
        )
