from __future__ import annotations

import base64
import hashlib
from pathlib import Path
from typing import cast

import pytest

pptx = pytest.importorskip("pptx", reason="python-pptx required for PPTX replacement workflow tests")
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.picture import Picture
from pptx.slide import Slide
from pptx.util import Inches

from counter_risk.writers.pptx_screenshots import replace_screenshot_pictures

_RED_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO6n8e0AAAAASUVORK5CYII="
)
_BLUE_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAwMBAX0X6XQAAAAASUVORK5CYII="
)
_GREEN_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP4DwQACfsD/QX+vJ8AAAAASUVORK5CYII="
)

_TARGET_PICTURE_GEOMETRY = {
    "All Programs": (0, 811705, 9144000, 5817695),
    "Ex Trend": (0, 1304636, 9144000, 5172364),
    "Trend": (0, 2031298, 9144000, 2795403),
}


def _write_png(path: Path, payload: bytes) -> None:
    path.write_bytes(payload)


def _make_workflow_pptx(path: Path, base_image: Path) -> None:
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    for title in ("All Programs", "Ex Trend", "Trend"):
        slide = prs.slides.add_slide(blank_layout)
        textbox = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.25),
            Inches(8.5),
            Inches(0.6),
        )
        textbox.text_frame.text = title

        left, top, width, height = _TARGET_PICTURE_GEOMETRY[title]
        slide.shapes.add_picture(
            str(base_image),
            left,
            top,
            width=width,
            height=height,
        )

    prs.save(str(path))


def _make_near_match_workflow_pptx(path: Path, base_image: Path) -> None:
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    exact_slide = prs.slides.add_slide(blank_layout)
    exact_box = exact_slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(8.5), Inches(0.6))
    exact_box.text_frame.text = "All Programs"
    left, top, width, height = _TARGET_PICTURE_GEOMETRY["All Programs"]
    exact_slide.shapes.add_picture(str(base_image), left, top, width=width, height=height)

    near_slide = prs.slides.add_slide(blank_layout)
    near_box = near_slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(8.5), Inches(0.6))
    near_box.text_frame.text = "All Programs Summary"
    near_slide.shapes.add_picture(str(base_image), left, top + 1, width=width, height=height)

    prs.save(str(path))


def _slide_picture_state(
    slide: Slide,
) -> tuple[int, tuple[tuple[str, tuple[int, int, int, int]], ...]]:
    pictures = [
        cast(Picture, shape) for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    details = tuple(
        (
            hashlib.sha256(picture.image.blob).hexdigest(),
            (picture.left, picture.top, picture.width, picture.height),
        )
        for picture in pictures
    )
    return len(pictures), details


def test_replacement_workflow_output_reopens(tmp_path: Path) -> None:
    source = tmp_path / "input.pptx"
    output = tmp_path / "output.pptx"

    base_image = tmp_path / "base.png"
    all_programs_image = tmp_path / "all_programs.png"
    ex_trend_image = tmp_path / "ex_trend.png"
    trend_image = tmp_path / "trend.png"

    _write_png(base_image, _RED_PNG)
    _write_png(all_programs_image, _BLUE_PNG)
    _write_png(ex_trend_image, _GREEN_PNG)
    _write_png(trend_image, _BLUE_PNG)

    _make_workflow_pptx(source, base_image)
    before = Presentation(str(source))
    before_by_title = {
        slide.shapes[0].text_frame.text: _slide_picture_state(slide) for slide in before.slides
    }

    replace_screenshot_pictures(
        source,
        {
            "All Programs": all_programs_image,
            "Ex Trend": ex_trend_image,
            "Trend": trend_image,
        },
        output,
    )

    reopened = Presentation(str(output))
    assert len(reopened.slides) == 3
    after_by_title = {
        slide.shapes[0].text_frame.text: _slide_picture_state(slide) for slide in reopened.slides
    }

    for title in ("All Programs", "Ex Trend", "Trend"):
        before_count, before_details = before_by_title[title]
        after_count, after_details = after_by_title[title]
        assert before_count == after_count
        assert len(before_details) == len(after_details)
        assert any(
            before_hash != after_hash
            for (before_hash, _), (after_hash, _) in zip(before_details, after_details, strict=True)
        )
        assert [geometry for _, geometry in before_details] == [
            geometry for _, geometry in after_details
        ]


def test_replacement_workflow_near_match_slide_remains_unchanged(tmp_path: Path) -> None:
    source = tmp_path / "near-input.pptx"
    output = tmp_path / "near-output.pptx"

    base_image = tmp_path / "near-base.png"
    replacement_image = tmp_path / "near-replacement.png"
    _write_png(base_image, _RED_PNG)
    _write_png(replacement_image, _GREEN_PNG)

    _make_near_match_workflow_pptx(source, base_image)
    before = Presentation(str(source))
    before_by_title = {
        slide.shapes[0].text_frame.text: _slide_picture_state(slide) for slide in before.slides
    }

    replace_screenshot_pictures(source, {"All Programs": replacement_image}, output)
    reopened = Presentation(str(output))
    after_by_title = {
        slide.shapes[0].text_frame.text: _slide_picture_state(slide) for slide in reopened.slides
    }

    exact_before = before_by_title["All Programs"]
    exact_after = after_by_title["All Programs"]
    assert exact_before[0] == exact_after[0]
    assert any(
        before_hash != after_hash
        for (before_hash, _), (after_hash, _) in zip(exact_before[1], exact_after[1], strict=True)
    )
    assert [geometry for _, geometry in exact_before[1]] == [
        geometry for _, geometry in exact_after[1]
    ]

    near_before = before_by_title["All Programs Summary"]
    near_after = after_by_title["All Programs Summary"]
    assert near_before == near_after
