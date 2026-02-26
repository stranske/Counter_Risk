"""Tests for append_concentration_table_slide in concentration_table module."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from counter_risk.ppt.concentration_table import append_concentration_table_slide

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

_SAMPLE_METRICS: list[dict[str, Any]] = [
    {
        "variant": "all_programs",
        "segment": "total",
        "top5_share": 0.65,
        "top10_share": 0.90,
        "hhi": 0.12,
    },
    {
        "variant": "ex_trend",
        "segment": "total",
        "top5_share": 0.50,
        "top10_share": 0.80,
        "hhi": 0.08,
    },
]


def _make_minimal_pptx(path: Path) -> None:
    """Create a minimal valid PPTX file at *path* using python-pptx."""
    from pptx import Presentation

    prs = Presentation()
    prs.save(str(path))


def _slide_count(path: Path) -> int:
    """Return the number of slides in the PPTX at *path*."""
    from pptx import Presentation

    return len(Presentation(str(path)).slides)


def _last_slide_table_cell_texts(path: Path) -> list[list[str]]:
    """Return cell texts from the table on the last slide as a 2-D list."""
    from pptx import Presentation

    prs = Presentation(str(path))
    slide = prs.slides[-1]
    for shape in slide.shapes:
        if shape.has_table:
            tbl = shape.table
            return [[cell.text for cell in row.cells] for row in tbl.rows]
    return []


# ---------------------------------------------------------------------------
# Tests: slide is appended
# ---------------------------------------------------------------------------


def test_append_concentration_table_slide_increases_slide_count(tmp_path: Path) -> None:
    """Appending a table slide increments the total slide count by one."""
    pptx = tmp_path / "deck.pptx"
    _make_minimal_pptx(pptx)
    before = _slide_count(pptx)
    append_concentration_table_slide(pptx, _SAMPLE_METRICS)
    assert _slide_count(pptx) == before + 1


def test_append_concentration_table_slide_modifies_file_in_place(tmp_path: Path) -> None:
    """The PPTX is modified in-place at the same path."""
    pptx = tmp_path / "deck.pptx"
    _make_minimal_pptx(pptx)
    append_concentration_table_slide(pptx, _SAMPLE_METRICS)
    assert pptx.exists()
    assert pptx.stat().st_size > 0


def test_append_concentration_table_slide_with_empty_metrics(tmp_path: Path) -> None:
    """Empty metrics list still appends a slide (title only, no table)."""
    pptx = tmp_path / "deck.pptx"
    _make_minimal_pptx(pptx)
    before = _slide_count(pptx)
    append_concentration_table_slide(pptx, [])
    assert _slide_count(pptx) == before + 1


# ---------------------------------------------------------------------------
# Tests: table structure and content
# ---------------------------------------------------------------------------


def test_append_concentration_table_slide_has_expected_header_row(tmp_path: Path) -> None:
    """The table header row contains all required column names."""
    pptx = tmp_path / "deck.pptx"
    _make_minimal_pptx(pptx)
    append_concentration_table_slide(pptx, _SAMPLE_METRICS)
    rows = _last_slide_table_cell_texts(pptx)
    assert len(rows) > 0
    header = rows[0]
    assert "Variant" in header
    assert "Segment" in header
    assert "Top 5 Share" in header
    assert "Top 10 Share" in header
    assert "HHI" in header


def test_append_concentration_table_slide_row_count_matches_records(tmp_path: Path) -> None:
    """The table has one header row plus one data row per metrics record."""
    pptx = tmp_path / "deck.pptx"
    _make_minimal_pptx(pptx)
    append_concentration_table_slide(pptx, _SAMPLE_METRICS)
    rows = _last_slide_table_cell_texts(pptx)
    assert len(rows) == len(_SAMPLE_METRICS) + 1  # header + data rows


def test_append_concentration_table_slide_data_rows_contain_variant_and_segment(
    tmp_path: Path,
) -> None:
    """Data rows correctly reflect variant and segment values from records."""
    pptx = tmp_path / "deck.pptx"
    _make_minimal_pptx(pptx)
    append_concentration_table_slide(pptx, _SAMPLE_METRICS)
    rows = _last_slide_table_cell_texts(pptx)
    data_rows = rows[1:]
    variants = [row[0] for row in data_rows]
    segments = [row[1] for row in data_rows]
    assert "all_programs" in variants
    assert "ex_trend" in variants
    assert "total" in segments


def test_append_concentration_table_slide_shares_formatted_as_percent(tmp_path: Path) -> None:
    """Top 5 and Top 10 share cells use percentage formatting."""
    pptx = tmp_path / "deck.pptx"
    _make_minimal_pptx(pptx)
    append_concentration_table_slide(pptx, _SAMPLE_METRICS)
    rows = _last_slide_table_cell_texts(pptx)
    # First data row, column index 2 = top5_share, 3 = top10_share
    top5_text = rows[1][2]
    top10_text = rows[1][3]
    assert "%" in top5_text, f"Expected '%' in top5 cell: {top5_text!r}"
    assert "%" in top10_text, f"Expected '%' in top10 cell: {top10_text!r}"


def test_append_concentration_table_slide_hhi_formatted_as_decimal(tmp_path: Path) -> None:
    """HHI cells use decimal (not percentage) formatting."""
    pptx = tmp_path / "deck.pptx"
    _make_minimal_pptx(pptx)
    append_concentration_table_slide(pptx, _SAMPLE_METRICS)
    rows = _last_slide_table_cell_texts(pptx)
    hhi_text = rows[1][4]
    assert "%" not in hhi_text, f"HHI should not contain '%': {hhi_text!r}"
    # Should be parseable as a float
    float(hhi_text)


# ---------------------------------------------------------------------------
# Tests: toggle behaviour (via config)
# ---------------------------------------------------------------------------


def test_config_include_concentration_table_in_ppt_defaults_to_false() -> None:
    """WorkflowConfig.include_concentration_table_in_ppt defaults to False."""
    from pathlib import Path

    from counter_risk.config import WorkflowConfig

    cfg = WorkflowConfig(
        mosers_ex_trend_xlsx=Path("ex.xlsx"),
        mosers_trend_xlsx=Path("trend.xlsx"),
        hist_all_programs_3yr_xlsx=Path("h1.xlsx"),
        hist_ex_llc_3yr_xlsx=Path("h2.xlsx"),
        hist_llc_3yr_xlsx=Path("h3.xlsx"),
        monthly_pptx=Path("monthly.pptx"),
    )
    assert cfg.include_concentration_table_in_ppt is False


def test_config_include_concentration_table_in_ppt_can_be_enabled() -> None:
    """WorkflowConfig.include_concentration_table_in_ppt can be set to True."""
    from pathlib import Path

    from counter_risk.config import WorkflowConfig

    cfg = WorkflowConfig(
        mosers_ex_trend_xlsx=Path("ex.xlsx"),
        mosers_trend_xlsx=Path("trend.xlsx"),
        hist_all_programs_3yr_xlsx=Path("h1.xlsx"),
        hist_ex_llc_3yr_xlsx=Path("h2.xlsx"),
        hist_llc_3yr_xlsx=Path("h3.xlsx"),
        monthly_pptx=Path("monthly.pptx"),
        include_concentration_table_in_ppt=True,
    )
    assert cfg.include_concentration_table_in_ppt is True


# ---------------------------------------------------------------------------
# Tests: custom title
# ---------------------------------------------------------------------------


def test_append_concentration_table_slide_accepts_custom_title(tmp_path: Path) -> None:
    """A custom title string is accepted without error."""
    pptx = tmp_path / "deck.pptx"
    _make_minimal_pptx(pptx)
    before = _slide_count(pptx)
    append_concentration_table_slide(pptx, _SAMPLE_METRICS, title="Custom Concentration Summary")
    assert _slide_count(pptx) == before + 1


def test_append_concentration_table_slide_single_record(tmp_path: Path) -> None:
    """A single-record metrics list produces a table with one data row."""
    pptx = tmp_path / "deck.pptx"
    _make_minimal_pptx(pptx)
    single = [_SAMPLE_METRICS[0]]
    append_concentration_table_slide(pptx, single)
    rows = _last_slide_table_cell_texts(pptx)
    assert len(rows) == 2  # header + 1 data row


def test_append_concentration_table_slide_accepts_path_str(tmp_path: Path) -> None:
    """pptx_path can be a plain string as well as a Path object."""
    pptx = tmp_path / "deck.pptx"
    _make_minimal_pptx(pptx)
    before = _slide_count(pptx)
    append_concentration_table_slide(str(pptx), _SAMPLE_METRICS)
    assert _slide_count(pptx) == before + 1
