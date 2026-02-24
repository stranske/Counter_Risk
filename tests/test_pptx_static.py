from __future__ import annotations

import base64
from pathlib import Path

import pytest

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches
except ModuleNotFoundError:  # pragma: no cover - optional dependency in some environments
    pytest.skip(
        "python-pptx required for static pptx rebuild tests",
        allow_module_level=True,
    )

from counter_risk.ppt.pptx_static import _rebuild_pptx_from_slide_images

_ONE_PIXEL_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO6n8e0AAAAASUVORK5CYII="
)


def _write_png(path: Path) -> None:
    path.write_bytes(_ONE_PIXEL_PNG)


def _make_source_with_chart(path: Path, image_path: Path) -> None:
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    slide1 = prs.slides.add_slide(blank_layout)
    chart_data = CategoryChartData()  # type: ignore[no-untyped-call]
    chart_data.categories = ["A", "B"]
    chart_data.add_series("Series 1", (1, 2))  # type: ignore[no-untyped-call]
    slide1.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(5),
        Inches(3),
        chart_data,
    )

    slide2 = prs.slides.add_slide(blank_layout)
    slide2.shapes.add_picture(str(image_path), Inches(1), Inches(1), Inches(4), Inches(3))

    prs.save(str(path))


def test_rebuild_pptx_from_slide_images_outputs_picture_only_slides(tmp_path: Path) -> None:
    source = tmp_path / "source.pptx"
    output = tmp_path / "rebuilt.pptx"

    source_img = tmp_path / "source.png"
    _write_png(source_img)
    _make_source_with_chart(source, source_img)

    slide_image_1 = tmp_path / "slide_0001.png"
    slide_image_2 = tmp_path / "slide_0002.png"
    _write_png(slide_image_1)
    _write_png(slide_image_2)

    _rebuild_pptx_from_slide_images(
        source_pptx=source,
        slide_images=[slide_image_1, slide_image_2],
        output_path=output,
    )

    rebuilt = Presentation(str(output))
    assert len(rebuilt.slides) == 2
    for slide in rebuilt.slides:
        picture_shapes = [
            shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]
        chart_shapes = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.CHART]
        assert len(picture_shapes) == 1
        assert chart_shapes == []
