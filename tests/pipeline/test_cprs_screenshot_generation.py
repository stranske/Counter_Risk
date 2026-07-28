"""Tests for CPRS-CH/CPRS-FCM screenshot generation orchestration."""

from __future__ import annotations

from pathlib import Path
from typing import Any

import pytest

import counter_risk.pipeline.run as run_module
from counter_risk.config import WorkflowConfig
from counter_risk.integrations import excel_com


def _build_config(tmp_path: Path, **overrides: Any) -> WorkflowConfig:
    placeholder_files = (
        "all.xlsx",
        "ex.xlsx",
        "trend.xlsx",
        "hist_all.xlsx",
        "hist_ex.xlsx",
        "hist_llc.xlsx",
        "monthly.pptx",
    )
    for filename in placeholder_files:
        (tmp_path / filename).write_bytes(b"placeholder")

    defaults: dict[str, Any] = {
        "mosers_all_programs_xlsx": tmp_path / "all.xlsx",
        "mosers_ex_trend_xlsx": tmp_path / "ex.xlsx",
        "mosers_trend_xlsx": tmp_path / "trend.xlsx",
        "hist_all_programs_3yr_xlsx": tmp_path / "hist_all.xlsx",
        "hist_ex_llc_3yr_xlsx": tmp_path / "hist_ex.xlsx",
        "hist_llc_3yr_xlsx": tmp_path / "hist_llc.xlsx",
        "monthly_pptx": tmp_path / "monthly.pptx",
    }
    defaults.update(overrides)
    return WorkflowConfig(**defaults)


def test_generates_screenshot_for_each_configured_variant_and_table(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    config = _build_config(tmp_path)
    run_dir = tmp_path / "run"
    warnings: list[str] = []
    calls: list[dict[str, Any]] = []

    def _fake_export(*, workbook_path: Path, sheet_name: str, output_png: Path) -> Path:
        calls.append(
            {"workbook_path": workbook_path, "sheet_name": sheet_name, "output_png": output_png}
        )
        output_png.parent.mkdir(parents=True, exist_ok=True)
        output_png.write_bytes(b"fake-png")
        return output_png

    monkeypatch.setattr(excel_com, "export_worksheet_range_as_png", _fake_export)

    result = run_module._generate_cprs_screenshot_inputs(
        config=config, run_dir=run_dir, warnings=warnings
    )

    assert set(result) == {"slide1", "slide2", "slide6", "slide7", "slide16", "slide17"}
    assert warnings == []
    assert len(calls) == 6

    call_by_sheet_and_workbook = {(call["sheet_name"], call["workbook_path"]) for call in calls}
    assert (("CPRS - CH", config.mosers_all_programs_xlsx)) in call_by_sheet_and_workbook
    assert (("CPRS - FCM", config.mosers_all_programs_xlsx)) in call_by_sheet_and_workbook
    assert (("CPRS - CH", config.mosers_ex_trend_xlsx)) in call_by_sheet_and_workbook
    assert (("CPRS - FCM", config.mosers_ex_trend_xlsx)) in call_by_sheet_and_workbook
    assert (("CPRS - CH", config.mosers_trend_xlsx)) in call_by_sheet_and_workbook
    assert (("CPRS - FCM", config.mosers_trend_xlsx)) in call_by_sheet_and_workbook

    for path in result.values():
        assert path.is_relative_to(run_dir)
        assert path.read_bytes() == b"fake-png"


def test_skips_variant_without_configured_workbook(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    config = _build_config(tmp_path, mosers_trend_xlsx=None)
    run_dir = tmp_path / "run"
    warnings: list[str] = []

    def _fake_export(*, workbook_path: Path, sheet_name: str, output_png: Path) -> Path:
        output_png.parent.mkdir(parents=True, exist_ok=True)
        output_png.write_bytes(b"fake-png")
        return output_png

    monkeypatch.setattr(excel_com, "export_worksheet_range_as_png", _fake_export)

    result = run_module._generate_cprs_screenshot_inputs(
        config=config, run_dir=run_dir, warnings=warnings
    )

    assert set(result) == {"slide1", "slide2", "slide6", "slide7"}
    assert "slide16" not in result
    assert "slide17" not in result
    assert warnings == []


def test_export_failure_appends_warning_and_continues(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    config = _build_config(tmp_path)
    run_dir = tmp_path / "run"
    warnings: list[str] = []

    def _fake_export(*, workbook_path: Path, sheet_name: str, output_png: Path) -> Path:
        if sheet_name == "CPRS - FCM" and workbook_path == config.mosers_all_programs_xlsx:
            raise excel_com.ExcelComError("boom")
        output_png.parent.mkdir(parents=True, exist_ok=True)
        output_png.write_bytes(b"fake-png")
        return output_png

    monkeypatch.setattr(excel_com, "export_worksheet_range_as_png", _fake_export)

    result = run_module._generate_cprs_screenshot_inputs(
        config=config, run_dir=run_dir, warnings=warnings
    )

    assert "slide2" not in result
    assert set(result) == {"slide1", "slide6", "slide7", "slide16", "slide17"}
    assert len(warnings) == 1
    assert "all_programs/fcm" in warnings[0]
    assert "slide 2" in warnings[0]


def test_pad_image_to_aspect_ratio_widens_canvas_without_cropping(tmp_path: Path) -> None:
    from PIL import Image

    image_path = tmp_path / "table.png"
    Image.new("RGB", (200, 200), (10, 20, 30)).save(image_path, "PNG")

    run_module._pad_image_to_aspect_ratio(image_path, target_ratio=2.0)

    with Image.open(image_path) as padded:
        assert padded.size == (400, 200)
        # Original content is preserved, centered, not cropped or squished.
        assert padded.getpixel((50, 100)) == (255, 255, 255)
        assert padded.getpixel((200, 100)) == (10, 20, 30)


def test_pad_image_to_aspect_ratio_heightens_canvas_without_cropping(tmp_path: Path) -> None:
    from PIL import Image

    image_path = tmp_path / "table.png"
    Image.new("RGB", (200, 200), (10, 20, 30)).save(image_path, "PNG")

    run_module._pad_image_to_aspect_ratio(image_path, target_ratio=0.5)

    with Image.open(image_path) as padded:
        assert padded.size == (200, 400)
        assert padded.getpixel((100, 50)) == (255, 255, 255)
        assert padded.getpixel((100, 200)) == (10, 20, 30)


def test_pad_image_to_aspect_ratio_noop_when_already_matching(tmp_path: Path) -> None:
    from PIL import Image

    image_path = tmp_path / "table.png"
    Image.new("RGB", (300, 200), (10, 20, 30)).save(image_path, "PNG")
    original_bytes = image_path.read_bytes()

    run_module._pad_image_to_aspect_ratio(image_path, target_ratio=1.5)

    assert image_path.read_bytes() == original_bytes


def test_picture_shape_aspect_ratio_returns_none_for_unreadable_pptx(tmp_path: Path) -> None:
    bad_pptx = tmp_path / "not_a_pptx.pptx"
    bad_pptx.write_bytes(b"not a real pptx")

    assert run_module._picture_shape_aspect_ratio(pptx_path=bad_pptx, slide_number=1) is None


def test_picture_shape_aspect_ratio_reads_real_template(tmp_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Emu

    pptx_path = tmp_path / "template.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    placeholder_png = tmp_path / "placeholder.png"
    from PIL import Image

    Image.new("RGB", (10, 10), (255, 255, 255)).save(placeholder_png, "PNG")
    slide.shapes.add_picture(
        str(placeholder_png), Emu(0), Emu(0), width=Emu(914400 * 4), height=Emu(914400 * 2)
    )
    presentation.save(str(pptx_path))

    ratio = run_module._picture_shape_aspect_ratio(pptx_path=pptx_path, slide_number=1)

    assert ratio == pytest.approx(2.0)
