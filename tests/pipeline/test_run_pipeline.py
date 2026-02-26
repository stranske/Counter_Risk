"""Integration-style tests for pipeline orchestration."""

from __future__ import annotations

import csv
import hashlib
import json
import platform
import re
import sys
import types
from datetime import date
from pathlib import Path
from typing import Any, Literal, cast

import pytest

import counter_risk.pipeline.run as run_module
from counter_risk.config import ReconciliationConfig, WorkflowConfig
from counter_risk.pipeline.parsing_types import UnmappedCounterpartyError
from counter_risk.pipeline.run import run_pipeline


class _FakeDataFrame:
    def __init__(
        self,
        records: list[dict[str, Any]] | None = None,
        columns: list[str] | tuple[str, ...] | None = None,
    ) -> None:
        self._rows = [dict(row) for row in (records or [])]
        if columns is not None:
            self.columns: list[str] = list(columns)
        elif self._rows:
            self.columns = list(self._rows[0].keys())
        else:
            self.columns = []

    @property
    def empty(self) -> bool:
        return len(self._rows) == 0

    @property
    def loc(self) -> _LocIndexer:
        return _LocIndexer(self)

    def __setitem__(self, key: str, value: Any) -> None:
        if key not in self.columns:
            self.columns.append(key)
        for row in self._rows:
            row[key] = value

    def astype(self, dtypes: dict[str, str]) -> _FakeDataFrame:
        for row in self._rows:
            for column, dtype in dtypes.items():
                if column not in row:
                    continue
                if dtype == "float64":
                    row[column] = float(row[column])
                elif dtype == "int64":
                    row[column] = int(row[column])
                elif dtype == "string":
                    row[column] = str(row[column])
        return self

    def to_dict(self, orient: str = "records") -> list[dict[str, Any]]:
        if orient != "records":
            raise ValueError("Only orient='records' is supported")
        return [dict(row) for row in self._rows]


class _LocIndexer:
    def __init__(self, frame: _FakeDataFrame) -> None:
        self._frame = frame

    def __getitem__(self, key: tuple[slice, list[str]]) -> _FakeDataFrame:
        _rows_slice, columns = key
        records = [{column: row.get(column) for column in columns} for row in self._frame._rows]
        return _FakeDataFrame(records=records, columns=columns)


class _FakeCell:
    def __init__(self, value: Any = None) -> None:
        self.value = value


class _FakeWorksheet:
    def __init__(self, title: str) -> None:
        self.title = title
        self.max_row = 1
        self.max_column = 1
        self._cells: dict[tuple[int, int], _FakeCell] = {}

    def cell(self, row: int, column: int) -> _FakeCell:
        self.max_row = max(self.max_row, row)
        self.max_column = max(self.max_column, column)
        key = (row, column)
        if key not in self._cells:
            self._cells[key] = _FakeCell()
        return self._cells[key]

    def set_value(self, row: int, column: int, value: Any) -> None:
        self.cell(row=row, column=column).value = value


class _FakeWorkbook:
    def __init__(self, sheets: dict[str, _FakeWorksheet]) -> None:
        self._sheets = dict(sheets)
        self.sheetnames = list(sheets)
        self.saved_paths: list[Path] = []
        self.closed = False

    @property
    def active(self) -> _FakeWorksheet:
        return self._sheets[self.sheetnames[0]]

    def __getitem__(self, item: str) -> _FakeWorksheet:
        return self._sheets[item]

    def save(self, path: Path) -> None:
        self.saved_paths.append(path)

    def close(self) -> None:
        self.closed = True


@pytest.fixture
def fake_pandas(monkeypatch: pytest.MonkeyPatch) -> None:
    fake_module = types.SimpleNamespace(DataFrame=_FakeDataFrame)
    monkeypatch.setitem(sys.modules, "pandas", fake_module)


@pytest.fixture(autouse=True)
def patch_repo_root(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setattr("counter_risk.pipeline.run._resolve_repo_root", lambda: tmp_path)


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    digest.update(path.read_bytes())
    return digest.hexdigest()


def _minimal_parsed_by_variant() -> dict[str, dict[str, _FakeDataFrame]]:
    totals = _FakeDataFrame(
        records=[{"counterparty": "Counterparty A", "Notional": 1.0, "NotionalChange": 0.5}]
    )
    futures = _FakeDataFrame(
        records=[
            {
                "account": "account-a",
                "description": "desc",
                "class": "class-a",
                "fcm": "fcm-a",
                "clearing_house": "ch-a",
                "notional": 1.0,
            }
        ]
    )
    return {
        "all_programs": {"totals": totals, "futures": futures},
        "ex_trend": {"totals": totals, "futures": futures},
        "trend": {"totals": totals, "futures": futures},
    }


def _minimal_workflow_config(
    tmp_path: Path, *, fail_policy: Literal["warn", "strict"] = "warn"
) -> WorkflowConfig:
    return WorkflowConfig(
        as_of_date=date(2025, 12, 31),
        output_root=tmp_path / "runs",
        reconciliation=ReconciliationConfig(fail_policy=fail_policy),
        mosers_all_programs_xlsx=tmp_path / "all.xlsx",
        mosers_ex_trend_xlsx=tmp_path / "ex.xlsx",
        mosers_trend_xlsx=tmp_path / "trend.xlsx",
        hist_all_programs_3yr_xlsx=tmp_path / "hist-all.xlsx",
        hist_ex_llc_3yr_xlsx=tmp_path / "hist-ex.xlsx",
        hist_llc_3yr_xlsx=tmp_path / "hist-trend.xlsx",
        monthly_pptx=tmp_path / "monthly.pptx",
    )


def test_build_parsed_data_by_sheet_uses_historical_sheet_names_without_total_key() -> None:
    parsed_sections = {
        "totals": _FakeDataFrame(
            records=[
                {"counterparty": "Counterparty A", "Notional": 10.0},
                {"counterparty": "Counterparty B", "Notional": 20.0},
            ]
        ),
        "futures": _FakeDataFrame(
            records=[
                {
                    "account": "acct",
                    "description": "desc",
                    "class": "cls",
                    "fcm": "fcm",
                    "clearing_house": "Clearing B",
                    "notional": 5.0,
                }
            ]
        ),
    }
    historical_headers_by_sheet = {
        "Sheet A": ("Counterparty A",),
        "Sheet B": ("Counterparty B", "Clearing B"),
    }

    parsed_data_by_sheet = run_module._build_parsed_data_by_sheet(
        parsed_sections=parsed_sections,
        historical_series_headers_by_sheet=historical_headers_by_sheet,
    )

    assert set(parsed_data_by_sheet) == {"Sheet A", "Sheet B"}
    assert "Total" not in parsed_data_by_sheet
    assert [row["counterparty"] for row in parsed_data_by_sheet["Sheet A"]["totals"]] == [
        "Counterparty A"
    ]
    assert [row["counterparty"] for row in parsed_data_by_sheet["Sheet B"]["totals"]] == [
        "Counterparty B"
    ]
    assert [row["clearing_house"] for row in parsed_data_by_sheet["Sheet B"]["futures"]] == [
        "Clearing B"
    ]


def test_run_reconciliation_checks_reports_gap_only_for_impacted_sheet(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    config = _minimal_workflow_config(tmp_path)
    warnings: list[str] = []
    parsed_by_variant = {
        "all_programs": {
            "totals": _FakeDataFrame(records=[{"counterparty": "Counterparty A", "Notional": 1.0}]),
            "futures": _FakeDataFrame(records=[]),
        },
        "ex_trend": {"totals": _FakeDataFrame(records=[]), "futures": _FakeDataFrame(records=[])},
        "trend": {"totals": _FakeDataFrame(records=[]), "futures": _FakeDataFrame(records=[])},
    }

    monkeypatch.setattr(
        run_module,
        "_extract_historical_series_headers_by_sheet",
        lambda _: {"Sheet A": ("Counterparty A",), "Sheet B": ("Counterparty B",)},
    )

    run_module._run_reconciliation_checks(
        run_dir=tmp_path,
        config=config,
        parsed_by_variant=parsed_by_variant,
        warnings=warnings,
    )

    assert any("sheet 'Sheet B'" in warning for warning in warnings)
    assert not any("sheet 'Sheet A'" in warning for warning in warnings)


def test_run_reconciliation_checks_counts_only_rows_tied_to_missing_series(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    config = _minimal_workflow_config(tmp_path)
    warnings: list[str] = []
    parsed_by_variant = {
        "all_programs": {
            "totals": _FakeDataFrame(
                records=[
                    {"counterparty": "Counterparty A", "Notional": 1.0},
                    {"counterparty": "Counterparty B", "Notional": 2.0},
                    {"counterparty": "Counterparty C", "Notional": 3.0},
                ]
            ),
            "futures": _FakeDataFrame(records=[]),
        },
        "ex_trend": {"totals": _FakeDataFrame(records=[]), "futures": _FakeDataFrame(records=[])},
        "trend": {"totals": _FakeDataFrame(records=[]), "futures": _FakeDataFrame(records=[])},
    }

    monkeypatch.setattr(
        run_module,
        "_extract_historical_series_headers_by_sheet",
        lambda _: {"Sheet A": ("Counterparty A", "Counterparty C")},
    )

    run_module._run_reconciliation_checks(
        run_dir=tmp_path,
        config=config,
        parsed_by_variant=parsed_by_variant,
        warnings=warnings,
    )

    mapping_updates = tmp_path / "NEEDS_MAPPING_UPDATES.txt"
    assert mapping_updates.exists()
    assert "impacted_series_count: 1" in mapping_updates.read_text(encoding="utf-8")
    assert "impacted_rows_count: 1" in mapping_updates.read_text(encoding="utf-8")
    assert any("impacted_series=1" in warning for warning in warnings)
    assert any("impacted_rows=1" in warning for warning in warnings)


def test_write_needs_mapping_updates_ignores_unmapped_counterparty_metadata_for_legacy_output(
    tmp_path: Path,
) -> None:
    output_path = run_module._write_needs_mapping_updates(
        run_dir=tmp_path,
        fail_policy="warn",
        reconciliation_by_variant={
            "all_programs": {
                "missing_series": [
                    {
                        "sheet": "Total",
                        "missing_from_historical_headers": ["Counterparty A"],
                    },
                    {
                        "sheet": "Total",
                        "error_type": "unmapped_counterparty",
                        "raw_counterparties": [" ACME  LTD "],
                        "normalized_counterparties": ["ACME LTD"],
                    },
                ]
            }
        },
        total_gap_count=2,
        impacted_series_count=2,
        impacted_rows_count=1,
    )

    text = output_path.read_text(encoding="utf-8")
    assert "missing_from_historical_headers=Counterparty A" in text
    assert "raw_counterparties" not in text
    assert "ACME LTD" not in text


def test_manifest_impacted_rows_counts_only_matching_normalized_label() -> None:
    parsed_sections = {
        "totals": [
            {"counterparty": " ACME  LTD ", "normalized_label": "ACME LTD", "Notional": 1.0},
            {"counterparty": "ACME LTD", "normalized_label": "ACME LTD", "Notional": 2.0},
            {"counterparty": "Beta LLC", "normalized_label": "BETA LLC", "Notional": 3.0},
        ],
        "futures": [],
    }
    reconciliation_sheet_result = {
        "missing_from_historical_headers": [],
        "missing_from_data": [],
        "missing_normalized_counterparties": ["ACME LTD"],
    }

    impacted_series, impacted_rows = run_module._calculate_impacted_scope_for_sheet(
        parsed_sections=parsed_sections,
        reconciliation_sheet_result=reconciliation_sheet_result,
    )

    assert impacted_series == 1
    assert impacted_rows == 2


def test_run_reconciliation_checks_counts_only_impacted_series_when_other_rows_unaffected(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    config = _minimal_workflow_config(tmp_path)
    warnings: list[str] = []
    parsed_by_variant = {
        "all_programs": {
            "totals": _FakeDataFrame(
                records=[
                    {"counterparty": "Counterparty A", "Notional": 1.0},
                    {"counterparty": "Counterparty B", "Notional": 2.0},
                    {"counterparty": "Counterparty C", "Notional": 3.0},
                    {"counterparty": "Counterparty D", "Notional": 4.0},
                ]
            ),
            "futures": _FakeDataFrame(
                records=[
                    {
                        "account": "acct",
                        "description": "desc",
                        "class": "cls",
                        "fcm": "fcm",
                        "clearing_house": "CME",
                        "notional": 5.0,
                    }
                ]
            ),
        },
        "ex_trend": {"totals": _FakeDataFrame(records=[]), "futures": _FakeDataFrame(records=[])},
        "trend": {"totals": _FakeDataFrame(records=[]), "futures": _FakeDataFrame(records=[])},
    }

    monkeypatch.setattr(
        run_module,
        "_extract_historical_series_headers_by_sheet",
        lambda _: {"Sheet A": ("Counterparty A", "Counterparty D", "CME")},
    )

    run_module._run_reconciliation_checks(
        run_dir=tmp_path,
        config=config,
        parsed_by_variant=parsed_by_variant,
        warnings=warnings,
    )

    mapping_updates = tmp_path / "NEEDS_MAPPING_UPDATES.txt"
    assert mapping_updates.exists()
    text = mapping_updates.read_text(encoding="utf-8")
    assert "impacted_series_count: 2" in text
    assert "impacted_rows_count: 2" in text
    assert any("impacted_series=2" in warning for warning in warnings)
    assert any("impacted_rows=2" in warning for warning in warnings)


def test_run_reconciliation_checks_counts_missing_from_data_as_impacted_series_only(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    config = _minimal_workflow_config(tmp_path)
    warnings: list[str] = []
    parsed_by_variant = {
        "all_programs": {
            "totals": _FakeDataFrame(
                records=[
                    {"counterparty": "Counterparty A", "Notional": 1.0},
                ]
            ),
            "futures": _FakeDataFrame(records=[]),
        },
        "ex_trend": {"totals": _FakeDataFrame(records=[]), "futures": _FakeDataFrame(records=[])},
        "trend": {"totals": _FakeDataFrame(records=[]), "futures": _FakeDataFrame(records=[])},
    }

    monkeypatch.setattr(
        run_module,
        "_extract_historical_series_headers_by_sheet",
        lambda _: {"Sheet A": ("Counterparty A", "Counterparty B", "Counterparty C")},
    )

    run_module._run_reconciliation_checks(
        run_dir=tmp_path,
        config=config,
        parsed_by_variant=parsed_by_variant,
        warnings=warnings,
    )

    mapping_updates = tmp_path / "NEEDS_MAPPING_UPDATES.txt"
    assert mapping_updates.exists()
    text = mapping_updates.read_text(encoding="utf-8")
    assert "impacted_series_count: 2" in text
    assert "impacted_rows_count: 0" in text
    assert any("impacted_series=2" in warning for warning in warnings)
    assert any("impacted_rows=0" in warning for warning in warnings)


def test_run_reconciliation_checks_segment_gaps_do_not_increase_impacted_scope(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    config = _minimal_workflow_config(tmp_path)
    config.reconciliation = ReconciliationConfig(
        fail_policy="warn",
        expected_segments_by_variant={"all_programs": ["swaps", "repo"]},
    )
    warnings: list[str] = []
    parsed_by_variant = {
        "all_programs": {
            "totals": _FakeDataFrame(records=[{"counterparty": "Counterparty A", "Notional": 1.0}]),
            "futures": _FakeDataFrame(
                records=[{"clearing_house": "CME", "segment": "swaps", "notional": 2.0}]
            ),
        },
        "ex_trend": {"totals": _FakeDataFrame(records=[]), "futures": _FakeDataFrame(records=[])},
        "trend": {"totals": _FakeDataFrame(records=[]), "futures": _FakeDataFrame(records=[])},
    }

    monkeypatch.setattr(
        run_module,
        "_extract_historical_series_headers_by_sheet",
        lambda _: {"Sheet A": ("Counterparty A", "CME")},
    )

    run_module._run_reconciliation_checks(
        run_dir=tmp_path,
        config=config,
        parsed_by_variant=parsed_by_variant,
        warnings=warnings,
    )

    mapping_updates = tmp_path / "NEEDS_MAPPING_UPDATES.txt"
    assert mapping_updates.exists()
    text = mapping_updates.read_text(encoding="utf-8")
    assert "impacted_series_count: 0" in text
    assert "impacted_rows_count: 0" in text
    assert any("impacted_series=0" in warning for warning in warnings)
    assert any("impacted_rows=0" in warning for warning in warnings)


def test_run_reconciliation_checks_segment_gaps_do_not_change_direct_impacted_count(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    config = _minimal_workflow_config(tmp_path)
    config.reconciliation = ReconciliationConfig(
        fail_policy="warn",
        expected_segments_by_variant={"all_programs": ["swaps", "repo"]},
    )
    warnings: list[str] = []
    parsed_by_variant = {
        "all_programs": {
            "totals": _FakeDataFrame(
                records=[
                    {"counterparty": "Counterparty A", "Notional": 1.0},
                    {"counterparty": "Counterparty B", "Notional": 2.0},
                ]
            ),
            "futures": _FakeDataFrame(
                records=[{"clearing_house": "CME", "segment": "swaps", "notional": 3.0}]
            ),
        },
        "ex_trend": {"totals": _FakeDataFrame(records=[]), "futures": _FakeDataFrame(records=[])},
        "trend": {"totals": _FakeDataFrame(records=[]), "futures": _FakeDataFrame(records=[])},
    }

    monkeypatch.setattr(
        run_module,
        "_extract_historical_series_headers_by_sheet",
        lambda _: {"Sheet A": ("Counterparty A", "CME")},
    )

    run_module._run_reconciliation_checks(
        run_dir=tmp_path,
        config=config,
        parsed_by_variant=parsed_by_variant,
        warnings=warnings,
    )

    mapping_updates = tmp_path / "NEEDS_MAPPING_UPDATES.txt"
    assert mapping_updates.exists()
    text = mapping_updates.read_text(encoding="utf-8")
    assert "impacted_series_count: 1" in text
    assert "impacted_rows_count: 1" in text
    assert any("impacted_series=1" in warning for warning in warnings)
    assert any("impacted_rows=1" in warning for warning in warnings)


def test_run_reconciliation_checks_strict_raises_unmapped_counterparty_error_from_result(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    config = _minimal_workflow_config(tmp_path, fail_policy="strict")
    warnings: list[str] = []
    parsed_by_variant = {
        "all_programs": {
            "totals": _FakeDataFrame(records=[{"counterparty": " ACME  LTD ", "Notional": 1.0}]),
            "futures": _FakeDataFrame(records=[]),
        },
        "ex_trend": {"totals": _FakeDataFrame(records=[]), "futures": _FakeDataFrame(records=[])},
        "trend": {"totals": _FakeDataFrame(records=[]), "futures": _FakeDataFrame(records=[])},
    }
    unmapped_error = UnmappedCounterpartyError(
        normalized_counterparty="ACME LTD",
        raw_counterparty=" ACME  LTD ",
        sheet="Total",
    )

    monkeypatch.setattr(
        run_module,
        "_extract_historical_series_headers_by_sheet",
        lambda _: {"Total": ("Legacy Counterparty",)},
    )
    monkeypatch.setattr(
        run_module,
        "reconcile_series_coverage",
        lambda **_: {
            "gap_count": 1,
            "by_sheet": {},
            "warnings": [],
            "missing_series": [],
            "exceptions": [unmapped_error],
        },
    )

    with pytest.raises(UnmappedCounterpartyError) as exc_info:
        run_module._run_reconciliation_checks(
            run_dir=tmp_path,
            config=config,
            parsed_by_variant=parsed_by_variant,
            warnings=warnings,
        )

    assert exc_info.value is unmapped_error
    assert exc_info.value.raw_counterparty == " ACME  LTD "
    assert exc_info.value.normalized_counterparty == "ACME LTD"


def test_run_pipeline_writes_expected_outputs_and_manifest(
    tmp_path: Path, fake_pandas: None
) -> None:
    fixtures = Path("tests/fixtures")
    output_root = tmp_path / "runs"

    config_path = tmp_path / "config.yml"
    config_path.write_text(
        "\n".join(
            [
                "as_of_date: 2025-12-31",
                f"mosers_all_programs_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - All Programs.xlsx'}",
                f"mosers_ex_trend_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - Ex Trend.xlsx'}",
                f"mosers_trend_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - Trend.xlsx'}",
                f"hist_all_programs_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - All Programs 3 Year.xlsx'}",
                f"hist_ex_llc_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - ex LLC 3 Year.xlsx'}",
                f"hist_llc_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - LLC 3 Year.xlsx'}",
                f"monthly_pptx: {fixtures / 'Monthly Counterparty Exposure Report.pptx'}",
                f"output_root: {output_root}",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    run_dir = run_pipeline(config_path)

    assert run_dir == tmp_path / "runs" / "2025-12-31"
    assert run_dir.exists()

    manifest_path = run_dir / "manifest.json"
    assert manifest_path.exists()

    expected_outputs = [
        run_dir / "Historical Counterparty Risk Graphs - All Programs 3 Year.xlsx",
        run_dir / "Historical Counterparty Risk Graphs - ex LLC 3 Year.xlsx",
        run_dir / "Historical Counterparty Risk Graphs - LLC 3 Year.xlsx",
        run_dir / "all_programs-mosers-input.xlsx",
        run_dir / "ex_trend-mosers-input.xlsx",
        run_dir / "trend-mosers-input.xlsx",
        run_dir / "Monthly Counterparty Exposure Report (Master) - 2025-12-31.pptx",
        run_dir / "Monthly Counterparty Exposure Report - 2025-12-31.pptx",
    ]
    for output_file in expected_outputs:
        assert output_file.exists(), f"Missing output file: {output_file}"

    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    assert manifest["as_of_date"] == "2025-12-31"
    assert re.fullmatch(r"\d{4}-\d{2}-\d{2}", manifest["run_date"])
    assert manifest["config_snapshot"]["output_root"] == str(output_root)

    for output_path in manifest["output_paths"]:
        assert not Path(output_path).is_absolute()
        assert (run_dir / output_path).exists(), f"Manifest references missing path: {output_path}"

    assert "PPT links not refreshed; COM refresh skipped" in manifest["warnings"]

    expected_hashes = {
        "mosers_all_programs_xlsx": _sha256(
            fixtures / "MOSERS Counterparty Risk Summary 12-31-2025 - All Programs.xlsx"
        ),
        "mosers_ex_trend_xlsx": _sha256(
            fixtures / "MOSERS Counterparty Risk Summary 12-31-2025 - Ex Trend.xlsx"
        ),
        "mosers_trend_xlsx": _sha256(
            fixtures / "MOSERS Counterparty Risk Summary 12-31-2025 - Trend.xlsx"
        ),
        "hist_all_programs_3yr_xlsx": _sha256(
            fixtures / "Historical Counterparty Risk Graphs - All Programs 3 Year.xlsx"
        ),
        "hist_ex_llc_3yr_xlsx": _sha256(
            fixtures / "Historical Counterparty Risk Graphs - ex LLC 3 Year.xlsx"
        ),
        "hist_llc_3yr_xlsx": _sha256(
            fixtures / "Historical Counterparty Risk Graphs - LLC 3 Year.xlsx"
        ),
        "monthly_pptx": _sha256(fixtures / "Monthly Counterparty Exposure Report.pptx"),
    }
    assert manifest["input_hashes"] == expected_hashes

    for variant in ("all_programs", "ex_trend", "trend"):
        assert variant in manifest["top_exposures"]
        assert variant in manifest["top_changes_per_variant"]


def test_write_risk_outputs_writes_rankings_and_top_movers(tmp_path: Path) -> None:
    warnings: list[str] = []
    parsed_by_variant = {
        "all_programs": {
            "totals": _FakeDataFrame(
                records=[
                    {
                        "counterparty": "A",
                        "Notional": 100.0,
                        "AnnualizedVolatility": 0.3,
                        "NotionalChange": 10.0,
                        "PositionUSD": 200.0,
                        "Vol": 0.2,
                        "PositionUSDChange": 20.0,
                    },
                    {
                        "counterparty": "B",
                        "Notional": 50.0,
                        "AnnualizedVolatility": 0.1,
                        "NotionalChange": 5.0,
                        "PositionUSD": 300.0,
                        "Vol": 0.1,
                        "PositionUSDChange": -10.0,
                    },
                ]
            )
        }
    }

    output_paths = run_module._write_risk_outputs(
        run_dir=tmp_path, parsed_by_variant=parsed_by_variant, warnings=warnings
    )

    assert tmp_path / "risk_rankings.csv" in output_paths
    assert tmp_path / "risk_top_movers.csv" in output_paths

    with (tmp_path / "risk_rankings.csv").open("r", encoding="utf-8", newline="") as stream:
        ranking_rows = list(csv.DictReader(stream))
    assert any(
        row["proxy_name"] == "risk_proxy_notional_annualized_volatility" and row["rank"] == "1"
        for row in ranking_rows
    )
    assert any(
        row["proxy_name"] == "risk_proxy_position_usd_vol" and row["rank"] == "1"
        for row in ranking_rows
    )

    with (tmp_path / "risk_top_movers.csv").open("r", encoding="utf-8", newline="") as stream:
        mover_rows = list(csv.DictReader(stream))
    assert {row["proxy_name"] for row in mover_rows} == {
        "risk_proxy_notional_annualized_volatility",
        "risk_proxy_position_usd_vol",
    }
    assert "risk_rankings.csv skipped" not in "\n".join(warnings)


def test_write_risk_outputs_warns_and_skips_rankings_when_proxy_columns_missing(
    tmp_path: Path,
) -> None:
    warnings: list[str] = []
    parsed_by_variant = {
        "all_programs": {
            "totals": _FakeDataFrame(records=[{"counterparty": "A", "Notional": 100.0}])
        }
    }

    output_paths = run_module._write_risk_outputs(
        run_dir=tmp_path, parsed_by_variant=parsed_by_variant, warnings=warnings
    )

    assert output_paths == []
    assert not (tmp_path / "risk_rankings.csv").exists()
    assert any(
        "requires Notional and AnnualizedVolatility columns" in warning for warning in warnings
    )
    assert any("requires PositionUSD and Vol columns" in warning for warning in warnings)
    assert any("risk_rankings.csv skipped" in warning for warning in warnings)


def test_write_risk_outputs_creates_partial_outputs_when_only_notional_proxy_exists(
    tmp_path: Path,
) -> None:
    warnings: list[str] = []
    parsed_by_variant = {
        "all_programs": {
            "totals": _FakeDataFrame(
                records=[
                    {
                        "counterparty": "A",
                        "Notional": 120.0,
                        "AnnualizedVolatility": 0.25,
                        "NotionalChange": 12.0,
                    }
                ]
            )
        }
    }

    output_paths = run_module._write_risk_outputs(
        run_dir=tmp_path, parsed_by_variant=parsed_by_variant, warnings=warnings
    )

    assert tmp_path / "risk_rankings.csv" in output_paths
    assert tmp_path / "risk_top_movers.csv" in output_paths
    with (tmp_path / "risk_rankings.csv").open("r", encoding="utf-8", newline="") as stream:
        ranking_rows = list(csv.DictReader(stream))
    assert {row["proxy_name"] for row in ranking_rows} == {
        "risk_proxy_notional_annualized_volatility"
    }
    assert any("requires PositionUSD and Vol columns" in warning for warning in warnings)


def test_run_pipeline_writes_risk_outputs_when_proxy_inputs_available(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    config_path = tmp_path / "config.yml"
    for required_file in (
        "all.xlsx",
        "ex.xlsx",
        "trend.xlsx",
        "hist-all.xlsx",
        "hist-ex.xlsx",
        "hist-trend.xlsx",
        "monthly.pptx",
    ):
        (tmp_path / required_file).write_text("fixture", encoding="utf-8")
    config_path.write_text(
        "\n".join(
            [
                "as_of_date: 2025-12-31",
                f"mosers_all_programs_xlsx: {tmp_path / 'all.xlsx'}",
                f"mosers_ex_trend_xlsx: {tmp_path / 'ex.xlsx'}",
                f"mosers_trend_xlsx: {tmp_path / 'trend.xlsx'}",
                f"hist_all_programs_3yr_xlsx: {tmp_path / 'hist-all.xlsx'}",
                f"hist_ex_llc_3yr_xlsx: {tmp_path / 'hist-ex.xlsx'}",
                f"hist_llc_3yr_xlsx: {tmp_path / 'hist-trend.xlsx'}",
                f"monthly_pptx: {tmp_path / 'monthly.pptx'}",
                f"output_root: {tmp_path / 'runs'}",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    parsed_by_variant = {
        "all_programs": {
            "totals": _FakeDataFrame(
                records=[
                    {
                        "counterparty": "A",
                        "Notional": 100.0,
                        "AnnualizedVolatility": 0.2,
                        "NotionalChange": 20.0,
                    }
                ]
            ),
            "futures": _FakeDataFrame(
                records=[
                    {
                        "account": "account-a",
                        "description": "desc",
                        "class": "class-a",
                        "fcm": "fcm-a",
                        "clearing_house": "ch-a",
                        "notional": 1.0,
                    }
                ]
            ),
        },
        "ex_trend": {
            "totals": _FakeDataFrame(
                records=[],
                columns=["counterparty", "Notional", "NotionalChange"],
            ),
            "futures": _FakeDataFrame(
                records=[],
                columns=["account", "description", "class", "fcm", "clearing_house", "notional"],
            ),
        },
        "trend": {
            "totals": _FakeDataFrame(
                records=[],
                columns=["counterparty", "Notional", "NotionalChange"],
            ),
            "futures": _FakeDataFrame(
                records=[],
                columns=["account", "description", "class", "fcm", "clearing_house", "notional"],
            ),
        },
    }

    monkeypatch.setattr("counter_risk.pipeline.run._parse_inputs", lambda _: parsed_by_variant)
    monkeypatch.setattr(
        "counter_risk.pipeline.run._run_reconciliation_checks",
        lambda *, run_dir, config, parsed_by_variant, warnings: None,
    )
    monkeypatch.setattr(
        "counter_risk.pipeline.run._update_historical_outputs",
        lambda *, run_dir, config, parsed_by_variant, as_of_date, warnings: [],
    )
    monkeypatch.setattr(
        "counter_risk.pipeline.run._write_outputs",
        lambda *, run_dir, config, as_of_date, warnings: (
            [],
            run_module.PptProcessingResult(status=run_module.PptProcessingStatus.SUCCESS),
        ),
    )

    run_dir = run_pipeline(config_path)

    assert (run_dir / "risk_rankings.csv").exists()
    assert (run_dir / "risk_top_movers.csv").exists()


def test_run_pipeline_writes_limit_breaches_csv_when_breaches_exist(
    tmp_path: Path, fake_pandas: None, monkeypatch: pytest.MonkeyPatch
) -> None:
    fixtures = Path("tests/fixtures")
    config_path = tmp_path / "config.yml"
    config_path.write_text(
        "\n".join(
            [
                "as_of_date: 2025-12-31",
                f"mosers_all_programs_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - All Programs.xlsx'}",
                f"mosers_ex_trend_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - Ex Trend.xlsx'}",
                f"mosers_trend_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - Trend.xlsx'}",
                f"hist_all_programs_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - All Programs 3 Year.xlsx'}",
                f"hist_ex_llc_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - ex LLC 3 Year.xlsx'}",
                f"hist_llc_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - LLC 3 Year.xlsx'}",
                f"monthly_pptx: {fixtures / 'Monthly Counterparty Exposure Report.pptx'}",
                f"output_root: {tmp_path / 'runs'}",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    limits_path = tmp_path / "config" / "limits.yml"
    limits_path.parent.mkdir(parents=True, exist_ok=True)
    limits_path.write_text(
        "\n".join(
            [
                "schema_version: 1",
                "limits:",
                "  - entity_type: counterparty",
                "    entity_name: citibank",
                "    limit_value: 250000000",
                "    limit_kind: absolute_notional",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    totals = _FakeDataFrame(
        records=[
            {
                "counterparty": "Citibank",
                "Notional": 300_000_000.0,
                "NotionalChange": 10_000_000.0,
            }
        ]
    )
    futures = _FakeDataFrame(
        records=[
            {
                "account": "acct-1",
                "description": "desc",
                "class": "Treasury",
                "fcm": "fcm-a",
                "clearing_house": "ch-a",
                "notional": 50_000_000.0,
            }
        ]
    )
    parsed = {
        "all_programs": {"totals": totals, "futures": futures},
        "ex_trend": {"totals": totals, "futures": futures},
        "trend": {"totals": totals, "futures": futures},
    }
    monkeypatch.setattr("counter_risk.pipeline.run._parse_inputs", lambda _: parsed)
    monkeypatch.setattr(
        "counter_risk.pipeline.run._update_historical_outputs",
        lambda *, run_dir, config, parsed_by_variant, as_of_date, warnings: [],
    )
    monkeypatch.setattr(
        "counter_risk.pipeline.run._write_outputs",
        lambda *, run_dir, config, as_of_date, warnings: (
            [],
            run_module.PptProcessingResult(status=run_module.PptProcessingStatus.SKIPPED),
        ),
    )

    run_dir = run_pipeline(config_path)

    limit_breaches_path = run_dir / "limit_breaches.csv"
    assert limit_breaches_path.exists()
    content = limit_breaches_path.read_text(encoding="utf-8")
    assert "entity_type,entity_name,limit_kind,actual_value,limit_value,breach_amount" in content
    assert "counterparty,citibank,absolute_notional" in content

    manifest = json.loads((run_dir / "manifest.json").read_text(encoding="utf-8"))
    assert "limit_breaches.csv" in manifest["output_paths"]
    assert manifest["limit_breach_summary"]["has_breaches"] is True
    assert manifest["limit_breach_summary"]["breach_count"] >= 1
    assert manifest["limit_breach_summary"]["report_path"] == "limit_breaches.csv"
    assert "limit_breaches.csv" in manifest["limit_breach_summary"]["warning_banner"]


def test_run_pipeline_warns_on_missing_limit_entities_by_default(
    tmp_path: Path, fake_pandas: None, monkeypatch: pytest.MonkeyPatch
) -> None:
    fixtures = Path("tests/fixtures")
    config_path = tmp_path / "config.yml"
    config_path.write_text(
        "\n".join(
            [
                "as_of_date: 2025-12-31",
                f"mosers_all_programs_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - All Programs.xlsx'}",
                f"mosers_ex_trend_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - Ex Trend.xlsx'}",
                f"mosers_trend_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - Trend.xlsx'}",
                f"hist_all_programs_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - All Programs 3 Year.xlsx'}",
                f"hist_ex_llc_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - ex LLC 3 Year.xlsx'}",
                f"hist_llc_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - LLC 3 Year.xlsx'}",
                f"monthly_pptx: {fixtures / 'Monthly Counterparty Exposure Report.pptx'}",
                f"output_root: {tmp_path / 'runs'}",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    limits_path = tmp_path / "config" / "limits.yml"
    limits_path.parent.mkdir(parents=True, exist_ok=True)
    limits_path.write_text(
        "\n".join(
            [
                "schema_version: 1",
                "strict_missing_entities: false",
                "limits:",
                "  - entity_type: counterparty",
                "    entity_name: not_present_counterparty",
                "    limit_value: 250000000",
                "    limit_kind: absolute_notional",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    totals = _FakeDataFrame(
        records=[
            {
                "counterparty": "Citibank",
                "Notional": 200_000_000.0,
                "NotionalChange": 10_000_000.0,
            }
        ]
    )
    futures = _FakeDataFrame(
        records=[
            {
                "account": "acct-1",
                "description": "desc",
                "class": "Treasury",
                "fcm": "fcm-a",
                "clearing_house": "ch-a",
                "notional": 50_000_000.0,
            }
        ]
    )
    parsed = {
        "all_programs": {"totals": totals, "futures": futures},
        "ex_trend": {"totals": totals, "futures": futures},
        "trend": {"totals": totals, "futures": futures},
    }
    monkeypatch.setattr("counter_risk.pipeline.run._parse_inputs", lambda _: parsed)
    monkeypatch.setattr(
        "counter_risk.pipeline.run._update_historical_outputs",
        lambda *, run_dir, config, parsed_by_variant, as_of_date, warnings: [],
    )
    monkeypatch.setattr(
        "counter_risk.pipeline.run._write_outputs",
        lambda *, run_dir, config, as_of_date, warnings: (
            [],
            run_module.PptProcessingResult(status=run_module.PptProcessingStatus.SKIPPED),
        ),
    )

    run_dir = run_pipeline(config_path)

    manifest = json.loads((run_dir / "manifest.json").read_text(encoding="utf-8"))
    assert any(
        "Limit check warning: configured limit entities missing from exposure data" in warning
        for warning in manifest["warnings"]
    )
    assert any(
        "counterparty:not_present_counterparty" in warning for warning in manifest["warnings"]
    )


def test_run_pipeline_strict_missing_limit_entities_fails(
    tmp_path: Path, fake_pandas: None, monkeypatch: pytest.MonkeyPatch
) -> None:
    fixtures = Path("tests/fixtures")
    config_path = tmp_path / "config.yml"
    config_path.write_text(
        "\n".join(
            [
                "as_of_date: 2025-12-31",
                f"mosers_all_programs_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - All Programs.xlsx'}",
                f"mosers_ex_trend_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - Ex Trend.xlsx'}",
                f"mosers_trend_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - Trend.xlsx'}",
                f"hist_all_programs_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - All Programs 3 Year.xlsx'}",
                f"hist_ex_llc_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - ex LLC 3 Year.xlsx'}",
                f"hist_llc_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - LLC 3 Year.xlsx'}",
                f"monthly_pptx: {fixtures / 'Monthly Counterparty Exposure Report.pptx'}",
                f"output_root: {tmp_path / 'runs'}",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    limits_path = tmp_path / "config" / "limits.yml"
    limits_path.parent.mkdir(parents=True, exist_ok=True)
    limits_path.write_text(
        "\n".join(
            [
                "schema_version: 1",
                "strict_missing_entities: true",
                "limits:",
                "  - entity_type: counterparty",
                "    entity_name: not_present_counterparty",
                "    limit_value: 250000000",
                "    limit_kind: absolute_notional",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    totals = _FakeDataFrame(
        records=[
            {
                "counterparty": "Citibank",
                "Notional": 200_000_000.0,
                "NotionalChange": 10_000_000.0,
            }
        ]
    )
    futures = _FakeDataFrame(
        records=[
            {
                "account": "acct-1",
                "description": "desc",
                "class": "Treasury",
                "fcm": "fcm-a",
                "clearing_house": "ch-a",
                "notional": 50_000_000.0,
            }
        ]
    )
    parsed = {
        "all_programs": {"totals": totals, "futures": futures},
        "ex_trend": {"totals": totals, "futures": futures},
        "trend": {"totals": totals, "futures": futures},
    }
    monkeypatch.setattr("counter_risk.pipeline.run._parse_inputs", lambda _: parsed)
    monkeypatch.setattr(
        "counter_risk.pipeline.run._update_historical_outputs",
        lambda *, run_dir, config, parsed_by_variant, as_of_date, warnings: [],
    )
    monkeypatch.setattr(
        "counter_risk.pipeline.run._write_outputs",
        lambda *, run_dir, config, as_of_date, warnings: (
            [],
            run_module.PptProcessingResult(status=run_module.PptProcessingStatus.SKIPPED),
        ),
    )

    with pytest.raises(RuntimeError, match="Pipeline failed during limit breach stage"):
        run_pipeline(config_path)


def test_run_pipeline_generates_all_programs_mosers_from_raw_nisa_input(
    tmp_path: Path, fake_pandas: None, monkeypatch: pytest.MonkeyPatch
) -> None:
    fixtures = Path("tests/fixtures")
    config_path = tmp_path / "config.yml"
    config_path.write_text(
        "\n".join(
            [
                "as_of_date: 2025-12-31",
                f"raw_nisa_all_programs_xlsx: {fixtures / 'NISA Monthly All Programs - Raw.xlsx'}",
                f"mosers_ex_trend_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - Ex Trend.xlsx'}",
                f"mosers_trend_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - Trend.xlsx'}",
                f"hist_all_programs_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - All Programs 3 Year.xlsx'}",
                f"hist_ex_llc_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - ex LLC 3 Year.xlsx'}",
                f"hist_llc_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - LLC 3 Year.xlsx'}",
                f"monthly_pptx: {fixtures / 'Monthly Counterparty Exposure Report.pptx'}",
                f"output_root: {tmp_path / 'runs'}",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    # Keep the raw-NISA generation path real, but stub downstream heavy stages
    # that are covered by dedicated integration tests.
    monkeypatch.setattr(
        "counter_risk.pipeline.run._parse_inputs", lambda _: _minimal_parsed_by_variant()
    )
    monkeypatch.setattr(
        "counter_risk.pipeline.run._update_historical_outputs",
        lambda *, run_dir, config, parsed_by_variant, as_of_date, warnings: [],
    )
    monkeypatch.setattr(
        "counter_risk.pipeline.run._write_outputs",
        lambda *, run_dir, config, as_of_date, warnings: (
            [],
            run_module.PptProcessingResult(status=run_module.PptProcessingStatus.SUCCESS),
        ),
    )

    run_dir = run_pipeline(config_path)

    generated_mosers_output = run_dir / "all_programs-mosers-input.xlsx"
    intermediate_generated_output = run_dir / "_generated" / "all_programs-generated-mosers.xlsx"
    assert generated_mosers_output.exists()
    assert intermediate_generated_output.exists()

    from openpyxl import load_workbook  # type: ignore[import-untyped]

    workbook = load_workbook(generated_mosers_output, read_only=True, data_only=True)
    try:
        assert workbook.sheetnames == ["CPRS - CH", "CPRS - FCM"]
    finally:
        workbook.close()

    manifest = json.loads((run_dir / "manifest.json").read_text(encoding="utf-8"))
    assert "raw_nisa_all_programs_xlsx" in manifest["input_hashes"]
    assert "Generated All Programs MOSERS workbook from raw NISA input" in manifest["warnings"]


def test_prepare_runtime_config_generates_and_copies_raw_nisa_mosers_output(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    run_dir = tmp_path / "run"
    raw_nisa = tmp_path / "raw.xlsx"
    raw_nisa.write_text("raw", encoding="utf-8")
    ex_trend = tmp_path / "ex.xlsx"
    ex_trend.write_text("ex", encoding="utf-8")
    trend = tmp_path / "trend.xlsx"
    trend.write_text("trend", encoding="utf-8")
    hist_all = tmp_path / "hist_all.xlsx"
    hist_all.write_text("hist_all", encoding="utf-8")
    hist_ex = tmp_path / "hist_ex.xlsx"
    hist_ex.write_text("hist_ex", encoding="utf-8")
    hist_trend = tmp_path / "hist_trend.xlsx"
    hist_trend.write_text("hist_trend", encoding="utf-8")
    monthly_pptx = tmp_path / "monthly.pptx"
    monthly_pptx.write_text("ppt", encoding="utf-8")

    generated_calls: list[dict[str, Any]] = []

    def _fake_generate_mosers_workbook(
        *, raw_nisa_path: Path, output_path: Path, as_of_date: date
    ) -> Path:
        generated_calls.append(
            {
                "raw_nisa_path": raw_nisa_path,
                "output_path": output_path,
                "as_of_date": as_of_date,
            }
        )
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_bytes(b"generated-workbook")
        return output_path

    monkeypatch.setattr(run_module, "generate_mosers_workbook", _fake_generate_mosers_workbook)

    config = WorkflowConfig(
        as_of_date=date(2025, 12, 31),
        raw_nisa_all_programs_xlsx=raw_nisa,
        mosers_all_programs_xlsx=None,
        mosers_ex_trend_xlsx=ex_trend,
        mosers_trend_xlsx=trend,
        hist_all_programs_3yr_xlsx=hist_all,
        hist_ex_llc_3yr_xlsx=hist_ex,
        hist_llc_3yr_xlsx=hist_trend,
        monthly_pptx=monthly_pptx,
        output_root=tmp_path / "runs",
    )
    warnings: list[str] = []

    runtime_config = run_module._prepare_runtime_config(
        config=config,
        run_dir=run_dir,
        as_of_date=date(2025, 12, 31),
        warnings=warnings,
    )

    generated_path = run_dir / "_generated" / "all_programs-generated-mosers.xlsx"
    canonical_path = run_dir / "all_programs-mosers-input.xlsx"

    assert generated_calls == [
        {
            "raw_nisa_path": raw_nisa,
            "output_path": generated_path,
            "as_of_date": date(2025, 12, 31),
        }
    ]
    assert generated_path.exists()
    assert canonical_path.exists()
    assert generated_path.read_bytes() == canonical_path.read_bytes()
    assert runtime_config.mosers_all_programs_xlsx == canonical_path
    assert "Generated All Programs MOSERS workbook from raw NISA input" in warnings


def _write_valid_config(tmp_path: Path, output_root: Path) -> Path:
    fixtures = Path("tests/fixtures")
    config_path = tmp_path / "config.yml"
    config_path.write_text(
        "\n".join(
            [
                "as_of_date: 2025-12-31",
                f"mosers_all_programs_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - All Programs.xlsx'}",
                f"mosers_ex_trend_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - Ex Trend.xlsx'}",
                f"mosers_trend_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - Trend.xlsx'}",
                f"hist_all_programs_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - All Programs 3 Year.xlsx'}",
                f"hist_ex_llc_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - ex LLC 3 Year.xlsx'}",
                f"hist_llc_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - LLC 3 Year.xlsx'}",
                f"monthly_pptx: {fixtures / 'Monthly Counterparty Exposure Report.pptx'}",
                f"output_root: {output_root}",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    return config_path


def test_run_pipeline_wraps_parse_errors(
    tmp_path: Path, fake_pandas: None, monkeypatch: pytest.MonkeyPatch
) -> None:
    config_path = _write_valid_config(tmp_path=tmp_path, output_root=tmp_path / "runs")

    def _boom(_: dict[str, Path]) -> dict[str, dict[str, Any]]:
        raise ValueError("bad parser input")

    monkeypatch.setattr("counter_risk.pipeline.run._parse_inputs", _boom)

    with pytest.raises(RuntimeError, match="Pipeline failed during parse stage") as exc_info:
        run_pipeline(config_path)

    assert isinstance(exc_info.value.__cause__, ValueError)
    assert "bad parser input" in str(exc_info.value.__cause__)


def test_run_pipeline_fails_when_as_of_date_is_not_derivable(
    tmp_path: Path, fake_pandas: None, monkeypatch: pytest.MonkeyPatch
) -> None:
    fixtures = Path("tests/fixtures")
    config_path = tmp_path / "config.yml"
    config_path.write_text(
        "\n".join(
            [
                f"mosers_all_programs_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - All Programs.xlsx'}",
                f"mosers_ex_trend_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - Ex Trend.xlsx'}",
                f"mosers_trend_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - Trend.xlsx'}",
                f"hist_all_programs_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - All Programs 3 Year.xlsx'}",
                f"hist_ex_llc_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - ex LLC 3 Year.xlsx'}",
                f"hist_llc_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - LLC 3 Year.xlsx'}",
                f"monthly_pptx: {fixtures / 'Monthly Counterparty Exposure Report.pptx'}",
                f"output_root: {tmp_path / 'runs'}",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    monkeypatch.setattr(
        "counter_risk.pipeline.run._collect_cprs_header_candidates", lambda config: []
    )

    with pytest.raises(
        RuntimeError, match="Pipeline failed during date derivation stage"
    ) as exc_info:
        run_pipeline(config_path)

    assert isinstance(exc_info.value.__cause__, ValueError)
    assert "Unable to derive as_of_date" in str(exc_info.value.__cause__)


def test_run_pipeline_wraps_parse_validation_errors(
    tmp_path: Path, fake_pandas: None, monkeypatch: pytest.MonkeyPatch
) -> None:
    config_path = _write_valid_config(tmp_path=tmp_path, output_root=tmp_path / "runs")

    malformed = {
        "all_programs": {
            "totals": _FakeDataFrame(records=[{"counterparty": "A", "Notional": 1.0}]),
            "futures": _FakeDataFrame(records=[]),
        },
        "ex_trend": {
            "totals": _FakeDataFrame(records=[{"counterparty": "B", "Notional": 2.0}]),
            "futures": _FakeDataFrame(records=[]),
        },
        "trend": {
            "totals": _FakeDataFrame(records=[]),
            "futures": _FakeDataFrame(records=[{"account": "Acct"}]),
        },
    }

    def _bad_parse(_: dict[str, Path]) -> dict[str, dict[str, Any]]:
        return malformed

    monkeypatch.setattr("counter_risk.pipeline.run._parse_inputs", _bad_parse)

    with pytest.raises(RuntimeError, match="Pipeline failed during parse stage") as exc_info:
        run_pipeline(config_path)

    assert isinstance(exc_info.value.__cause__, ValueError)
    assert "missing required columns" in str(exc_info.value.__cause__)


def test_run_pipeline_warn_mode_writes_mapping_updates_and_completes(
    tmp_path: Path, fake_pandas: None, monkeypatch: pytest.MonkeyPatch
) -> None:
    config_path = _write_valid_config(tmp_path=tmp_path, output_root=tmp_path / "runs")
    config_path.write_text(
        config_path.read_text(encoding="utf-8")
        + "\n".join(
            [
                "reconciliation:",
                "  fail_policy: warn",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    monkeypatch.setattr(
        "counter_risk.pipeline.run._parse_inputs", lambda _: _minimal_parsed_by_variant()
    )
    monkeypatch.setattr(
        "counter_risk.pipeline.run._extract_historical_series_headers_by_sheet",
        lambda _: {"Total": ("Legacy Counterparty",)},
    )
    monkeypatch.setattr(
        "counter_risk.pipeline.run._update_historical_outputs",
        lambda *, run_dir, config, parsed_by_variant, as_of_date, warnings: [],
    )
    monkeypatch.setattr(
        "counter_risk.pipeline.run._write_outputs",
        lambda *, run_dir, config, as_of_date, warnings: (
            [],
            run_module.PptProcessingResult(status=run_module.PptProcessingStatus.SUCCESS),
        ),
    )

    run_dir = run_pipeline(config_path)
    manifest = json.loads((run_dir / "manifest.json").read_text(encoding="utf-8"))
    mapping_updates = run_dir / "NEEDS_MAPPING_UPDATES.txt"

    assert mapping_updates.exists()
    text = mapping_updates.read_text(encoding="utf-8")
    assert "run_identifier: 2025-12-31" in text
    assert "fail_policy: warn" in text
    assert "missing_from_historical_headers" in text
    assert any("Reconciliation summary:" in warning for warning in manifest["warnings"])


def test_run_pipeline_strict_mode_fails_when_reconciliation_has_gaps(
    tmp_path: Path, fake_pandas: None, monkeypatch: pytest.MonkeyPatch
) -> None:
    config_path = _write_valid_config(tmp_path=tmp_path, output_root=tmp_path / "runs")
    config_path.write_text(
        config_path.read_text(encoding="utf-8")
        + "\n".join(
            [
                "reconciliation:",
                "  fail_policy: strict",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    monkeypatch.setattr(
        "counter_risk.pipeline.run._parse_inputs", lambda _: _minimal_parsed_by_variant()
    )
    monkeypatch.setattr(
        "counter_risk.pipeline.run._extract_historical_series_headers_by_sheet",
        lambda _: {"Total": ("Legacy Counterparty",)},
    )

    with pytest.raises(RuntimeError, match="Pipeline failed during parse stage") as exc_info:
        run_pipeline(config_path)

    assert isinstance(exc_info.value.__cause__, UnmappedCounterpartyError)
    assert "Unmapped normalized counterparty" in str(exc_info.value.__cause__)


def test_run_pipeline_wraps_output_write_errors(
    tmp_path: Path, fake_pandas: None, monkeypatch: pytest.MonkeyPatch
) -> None:
    config_path = _write_valid_config(tmp_path=tmp_path, output_root=tmp_path / "runs")
    monkeypatch.setattr(
        "counter_risk.pipeline.run._parse_inputs", lambda _: _minimal_parsed_by_variant()
    )
    monkeypatch.setattr(
        "counter_risk.pipeline.run._update_historical_outputs",
        lambda *, run_dir, config, parsed_by_variant, as_of_date, warnings: [],
    )

    def _boom(*, run_dir: Path, config: Any, as_of_date: date, warnings: list[str]) -> list[Path]:
        _ = (run_dir, config, as_of_date, warnings)
        raise OSError("disk full")

    monkeypatch.setattr("counter_risk.pipeline.run._write_outputs", _boom)

    with pytest.raises(RuntimeError, match="Pipeline failed during output write stage") as exc_info:
        run_pipeline(config_path)

    assert isinstance(exc_info.value.__cause__, OSError)
    assert "disk full" in str(exc_info.value.__cause__)


def test_run_pipeline_wraps_input_validation_errors(
    tmp_path: Path, fake_pandas: None, monkeypatch: pytest.MonkeyPatch
) -> None:
    config_path = _write_valid_config(tmp_path=tmp_path, output_root=tmp_path / "runs")

    def _boom(_: dict[str, Path]) -> None:
        raise FileNotFoundError("missing source workbook")

    monkeypatch.setattr("counter_risk.pipeline.run._validate_input_files", _boom)

    with pytest.raises(
        RuntimeError, match="Pipeline failed during input validation stage"
    ) as exc_info:
        run_pipeline(config_path)

    assert isinstance(exc_info.value.__cause__, FileNotFoundError)
    assert "missing source workbook" in str(exc_info.value.__cause__)


def test_run_pipeline_wraps_compute_errors(
    tmp_path: Path, fake_pandas: None, monkeypatch: pytest.MonkeyPatch
) -> None:
    config_path = _write_valid_config(tmp_path=tmp_path, output_root=tmp_path / "runs")
    monkeypatch.setattr(
        "counter_risk.pipeline.run._parse_inputs", lambda _: _minimal_parsed_by_variant()
    )

    def _boom(
        _: dict[str, dict[str, Any]],
    ) -> tuple[dict[str, list[dict[str, Any]]], dict[str, list[dict[str, Any]]]]:
        raise ValueError("bad compute inputs")

    monkeypatch.setattr("counter_risk.pipeline.run._compute_metrics", _boom)

    with pytest.raises(RuntimeError, match="Pipeline failed during compute stage") as exc_info:
        run_pipeline(config_path)

    assert isinstance(exc_info.value.__cause__, ValueError)
    assert "bad compute inputs" in str(exc_info.value.__cause__)


def test_run_pipeline_wraps_historical_update_errors(
    tmp_path: Path, fake_pandas: None, monkeypatch: pytest.MonkeyPatch
) -> None:
    config_path = _write_valid_config(tmp_path=tmp_path, output_root=tmp_path / "runs")
    monkeypatch.setattr(
        "counter_risk.pipeline.run._parse_inputs", lambda _: _minimal_parsed_by_variant()
    )

    def _boom(
        *,
        run_dir: Path,
        config: Any,
        parsed_by_variant: dict[str, dict[str, Any]],
        as_of_date: date,
        warnings: list[str],
    ) -> list[Path]:
        _ = (run_dir, config, parsed_by_variant, as_of_date, warnings)
        raise OSError("historical workbook write failed")

    monkeypatch.setattr("counter_risk.pipeline.run._update_historical_outputs", _boom)

    with pytest.raises(
        RuntimeError, match="Pipeline failed during historical update stage"
    ) as exc_info:
        run_pipeline(config_path)

    assert isinstance(exc_info.value.__cause__, OSError)
    assert "historical workbook write failed" in str(exc_info.value.__cause__)


def test_run_pipeline_wraps_manifest_generation_errors(
    tmp_path: Path, fake_pandas: None, monkeypatch: pytest.MonkeyPatch
) -> None:
    config_path = _write_valid_config(tmp_path=tmp_path, output_root=tmp_path / "runs")
    monkeypatch.setattr(
        "counter_risk.pipeline.run._parse_inputs", lambda _: _minimal_parsed_by_variant()
    )
    monkeypatch.setattr(
        "counter_risk.pipeline.run._update_historical_outputs",
        lambda *, run_dir, config, parsed_by_variant, as_of_date, warnings: [],
    )
    monkeypatch.setattr(
        "counter_risk.pipeline.run._write_outputs",
        lambda *, run_dir, config, as_of_date, warnings: (
            [],
            run_module.PptProcessingResult(status=run_module.PptProcessingStatus.SUCCESS),
        ),
    )

    def _boom(self: Any, *, run_dir: Path, manifest: dict[str, Any]) -> Path:
        _ = (self, run_dir, manifest)
        raise OSError("manifest disk error")

    monkeypatch.setattr("counter_risk.pipeline.run.ManifestBuilder.write", _boom)

    with pytest.raises(
        RuntimeError, match="Pipeline failed during manifest generation stage"
    ) as exc_info:
        run_pipeline(config_path)

    assert isinstance(exc_info.value.__cause__, OSError)
    assert "manifest disk error" in str(exc_info.value.__cause__)


def test_run_pipeline_passes_as_of_date_and_parsed_inputs_to_historical_update(
    tmp_path: Path, fake_pandas: None, monkeypatch: pytest.MonkeyPatch
) -> None:
    config_path = _write_valid_config(tmp_path=tmp_path, output_root=tmp_path / "runs")
    monkeypatch.setattr(
        "counter_risk.pipeline.run._parse_inputs", lambda _: _minimal_parsed_by_variant()
    )
    calls: list[dict[str, Any]] = []

    def _capture(
        *,
        run_dir: Path,
        config: Any,
        parsed_by_variant: dict[str, dict[str, Any]],
        as_of_date: date,
        warnings: list[str],
    ) -> list[Path]:
        _ = config
        calls.append(
            {
                "run_dir": run_dir,
                "variants": sorted(parsed_by_variant.keys()),
                "as_of_date": as_of_date,
                "warnings": warnings,
            }
        )
        return []

    monkeypatch.setattr("counter_risk.pipeline.run._update_historical_outputs", _capture)

    run_dir = run_pipeline(config_path)

    assert run_dir == tmp_path / "runs" / "2025-12-31"
    assert len(calls) == 1
    assert calls[0]["as_of_date"] == date(2025, 12, 31)
    assert calls[0]["variants"] == ["all_programs", "ex_trend", "trend"]


def test_run_pipeline_invokes_ppt_link_refresh(
    tmp_path: Path, fake_pandas: None, monkeypatch: pytest.MonkeyPatch
) -> None:
    config_path = _write_valid_config(tmp_path=tmp_path, output_root=tmp_path / "runs")
    monkeypatch.setattr(
        "counter_risk.pipeline.run._parse_inputs", lambda _: _minimal_parsed_by_variant()
    )
    monkeypatch.setattr(
        "counter_risk.pipeline.run._update_historical_outputs",
        lambda *, run_dir, config, parsed_by_variant, as_of_date, warnings: [],
    )
    seen: dict[str, Path] = {}

    def _refresh(pptx_path: Path) -> bool:
        seen["path"] = pptx_path
        return True

    monkeypatch.setattr("counter_risk.pipeline.run._refresh_ppt_links", _refresh)

    run_dir = run_pipeline(config_path)
    manifest = json.loads((run_dir / "manifest.json").read_text(encoding="utf-8"))

    assert (
        seen["path"] == run_dir / "Monthly Counterparty Exposure Report (Master) - 2025-12-31.pptx"
    )
    assert "PPT links not refreshed; COM refresh skipped" not in manifest["warnings"]


def test_run_pipeline_ignores_config_output_root_for_run_directory(
    tmp_path: Path, fake_pandas: None, monkeypatch: pytest.MonkeyPatch
) -> None:
    config_path = _write_valid_config(
        tmp_path=tmp_path, output_root=tmp_path / "different-output-root"
    )
    monkeypatch.setattr(
        "counter_risk.pipeline.run._parse_inputs", lambda _: _minimal_parsed_by_variant()
    )
    monkeypatch.setattr(
        "counter_risk.pipeline.run._update_historical_outputs",
        lambda *, run_dir, config, parsed_by_variant, as_of_date, warnings: [],
    )
    monkeypatch.setattr(
        "counter_risk.pipeline.run._write_outputs",
        lambda *, run_dir, config, as_of_date, warnings: (
            [],
            run_module.PptProcessingResult(status=run_module.PptProcessingStatus.SUCCESS),
        ),
    )

    run_dir = run_pipeline(config_path)

    assert run_dir == tmp_path / "runs" / "2025-12-31"
    assert not (tmp_path / "different-output-root" / "2025-12-31").exists()


def test_run_pipeline_wraps_config_validation_errors_for_output_root_file(
    tmp_path: Path, fake_pandas: None
) -> None:
    output_root_file = tmp_path / "runs"
    output_root_file.write_text("not-a-directory", encoding="utf-8")
    config_path = _write_valid_config(tmp_path=tmp_path, output_root=output_root_file)

    with pytest.raises(RuntimeError, match="Pipeline failed during config validation") as exc_info:
        run_pipeline(config_path)

    assert isinstance(exc_info.value.__cause__, ValueError)
    assert "output_root must be a directory path" in str(exc_info.value.__cause__)


def test_run_pipeline_wraps_config_validation_errors_for_invalid_ppt_extension(
    tmp_path: Path, fake_pandas: None
) -> None:
    fixtures = Path("tests/fixtures")
    output_root = tmp_path / "runs"
    config_path = tmp_path / "config.yml"
    config_path.write_text(
        "\n".join(
            [
                "as_of_date: 2025-12-31",
                f"mosers_all_programs_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - All Programs.xlsx'}",
                f"mosers_ex_trend_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - Ex Trend.xlsx'}",
                f"mosers_trend_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - Trend.xlsx'}",
                f"hist_all_programs_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - All Programs 3 Year.xlsx'}",
                f"hist_ex_llc_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - ex LLC 3 Year.xlsx'}",
                f"hist_llc_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - LLC 3 Year.xlsx'}",
                f"monthly_pptx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - All Programs.xlsx'}",
                f"output_root: {output_root}",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    with pytest.raises(RuntimeError, match="Pipeline failed during config validation") as exc_info:
        run_pipeline(config_path)

    assert isinstance(exc_info.value.__cause__, ValueError)
    assert "Invalid file type for monthly_pptx: expected .pptx" in str(exc_info.value.__cause__)


def test_run_pipeline_wraps_config_load_errors(
    tmp_path: Path, fake_pandas: None, monkeypatch: pytest.MonkeyPatch
) -> None:
    config_path = tmp_path / "missing.yml"

    def _boom(_: Path) -> Any:
        raise ValueError("config parse failed")

    monkeypatch.setattr("counter_risk.pipeline.run.load_config", _boom)

    with pytest.raises(RuntimeError, match="Pipeline failed during config load") as exc_info:
        run_pipeline(config_path)

    assert isinstance(exc_info.value.__cause__, ValueError)
    assert "config parse failed" in str(exc_info.value.__cause__)


def test_run_pipeline_wraps_run_directory_setup_errors(
    tmp_path: Path, fake_pandas: None, monkeypatch: pytest.MonkeyPatch
) -> None:
    config_path = _write_valid_config(tmp_path=tmp_path, output_root=tmp_path / "runs")

    def _boom(self: Path, parents: bool, exist_ok: bool) -> None:
        _ = (self, parents, exist_ok)
        raise OSError("permission denied")

    monkeypatch.setattr("counter_risk.pipeline.run.Path.mkdir", _boom)

    with pytest.raises(
        RuntimeError, match="Pipeline failed during run directory setup stage"
    ) as exc_info:
        run_pipeline(config_path)

    assert isinstance(exc_info.value.__cause__, OSError)
    assert "permission denied" in str(exc_info.value.__cause__)


def test_merge_historical_workbook_prefers_configured_total_sheet(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    workbook_path = tmp_path / "hist.xlsx"
    workbook_path.write_bytes(b"fixture")

    decoy = _FakeWorksheet("A Decoy")
    decoy.set_value(1, 1, "Date")
    decoy.set_value(1, 2, "Wrong")
    decoy.set_value(1, 3, "Wrong")
    decoy.set_value(2, 1, "2025-12-31")

    target = _FakeWorksheet("Total")
    target.set_value(1, 1, "Date")
    target.set_value(1, 2, "Barclays")
    target.set_value(1, 3, "Citibank")
    target.set_value(2, 1, "2025-12-31")

    workbook = _FakeWorkbook({"A Decoy": decoy, "Total": target})
    monkeypatch.setitem(
        sys.modules, "openpyxl", types.SimpleNamespace(load_workbook=lambda filename: workbook)
    )

    run_module._merge_historical_workbook(
        workbook_path=workbook_path,
        variant="all_programs",
        as_of_date=date(2026, 2, 13),
        totals_records=[{"Notional": 10.0, "counterparty": "A"}],
        warnings=[],
    )

    assert target.cell(row=3, column=1).value == date(2026, 2, 13)
    assert target.cell(row=3, column=2).value == pytest.approx(10.0)
    assert target.cell(row=3, column=3).value == 1
    assert decoy.cell(row=3, column=1).value is None


def test_merge_historical_workbook_uses_deterministic_fallback_sheet_when_preferred_missing(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    workbook_path = tmp_path / "hist.xlsx"
    workbook_path.write_bytes(b"fixture")

    alpha = _FakeWorksheet("Alpha")
    alpha.set_value(1, 1, "Date")
    alpha.set_value(1, 2, "Series A")
    alpha.set_value(1, 3, "Series B")
    alpha.set_value(2, 1, "2025-12-31")

    zulu = _FakeWorksheet("Zulu")
    zulu.set_value(1, 1, "Date")
    zulu.set_value(1, 2, "Series A")
    zulu.set_value(1, 3, "Series B")
    zulu.set_value(2, 1, "2025-12-31")

    workbook = _FakeWorkbook({"Zulu": zulu, "Alpha": alpha})
    monkeypatch.setitem(
        sys.modules, "openpyxl", types.SimpleNamespace(load_workbook=lambda filename: workbook)
    )

    run_module._merge_historical_workbook(
        workbook_path=workbook_path,
        variant="unknown_variant",
        as_of_date=date(2026, 2, 13),
        totals_records=[
            {"Notional": 20.0, "counterparty": "A"},
            {"Notional": 5.0, "counterparty": "B"},
        ],
        warnings=[],
    )

    assert alpha.cell(row=3, column=1).value == date(2026, 2, 13)
    assert alpha.cell(row=3, column=2).value == pytest.approx(25.0)
    assert alpha.cell(row=3, column=3).value == 2
    assert zulu.cell(row=3, column=1).value is None


def test_merge_historical_workbook_fails_fast_when_required_headers_missing(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    workbook_path = tmp_path / "hist.xlsx"
    workbook_path.write_bytes(b"fixture")

    broken = _FakeWorksheet("Total")
    broken.set_value(1, 1, "Date")
    broken.set_value(1, 2, "")
    broken.set_value(1, 3, "Series B")
    broken.set_value(2, 1, "2025-12-31")

    workbook = _FakeWorkbook({"Total": broken})
    monkeypatch.setitem(
        sys.modules, "openpyxl", types.SimpleNamespace(load_workbook=lambda filename: workbook)
    )

    with pytest.raises(RuntimeError, match="Failed to update historical workbook") as exc_info:
        run_module._merge_historical_workbook(
            workbook_path=workbook_path,
            variant="all_programs",
            as_of_date=date(2026, 2, 13),
            totals_records=[{"Notional": 10.0, "counterparty": "A"}],
            warnings=[],
        )

    message = str(exc_info.value.__cause__)
    assert "missing required columns" in message
    assert "Total" in message
    assert "value series 1" in message
    assert broken.cell(row=3, column=1).value is None


# ---------------------------------------------------------------------------
# Static distribution fallback tests
# ---------------------------------------------------------------------------


def _make_minimal_config(tmp_path: Path, *, distribution_static: bool = False) -> WorkflowConfig:
    """Return a minimal WorkflowConfig with placeholder files."""
    tmp_path.mkdir(parents=True, exist_ok=True)
    for name in (
        "all.xlsx",
        "ex.xlsx",
        "trend.xlsx",
        "hist_all.xlsx",
        "hist_ex.xlsx",
        "hist_llc.xlsx",
        "monthly.pptx",
    ):
        (tmp_path / name).write_bytes(b"placeholder")
    return WorkflowConfig(
        as_of_date=date(2025, 12, 31),
        mosers_all_programs_xlsx=tmp_path / "all.xlsx",
        mosers_ex_trend_xlsx=tmp_path / "ex.xlsx",
        mosers_trend_xlsx=tmp_path / "trend.xlsx",
        hist_all_programs_3yr_xlsx=tmp_path / "hist_all.xlsx",
        hist_ex_llc_3yr_xlsx=tmp_path / "hist_ex.xlsx",
        hist_llc_3yr_xlsx=tmp_path / "hist_llc.xlsx",
        monthly_pptx=tmp_path / "monthly.pptx",
        distribution_static=distribution_static,
        output_root=tmp_path / "runs",
    )


def test_create_static_distribution_is_noop_when_flag_is_false(
    tmp_path: Path,
) -> None:
    """When distribution_static=False the function returns no paths and no warnings."""
    source_pptx = tmp_path / "deck.pptx"
    source_pptx.write_bytes(b"fake-pptx")
    run_dir = tmp_path / "run"
    run_dir.mkdir()
    config = _make_minimal_config(tmp_path / "cfg", distribution_static=False)
    warnings: list[str] = []

    output = run_module._create_static_distribution(
        source_pptx=source_pptx,
        run_dir=run_dir,
        config=config,
        warnings=warnings,
    )

    assert output == []
    assert warnings == []


def test_create_static_distribution_warns_on_non_windows(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """On non-Windows platforms a clear warning is emitted and no paths are returned."""
    monkeypatch.setattr("counter_risk.pipeline.run.platform.system", lambda: "Linux")

    source_pptx = tmp_path / "deck.pptx"
    source_pptx.write_bytes(b"fake-pptx")
    run_dir = tmp_path / "run"
    run_dir.mkdir()
    config = _make_minimal_config(tmp_path / "cfg", distribution_static=True)
    warnings: list[str] = []

    output = run_module._create_static_distribution(
        source_pptx=source_pptx,
        run_dir=run_dir,
        config=config,
        warnings=warnings,
    )

    assert output == []
    assert len(warnings) == 1
    assert "distribution_static" in warnings[0]
    assert "Windows" in warnings[0]


def test_create_static_distribution_warns_when_win32com_missing(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """When win32com is absent on a simulated Windows host a warning is emitted."""
    monkeypatch.setattr("counter_risk.pipeline.run.platform.system", lambda: "Windows")
    # Ensure win32com.client cannot be imported.
    monkeypatch.setitem(sys.modules, "win32com", None)
    monkeypatch.setitem(sys.modules, "win32com.client", None)

    source_pptx = tmp_path / "deck.pptx"
    source_pptx.write_bytes(b"fake-pptx")
    run_dir = tmp_path / "run"
    run_dir.mkdir()
    config = _make_minimal_config(tmp_path / "cfg", distribution_static=True)
    warnings: list[str] = []

    output = run_module._create_static_distribution(
        source_pptx=source_pptx,
        run_dir=run_dir,
        config=config,
        warnings=warnings,
    )

    assert output == []
    assert len(warnings) == 1
    assert "distribution_static" in warnings[0]
    assert "win32com" in warnings[0]


def test_create_static_distribution_rebuilds_from_slide_images(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """On Windows with COM, static output is rebuilt from one exported image per slide."""
    from pptx import Presentation

    monkeypatch.setattr(platform, "system", lambda: "Windows")

    source_pptx = tmp_path / "source.pptx"
    source_prs = Presentation()
    blank_layout = source_prs.slide_layouts[6]
    source_prs.slides.add_slide(blank_layout)
    source_prs.slides.add_slide(blank_layout)
    source_prs.save(str(source_pptx))

    run_dir = tmp_path / "run"
    run_dir.mkdir()
    config = _make_minimal_config(tmp_path / "cfg", distribution_static=True)
    warnings: list[str] = []

    png_bytes = (
        b"\x89PNG\r\n\x1a\n"
        b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x06\x00\x00\x00\x1f\x15\xc4\x89"
        b"\x00\x00\x00\rIDATx\x9cc\xf8\xcf\xc0\xf0\x1f\x00\x05\x00\x01\xff\x89\x99=\x1d"
        b"\x00\x00\x00\x00IEND\xaeB`\x82"
    )

    class _FakeSlide:
        def __init__(self, slide_idx: int) -> None:
            self._slide_idx = slide_idx

        def Export(self, path: str, fmt: str) -> None:  # noqa: N802
            assert fmt == "PNG"
            Path(path).write_bytes(png_bytes)

    class _FakeSlides:
        def __init__(self, count: int) -> None:
            self.Count = count
            self._slides = {idx: _FakeSlide(idx) for idx in range(1, count + 1)}

        def __getitem__(self, idx: int) -> _FakeSlide:
            return self._slides[idx]

    class _FakePresentation:
        def __init__(self) -> None:
            self.Slides = _FakeSlides(2)

        def ExportAsFixedFormat(self, path: str, fmt: int) -> None:  # noqa: N802
            assert fmt == 2
            Path(path).write_bytes(b"%PDF-1.4\n")

        def Close(self) -> None:  # noqa: N802
            return None

    class _FakePowerPointApplication:
        def __init__(self) -> None:
            self.Visible = False
            self.Presentations = types.SimpleNamespace(
                Open=lambda *_args, **_kwargs: _FakePresentation()
            )

        def Quit(self) -> None:  # noqa: N802
            return None

    fake_client = types.SimpleNamespace(
        DispatchEx=lambda *_args, **_kwargs: _FakePowerPointApplication()
    )
    fake_win32com = types.ModuleType("win32com")
    cast(Any, fake_win32com).client = fake_client
    monkeypatch.setitem(sys.modules, "win32com", fake_win32com)
    monkeypatch.setitem(sys.modules, "win32com.client", fake_client)

    output = run_module._create_static_distribution(
        source_pptx=source_pptx,
        run_dir=run_dir,
        config=config,
        warnings=warnings,
    )

    static_paths = [path for path in output if path.suffix == ".pptx"]
    assert len(static_paths) == 1
    assert warnings == []

    output_prs = Presentation(str(static_paths[0]))
    assert len(output_prs.slides) == 2


def test_run_pipeline_manifest_includes_distribution_static_warning(
    tmp_path: Path, fake_pandas: None, monkeypatch: pytest.MonkeyPatch
) -> None:
    """Full pipeline: distribution_static=True on non-Windows emits manifest warning."""
    fixtures = Path("tests/fixtures")
    output_root = tmp_path / "runs"

    config_path = tmp_path / "config.yml"
    config_path.write_text(
        "\n".join(
            [
                "as_of_date: 2025-12-31",
                f"mosers_all_programs_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - All Programs.xlsx'}",
                f"mosers_ex_trend_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - Ex Trend.xlsx'}",
                f"mosers_trend_xlsx: {fixtures / 'MOSERS Counterparty Risk Summary 12-31-2025 - Trend.xlsx'}",
                f"hist_all_programs_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - All Programs 3 Year.xlsx'}",
                f"hist_ex_llc_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - ex LLC 3 Year.xlsx'}",
                f"hist_llc_3yr_xlsx: {fixtures / 'Historical Counterparty Risk Graphs - LLC 3 Year.xlsx'}",
                f"monthly_pptx: {fixtures / 'Monthly Counterparty Exposure Report.pptx'}",
                f"output_root: {output_root}",
                "distribution_static: true",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    # Force non-Windows so the fallback path is exercised.
    monkeypatch.setattr("counter_risk.pipeline.run.platform.system", lambda: "Linux")

    run_dir = run_pipeline(config_path)

    manifest = json.loads((run_dir / "manifest.json").read_text(encoding="utf-8"))

    # Regular PPT is still produced (distribution filename for as_of_date 2025-12-31).
    assert (run_dir / "Monthly Counterparty Exposure Report - 2025-12-31.pptx").exists()

    # Warning about static distribution being unavailable is present.
    assert any("distribution_static" in w for w in manifest["warnings"])

    # No static distribution files produced (no COM available).
    assert not (
        run_dir
        / "Monthly Counterparty Exposure Report (Master) - 2025-12-31_distribution_static.pptx"
    ).exists()
    assert not (
        run_dir / "Monthly Counterparty Exposure Report (Master) - 2025-12-31_distribution.pdf"
    ).exists()

    # Config snapshot captures the flag.
    assert manifest["config_snapshot"]["distribution_static"] is True


# ---------------------------------------------------------------------------
# _export_chart_shapes_as_images tests
# ---------------------------------------------------------------------------


def test_export_chart_shapes_as_images_returns_empty_when_no_ole_shapes(
    tmp_path: Path,
) -> None:
    """When a presentation has no OLE/chart shapes the result dict is empty."""
    # Build a fake COM presentation with one slide containing a plain shape.
    slide_images_dir = tmp_path / "imgs"
    slide_images_dir.mkdir()

    class _FakeShape:
        Type = 1  # msoAutoShape – not OLE
        HasChart = False
        Id = 1
        Name = "Shape 1"

        def Export(self, path: str, fmt: int) -> None:  # noqa: N802  # pragma: no cover
            raise AssertionError("Export should not be called for non-OLE shapes")

    class _FakeSlide:
        Shapes = [_FakeShape()]

    class _FakePres:
        Slides = [_FakeSlide()]

        @property
        def Count(self) -> int:  # noqa: N802  # pragma: no cover
            return 1

    # Build a Slides-like object with .Count and index access.
    class _SlideList:
        def __init__(self) -> None:
            self._slides = [_FakeSlide()]
            self.Count = 1

        def __iter__(self) -> Any:
            return iter(self._slides)

        def __getitem__(self, idx: int) -> Any:
            return self._slides[idx - 1]

    class _Presentation:
        def __init__(self) -> None:
            self.Slides = _SlideList()

    warnings: list[str] = []
    result = run_module._export_chart_shapes_as_images(
        com_presentation=_Presentation(),
        slide_images_dir=slide_images_dir,
        warnings=warnings,
    )

    assert result == {}
    assert warnings == []


def test_export_chart_shapes_as_images_records_warning_on_export_failure(
    tmp_path: Path,
) -> None:
    """When Export() raises, a warning is recorded and the shape is skipped."""
    slide_images_dir = tmp_path / "imgs"
    slide_images_dir.mkdir()

    class _FailingShape:
        Type = 7  # msoEmbeddedOLEObject
        HasChart = False
        Id = 5
        Name = "Chart 5"

        def Export(self, path: str, fmt: int) -> None:  # noqa: N802
            raise RuntimeError("COM export failed")

    class _SlideList:
        def __init__(self) -> None:
            self.Count = 1

        def __getitem__(self, idx: int):  # type: ignore[no-untyped-def]
            class _Slide:
                Shapes = [_FailingShape()]

            return _Slide()

    class _Presentation:
        def __init__(self) -> None:
            self.Slides = _SlideList()

    warnings: list[str] = []
    result = run_module._export_chart_shapes_as_images(
        com_presentation=_Presentation(),
        slide_images_dir=slide_images_dir,
        warnings=warnings,
    )

    assert result == {}
    assert len(warnings) == 1
    assert "Chart 5" in warnings[0]
    assert "distribution_static chart export failed" in warnings[0]


def test_export_chart_shapes_as_images_collects_ole_shapes(
    tmp_path: Path,
) -> None:
    """OLE shapes that export successfully are included in the result mapping."""
    slide_images_dir = tmp_path / "imgs"
    slide_images_dir.mkdir()
    written: list[str] = []

    class _OleShape:
        Type = 10  # msoLinkedOLEObject
        HasChart = False
        Id = 3
        Name = "Chart 3"

        def Export(self, path: str, fmt: int) -> None:  # noqa: N802
            Path(path).write_bytes(b"fake-png")
            written.append(path)

    class _SlideList:
        def __init__(self) -> None:
            self.Count = 1

        def __getitem__(self, idx: int):  # type: ignore[no-untyped-def]
            class _Slide:
                Shapes = [_OleShape()]

            return _Slide()

    class _Presentation:
        def __init__(self) -> None:
            self.Slides = _SlideList()

    warnings: list[str] = []
    result = run_module._export_chart_shapes_as_images(
        com_presentation=_Presentation(),
        slide_images_dir=slide_images_dir,
        warnings=warnings,
    )

    assert (1, "Chart 3") in result
    assert result[(1, "Chart 3")].exists()
    assert warnings == []


# ---------------------------------------------------------------------------
# _rebuild_pptx_replacing_charts tests
# ---------------------------------------------------------------------------


def _make_test_pptx_with_picture(tmp_path: Path, shape_name: str) -> Path:
    """Create a minimal PPTX containing one slide with a single picture shape."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    # Use blank slide layout.
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Create a tiny 1x1 PNG (minimal valid PNG bytes).
    png_bytes = (
        b"\x89PNG\r\n\x1a\n"  # PNG signature
        b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x02\x00\x00\x00\x90wS\xde"  # 1x1 RGB
        b"\x00\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00\x05\x18\xd8N"
        b"\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    img_path = tmp_path / "placeholder.png"
    img_path.write_bytes(png_bytes)

    pic = slide.shapes.add_picture(str(img_path), Inches(1), Inches(1), Inches(4), Inches(3))
    pic.name = shape_name

    out = tmp_path / "source.pptx"
    prs.save(str(out))
    return out


def test_rebuild_pptx_replacing_charts_empty_mapping_preserves_slide(
    tmp_path: Path,
) -> None:
    """With an empty chart_images mapping the output PPTX has the same slide count."""
    from pptx import Presentation

    source = _make_test_pptx_with_picture(tmp_path, "SomeShape")
    output = tmp_path / "out.pptx"

    run_module._rebuild_pptx_replacing_charts(
        source_pptx=source,
        output_path=output,
        chart_images={},
    )

    assert output.exists()
    prs = Presentation(str(output))
    assert len(prs.slides) == 1
    # Shape count unchanged when no replacements requested.
    assert len(list(prs.slides[0].shapes)) == 1


def test_rebuild_pptx_replacing_charts_replaces_named_shape(
    tmp_path: Path,
) -> None:
    """A shape matching a chart_images key is removed and replaced with a picture."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    shape_name = "Chart OLE 1"
    source = _make_test_pptx_with_picture(tmp_path, shape_name)

    # Replacement image (minimal valid PNG).
    png_bytes = (
        b"\x89PNG\r\n\x1a\n"
        b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x02\x00\x00\x00\x90wS\xde"
        b"\x00\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00\x05\x18\xd8N"
        b"\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    replacement = tmp_path / "replacement.png"
    replacement.write_bytes(png_bytes)

    output = tmp_path / "out.pptx"
    run_module._rebuild_pptx_replacing_charts(
        source_pptx=source,
        output_path=output,
        chart_images={(1, shape_name): replacement},
    )

    assert output.exists()
    prs = Presentation(str(output))
    slide = prs.slides[0]
    # The slide should still have exactly one shape (the replacement picture).
    shapes = list(slide.shapes)
    assert len(shapes) == 1
    assert shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE


def test_rebuild_pptx_replacing_charts_preserves_position(
    tmp_path: Path,
) -> None:
    """The replacement picture is placed at the same geometry as the original shape."""
    from pptx import Presentation

    shape_name = "OLE Chart"
    source = _make_test_pptx_with_picture(tmp_path, shape_name)

    # Record original shape geometry.
    original_prs = Presentation(str(source))
    orig_shape = list(original_prs.slides[0].shapes)[0]
    orig_left, orig_top = orig_shape.left, orig_shape.top
    orig_width, orig_height = orig_shape.width, orig_shape.height

    png_bytes = (
        b"\x89PNG\r\n\x1a\n"
        b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x02\x00\x00\x00\x90wS\xde"
        b"\x00\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00\x05\x18\xd8N"
        b"\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    replacement = tmp_path / "replacement.png"
    replacement.write_bytes(png_bytes)

    output = tmp_path / "out.pptx"
    run_module._rebuild_pptx_replacing_charts(
        source_pptx=source,
        output_path=output,
        chart_images={(1, shape_name): replacement},
    )

    prs = Presentation(str(output))
    new_shape = list(prs.slides[0].shapes)[0]
    assert new_shape.left == orig_left
    assert new_shape.top == orig_top
    assert new_shape.width == orig_width
    assert new_shape.height == orig_height


def test_shape_match_confidence_high_when_name_is_unique() -> None:
    class _Shape:
        def __init__(self) -> None:
            self.name = "Chart 1"
            self.left = 10
            self.top = 10
            self.width = 100
            self.height = 100

    confidence = run_module._shape_match_confidence(
        target_name="Chart 1",
        candidate_shapes=[_Shape()],
    )

    assert confidence >= 0.95


def test_shape_match_confidence_low_when_name_is_duplicated() -> None:
    class _Shape:
        def __init__(self) -> None:
            self.name = "Chart 1"
            self.left = 10
            self.top = 10
            self.width = 100
            self.height = 100

    confidence = run_module._shape_match_confidence(
        target_name="Chart 1",
        candidate_shapes=[_Shape(), _Shape()],
    )

    assert confidence < 0.8


def test_rebuild_pptx_replacing_charts_replaces_whole_slide_on_low_confidence(
    tmp_path: Path,
) -> None:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches

    # Build a slide with duplicate names to force low-confidence per-shape matching.
    source = _make_test_pptx_with_picture(tmp_path, "Chart OLE")
    prs = Presentation(str(source))
    slide = prs.slides[0]
    second_shape_image = tmp_path / "shape2.png"
    second_shape_image.write_bytes(
        b"\x89PNG\r\n\x1a\n"
        b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x02\x00\x00\x00\x90wS\xde"
        b"\x00\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00\x05\x18\xd8N"
        b"\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    duplicate = slide.shapes.add_picture(str(second_shape_image), Inches(2), Inches(2))
    duplicate.name = "Chart OLE"
    prs.save(str(source))

    replacement = tmp_path / "replacement.png"
    replacement.write_bytes(second_shape_image.read_bytes())
    slide_fallback = tmp_path / "slide_fallback.png"
    slide_fallback.write_bytes(second_shape_image.read_bytes())

    output = tmp_path / "out.pptx"
    run_module._rebuild_pptx_replacing_charts(
        source_pptx=source,
        output_path=output,
        chart_images={(1, "Chart OLE"): replacement},
        fallback_slide_images={1: slide_fallback},
    )

    rebuilt = Presentation(str(output))
    shapes = list(rebuilt.slides[0].shapes)
    assert len(shapes) == 1
    assert shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE


def test_rebuild_pptx_replacing_charts_raises_when_low_confidence_has_no_slide_fallback(
    tmp_path: Path,
) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    source = _make_test_pptx_with_picture(tmp_path, "Chart OLE")
    prs = Presentation(str(source))
    slide = prs.slides[0]
    second_shape_image = tmp_path / "shape2.png"
    second_shape_image.write_bytes(
        b"\x89PNG\r\n\x1a\n"
        b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x02\x00\x00\x00\x90wS\xde"
        b"\x00\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00\x05\x18\xd8N"
        b"\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    duplicate = slide.shapes.add_picture(str(second_shape_image), Inches(2), Inches(2))
    duplicate.name = "Chart OLE"
    prs.save(str(source))

    replacement = tmp_path / "replacement.png"
    replacement.write_bytes(second_shape_image.read_bytes())
    output = tmp_path / "out.pptx"

    with pytest.raises(RuntimeError, match="no slide-image fallback"):
        run_module._rebuild_pptx_replacing_charts(
            source_pptx=source,
            output_path=output,
            chart_images={(1, "Chart OLE"): replacement},
        )


def test_rebuild_pptx_replacing_charts_raises_for_full_deck_rebuild_on_low_confidence(
    tmp_path: Path,
) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    source = _make_test_pptx_with_picture(tmp_path, "Chart OLE")
    prs = Presentation(str(source))
    slide = prs.slides[0]
    second_shape_image = tmp_path / "shape2.png"
    second_shape_image.write_bytes(
        b"\x89PNG\r\n\x1a\n"
        b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x02\x00\x00\x00\x90wS\xde"
        b"\x00\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00\x05\x18\xd8N"
        b"\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    duplicate = slide.shapes.add_picture(str(second_shape_image), Inches(2), Inches(2))
    duplicate.name = "Chart OLE"
    prs.save(str(source))

    replacement = tmp_path / "replacement.png"
    replacement.write_bytes(second_shape_image.read_bytes())
    slide_fallback = tmp_path / "slide_fallback.png"
    slide_fallback.write_bytes(second_shape_image.read_bytes())
    output = tmp_path / "out.pptx"

    with pytest.raises(RuntimeError, match="full-deck static rebuild required"):
        run_module._rebuild_pptx_replacing_charts(
            source_pptx=source,
            output_path=output,
            chart_images={(1, "Chart OLE"): replacement},
            fallback_slide_images={1: slide_fallback},
            fallback_to_full_deck_rebuild=True,
        )

    assert not output.exists()


# ---------------------------------------------------------------------------
# _apply_chart_replacement tests
# ---------------------------------------------------------------------------


def test_apply_chart_replacement_skips_on_non_windows(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """On non-Windows, chart replacement is silently skipped."""
    monkeypatch.setattr("counter_risk.pipeline.run.platform.system", lambda: "Linux")
    result = run_module._apply_chart_replacement(
        master_pptx_path=tmp_path / "master.pptx",
        output_path=tmp_path / "out.pptx",
        run_dir=tmp_path,
        static_mode=False,
        warnings=[],
    )
    assert result is False


def test_apply_chart_replacement_skips_when_no_chart_shapes(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """When COM finds no chart shapes, replacement is skipped."""
    monkeypatch.setattr("counter_risk.pipeline.run.platform.system", lambda: "Windows")

    class _FakeShape:
        Type = 1
        HasChart = False
        Id = 1
        Name = "TextBox"

        def Export(self, path: str, fmt: str) -> None:  # noqa: N802
            raise AssertionError("should not be called")

    class _SlideList:
        Count = 1

        def __getitem__(self, idx: int) -> Any:
            class _Slide:
                Shapes = [_FakeShape()]

            return _Slide()

    class _FakePres:
        Slides = _SlideList()

        def Close(self) -> None:  # noqa: N802
            pass

    class _FakeApp:
        Visible = False
        Presentations = types.SimpleNamespace(Open=lambda *a, **kw: _FakePres())

        def Quit(self) -> None:  # noqa: N802
            pass

    fake_client = types.SimpleNamespace(DispatchEx=lambda *a, **kw: _FakeApp())
    fake_win32com = types.ModuleType("win32com")
    cast(Any, fake_win32com).client = fake_client
    monkeypatch.setitem(sys.modules, "win32com", fake_win32com)
    monkeypatch.setitem(sys.modules, "win32com.client", fake_client)

    source = tmp_path / "master.pptx"
    source.write_bytes(b"fake")
    warnings: list[str] = []
    result = run_module._apply_chart_replacement(
        master_pptx_path=source,
        output_path=tmp_path / "out.pptx",
        run_dir=tmp_path,
        static_mode=False,
        warnings=warnings,
    )
    assert result is False


def test_apply_chart_replacement_returns_false_on_com_failure(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """When COM session raises, replacement returns False and records a warning."""
    monkeypatch.setattr("counter_risk.pipeline.run.platform.system", lambda: "Windows")

    class _FailApp:
        Visible = False
        Presentations = types.SimpleNamespace(
            Open=lambda *a, **kw: (_ for _ in ()).throw(RuntimeError("COM init failed"))
        )

        def Quit(self) -> None:  # noqa: N802
            pass

    fake_client = types.SimpleNamespace(DispatchEx=lambda *a, **kw: _FailApp())
    fake_win32com = types.ModuleType("win32com")
    cast(Any, fake_win32com).client = fake_client
    monkeypatch.setitem(sys.modules, "win32com", fake_win32com)
    monkeypatch.setitem(sys.modules, "win32com.client", fake_client)

    source = tmp_path / "master.pptx"
    source.write_bytes(b"fake")
    warnings: list[str] = []
    result = run_module._apply_chart_replacement(
        master_pptx_path=source,
        output_path=tmp_path / "out.pptx",
        run_dir=tmp_path,
        static_mode=False,
        warnings=warnings,
    )
    assert result is False
    assert any("chart_replacement COM session failed" in w for w in warnings)
