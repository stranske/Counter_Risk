"""Utilities for deterministic PowerPoint assertions in tests."""

from __future__ import annotations

from pathlib import Path
from typing import Any


def _as_presentation(pptx_obj_or_path: Any) -> Any:
    if isinstance(pptx_obj_or_path, (str, Path)):
        from pptx import Presentation  # type: ignore[import-not-found]

        return Presentation(str(pptx_obj_or_path))
    return pptx_obj_or_path


def assert_slide_count(pptx_obj_or_path: Any, expected_count: int) -> None:
    """Assert a presentation has exactly ``expected_count`` slides."""
    presentation = _as_presentation(pptx_obj_or_path)
    actual_count = len(presentation.slides)
    if actual_count != expected_count:
        raise AssertionError(f"Expected {expected_count} slides but found {actual_count}.")


def assert_chart_present(pptx_obj_or_path: Any) -> None:
    """Assert that at least one chart shape exists in the presentation."""
    presentation = _as_presentation(pptx_obj_or_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_chart", False):
                return
    raise AssertionError("Expected at least one chart shape in presentation.")


def assert_picture_present(pptx_obj_or_path: Any) -> None:
    """Assert that at least one picture shape exists in the presentation."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import-not-found]

    presentation = _as_presentation(pptx_obj_or_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return
    raise AssertionError("Expected at least one picture shape in presentation.")
