"""Screenshot replacement helpers for PowerPoint reports."""

from __future__ import annotations

from collections.abc import Mapping
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from pptx.slide import Slide

_SECTION_TITLE_SUBSTRINGS: dict[str, tuple[str, ...]] = {
    "allprograms": ("allprograms",),
    "extrend": ("extrend", "excludingtrend"),
    "trend": ("trend",),
}

_EXPECTED_PICTURE_GEOMETRY_BY_SECTION: dict[str, tuple[tuple[int, int, int, int], ...]] = {
    "allprograms": ((0, 811705, 9144000, 5817695),),
    "extrend": ((0, 1304636, 9144000, 5172364),),
    "trend": ((0, 2031298, 9144000, 2795403),),
}


def _normalize_key(value: str) -> str:
    return "".join(ch.lower() for ch in value if ch.isalnum())


def _slide_title(slide: Slide) -> str:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text_frame = getattr(shape, "text_frame", None)
        if text_frame is None:
            continue
        text = str(getattr(text_frame, "text", "")).strip()
        if text:
            return text
    return ""


def _picture_shapes(slide: Slide) -> list[BaseShape]:
    return [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]


def _picture_geometry(picture: BaseShape) -> tuple[int, int, int, int]:
    return picture.left, picture.top, picture.width, picture.height


def _canonical_section_key(section: str) -> str:
    normalized = _normalize_key(section)
    for canonical_key, title_substrings in _SECTION_TITLE_SUBSTRINGS.items():
        if normalized == canonical_key or normalized in title_substrings:
            return canonical_key
    return normalized


def _section_matches_title(section_key: str, normalized_title: str) -> bool:
    for title_substring in _SECTION_TITLE_SUBSTRINGS.get(section_key, (section_key,)):
        if title_substring in normalized_title:
            return True
    return False


def _slide_matches_expected_picture_geometry(slide: Slide, section_key: str) -> bool:
    expected_geometry = _EXPECTED_PICTURE_GEOMETRY_BY_SECTION.get(section_key)
    if expected_geometry is None:
        return bool(_picture_shapes(slide))

    pictures = _picture_shapes(slide)
    if len(pictures) != len(expected_geometry):
        return False

    actual_geometry = tuple(_picture_geometry(picture) for picture in pictures)
    return actual_geometry == expected_geometry


def _resolve_image_path(value: Path | str) -> Path:
    image_path = Path(value)
    if not image_path.exists() or not image_path.is_file():
        raise ValueError(f"replacement image does not exist: {image_path}")
    return image_path


def _swap_picture_image_part(picture: BaseShape, replacement: Path) -> None:
    """Swap a picture's embedded image part while preserving shape geometry.

    This relies on private python-pptx API (`picture._element`) because public
    APIs do not provide an in-place image swap; updating `a:blip/@r:embed`
    allows replacing image content without appending/removing shapes.
    """
    _, new_rid = picture.part.get_or_add_image_part(str(replacement))
    blip = picture._element.blipFill.blip
    if blip is None:
        raise ValueError("picture shape is missing blip image content")

    old_rid = blip.rEmbed
    blip.rEmbed = new_rid

    if old_rid and old_rid != new_rid and picture.part._rel_ref_count(old_rid) == 0:
        picture.part.drop_rel(old_rid)


def replace_screenshot_pictures(
    pptx_in: Path | str,
    images_by_section: Mapping[str, Path | str],
    pptx_out: Path | str,
) -> None:
    """Replace existing screenshot pictures by matching section titles.

    Matching is based on section-title text plus existing picture shapes in the
    matched slide. Existing picture shapes are replaced in-place by preserving
    their count, location, and size.
    """

    source = Path(pptx_in)
    destination = Path(pptx_out)
    if not source.exists() or not source.is_file():
        raise ValueError(f"input pptx does not exist: {source}")

    normalized_targets: dict[str, Path] = {}
    for section, image in images_by_section.items():
        normalized_targets[_canonical_section_key(section)] = _resolve_image_path(image)

    if not normalized_targets:
        raise ValueError("images_by_section must contain at least one section mapping")

    presentation = Presentation(str(source))
    matched_sections: set[str] = set()

    for slide in presentation.slides:
        title_key = _normalize_key(_slide_title(slide))
        if not title_key:
            continue

        matched_target: str | None = None
        for section_key in normalized_targets:
            if _section_matches_title(section_key, title_key):
                matched_target = section_key
                break

        if matched_target is None:
            continue

        pictures = _picture_shapes(slide)
        if not pictures:
            raise ValueError(f"matched slide has no picture shapes for section '{matched_target}'")

        if not _slide_matches_expected_picture_geometry(slide, matched_target):
            continue
        replacement = normalized_targets[matched_target]
        for picture in list(pictures):
            _swap_picture_image_part(picture, replacement)

        if len(_picture_shapes(slide)) != len(pictures):
            raise ValueError(
                f"picture replacement changed shape count for section '{matched_target}'"
            )

        matched_sections.add(matched_target)

    missing = sorted(set(normalized_targets) - matched_sections)
    if missing:
        raise ValueError("no matching slide title found for sections: " + ", ".join(missing))

    destination.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(destination))
