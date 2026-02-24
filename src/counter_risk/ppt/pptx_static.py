"""Helpers for producing static, image-only PowerPoint decks."""

from __future__ import annotations

from pathlib import Path


class PngValidationError(ValueError):
    """Raised when a slide PNG fails static rebuild validation."""


def _validate_slide_png(png_path: Path) -> None:
    """Validate a slide PNG before inserting it into a rebuilt PPTX."""

    if not png_path.exists():
        raise PngValidationError(f"PNG validation failed for '{png_path}': file not found")
    if not png_path.is_file():
        raise PngValidationError(f"PNG validation failed for '{png_path}': not a file")
    if png_path.stat().st_size <= 0:
        raise PngValidationError(f"PNG validation failed for '{png_path}': empty file")

    try:
        from PIL import Image, UnidentifiedImageError

        with Image.open(png_path) as image:
            image.verify()
    except (UnidentifiedImageError, OSError, SyntaxError, ValueError) as exc:
        raise PngValidationError(
            f"PNG validation failed for '{png_path}': invalid PNG ({exc})"
        ) from exc


def _rebuild_pptx_from_slide_images(
    *,
    source_pptx: Path,
    slide_images: list[Path],
    output_path: Path,
) -> None:
    """Create a new PPTX where every slide is a single full-bleed PNG image."""
    from pptx import Presentation

    for img_path in slide_images:
        _validate_slide_png(img_path)

    source_prs = Presentation(str(source_pptx))
    assert source_prs.slide_width is not None, "source PPT has no slide width"
    assert source_prs.slide_height is not None, "source PPT has no slide height"
    slide_width = source_prs.slide_width
    slide_height = source_prs.slide_height

    new_prs = Presentation()
    new_prs.slide_width = slide_width
    new_prs.slide_height = slide_height
    blank_layout = new_prs.slide_layouts[6]

    for img_path in slide_images:
        new_slide = new_prs.slides.add_slide(blank_layout)
        new_slide.shapes.add_picture(
            str(img_path), left=0, top=0, width=slide_width, height=slide_height
        )

    new_prs.save(str(output_path))
