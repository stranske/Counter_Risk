"""Append a concentration metrics summary table slide to a PPTX file."""

from __future__ import annotations

from pathlib import Path
from typing import Any


def append_concentration_table_slide(
    pptx_path: Path | str,
    metrics_records: list[dict[str, Any]],
    *,
    title: str = "Concentration Metrics",
) -> None:
    """Append a slide with a concentration metrics table to an existing PPTX.

    The new slide is appended after all existing slides and contains a native
    python-pptx table with one row per ``(variant, segment)`` group and columns
    for Top 5 share, Top 10 share, and HHI.  Share values are formatted as
    percentages; HHI is formatted to four decimal places.

    Parameters
    ----------
    pptx_path:
        Path to the PPTX file to modify in-place.
    metrics_records:
        List of dicts, each with keys ``variant``, ``segment``,
        ``top5_share``, ``top10_share``, and ``hhi``.
    title:
        Text label placed at the top of the new slide.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    assert slide_height is not None

    title_left = Inches(0.5)
    title_top = Inches(0.25)
    title_width = slide_width - Inches(1)
    title_height = Inches(0.45)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].runs[0].font.size = Pt(18)
    tf.paragraphs[0].runs[0].font.bold = True

    if not metrics_records:
        prs.save(str(pptx_path))
        return

    _headers = ("Variant", "Segment", "Top 5 Share", "Top 10 Share", "HHI")
    _keys = ("variant", "segment", "top5_share", "top10_share", "hhi")

    rows_count = len(metrics_records) + 1  # header row + data rows
    cols_count = len(_headers)

    table_left = Inches(0.5)
    table_top = Inches(0.85)
    table_width = slide_width - Inches(1)
    table_height = slide_height - Inches(1.1)

    tbl = slide.shapes.add_table(
        rows_count, cols_count, table_left, table_top, table_width, table_height
    ).table

    for col_idx, header in enumerate(_headers):
        cell = tbl.cell(0, col_idx)
        cell.text = header
        para = cell.text_frame.paragraphs[0]
        if para.runs:
            para.runs[0].font.bold = True
            para.runs[0].font.size = Pt(10)

    for row_idx, record in enumerate(metrics_records):
        for col_idx, key in enumerate(_keys):
            cell = tbl.cell(row_idx + 1, col_idx)
            value = record.get(key, "")
            if key in ("top5_share", "top10_share"):
                try:
                    cell.text = f"{float(value):.2%}"
                except (TypeError, ValueError):
                    cell.text = str(value)
            elif key == "hhi":
                try:
                    cell.text = f"{float(value):.4f}"
                except (TypeError, ValueError):
                    cell.text = str(value)
            else:
                cell.text = str(value)
            para = cell.text_frame.paragraphs[0]
            if para.runs:
                para.runs[0].font.size = Pt(9)

    prs.save(str(pptx_path))
