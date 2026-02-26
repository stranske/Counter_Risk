"""Counter Risk pipeline orchestration."""

from __future__ import annotations

import contextlib
import hashlib
import logging
import platform
import shutil
from collections.abc import Callable, Mapping
from dataclasses import dataclass
from datetime import date
from enum import StrEnum
from pathlib import Path
from typing import TYPE_CHECKING, Any, Literal, NoReturn
from zipfile import BadZipFile

from counter_risk.config import WorkflowConfig, load_config
from counter_risk.dates import derive_as_of_date, derive_run_date
from counter_risk.normalize import (
    canonicalize_name,
    normalize_counterparty,
    normalize_counterparty_with_source,
)
from counter_risk.outputs.base import OutputContext, OutputGenerator
from counter_risk.parsers import parse_fcm_totals, parse_futures_detail
from counter_risk.pipeline.manifest import ManifestBuilder
from counter_risk.pipeline.parsing_types import (
    ParsedDataInvalidShapeError,
    ParsedDataMissingKeyError,
    UnmappedCounterpartyError,
)
from counter_risk.pipeline.ppt_naming import PptOutputNames, resolve_ppt_output_names
from counter_risk.pipeline.ppt_validation import validate_distribution_ppt_standalone
from counter_risk.pipeline.run_folder_outputs import (
    RunFolderReadmePptOutputs,
    build_run_folder_readme_content,
)
from counter_risk.pipeline.time_utils import utc_now_isoformat
from counter_risk.ppt.pptx_postprocess import (
    list_external_relationship_targets,
    scrub_external_relationships_from_pptx,
)
from counter_risk.ppt.pptx_static import _rebuild_pptx_from_slide_images
from counter_risk.writers import generate_mosers_workbook

LOGGER = logging.getLogger(__name__)

if TYPE_CHECKING:
    from counter_risk.outputs.ppt_link_refresh import PptLinkRefreshOutputGenerator

_EXPECTED_VARIANTS: tuple[str, ...] = ("all_programs", "ex_trend", "trend")

# PowerPoint COM shape type constants used when identifying OLE/chart shapes.
_COM_SHAPE_TYPE_EMBEDDED_OLE: int = 7  # msoEmbeddedOLEObject
_COM_SHAPE_TYPE_LINKED_OLE: int = 10  # msoLinkedOLEObject
_COM_SHAPE_FORMAT_PNG: int = 2  # ppShapeFormatPNG
_SHAPE_MATCH_CONFIDENCE_THRESHOLD: float = 0.8
_REQUIRED_TOTAL_COLUMNS: tuple[str, ...] = ("counterparty", "Notional", "NotionalChange")
_REQUIRED_FUTURES_COLUMNS: tuple[str, ...] = (
    "account",
    "description",
    "class",
    "fcm",
    "clearing_house",
    "notional",
)
_PREFERRED_HISTORICAL_SHEET_BY_VARIANT: dict[str, str] = {
    "all_programs": "Total",
    "ex_trend": "Total",
    "trend": "Total",
}
_DATE_HEADER_CANDIDATES: tuple[str, ...] = ("date", "as of date", "as-of date")
_REQUIRED_HISTORICAL_APPEND_HEADERS: tuple[str, ...] = (
    "date",
    "value series 1",
    "value series 2",
)


@dataclass(frozen=True)
class _VariantInputs:
    name: str
    workbook_path: Path
    historical_path: Path


class PptProcessingStatus(StrEnum):
    """Machine-readable statuses for PPT processing."""

    SUCCESS = "success"
    SKIPPED = "skipped"
    FAILED = "failed"


@dataclass(frozen=True)
class PptProcessingResult:
    """Result envelope for PPT processing and link refresh."""

    status: PptProcessingStatus
    error_detail: str | None = None


ScreenshotReplacer = Callable[[Path, Path, dict[str, Path]], None]


def reconcile_series_coverage(
    *,
    parsed_data_by_sheet: Mapping[str, Mapping[str, Any]],
    historical_series_headers_by_sheet: Mapping[str, tuple[str, ...] | list[str] | set[str]],
    variant: str | None = None,
    expected_segments_by_variant: (
        Mapping[str, tuple[str, ...] | list[str] | set[str]] | None
    ) = None,
    fail_policy: Literal["warn", "strict"] = "warn",
) -> dict[str, Any]:
    """Reconcile parsed series labels against historical workbook headers per sheet.

    Compares current-month series labels from parsed tables against historical workbook
    headers and optionally validates variant-specific segment expectations.
    """
    _validate_reconciliation_parsed_data_by_sheet(parsed_data_by_sheet=parsed_data_by_sheet)

    by_sheet: dict[str, dict[str, Any]] = {}
    missing_series: list[dict[str, Any]] = []
    missing_segments: list[dict[str, Any]] = []
    exceptions: list[UnmappedCounterpartyError] = []
    warnings: list[str] = []
    gap_count = 0

    expected_segments = _expected_segments_for_variant(
        variant=variant,
        expected_segments_by_variant=expected_segments_by_variant,
    )
    sheet_names = sorted(
        set(parsed_data_by_sheet).union(historical_series_headers_by_sheet), key=str.casefold
    )
    for sheet_name in sheet_names:
        parsed_sections = parsed_data_by_sheet.get(sheet_name, {})
        totals_records = _records(parsed_sections.get("totals", []))
        futures_records = _records(parsed_sections.get("futures", []))
        historical_series_headers = sorted(
            {
                value
                for value in (
                    canonicalize_name(str(header))
                    for header in historical_series_headers_by_sheet.get(sheet_name, ())
                )
                if value
            }
        )

        counterparties_in_data = sorted(
            {
                value
                for value in (
                    str(record.get("counterparty", "")).strip() for record in totals_records
                )
                if value
            }
        )
        (
            normalized_counterparties_in_data,
            counterparty_sources_by_raw_name,
        ) = _counterparty_resolution_maps_from_records(totals_records)
        raw_counterparties_by_normalized = _raw_counterparties_by_normalized_from_parsed_data(
            parsed_sections
        )
        clearing_houses_in_data = sorted(
            {
                value
                for value in (
                    canonicalize_name(str(record.get("clearing_house", "")))
                    for record in futures_records
                )
                if value
            }
        )
        normalized_historical_series_headers = {
            normalize_counterparty(header) for header in historical_series_headers
        }
        missing_normalized_counterparties = sorted(
            set(normalized_counterparties_in_data).difference(normalized_historical_series_headers),
            key=str.casefold,
        )
        missing_counterparty_labels = sorted(
            {
                canonicalize_name(raw_name)
                for normalized_name in missing_normalized_counterparties
                for raw_name in normalized_counterparties_in_data.get(normalized_name, ())
                if canonicalize_name(raw_name)
            },
            key=str.casefold,
        )
        missing_clearing_houses = sorted(
            set(clearing_houses_in_data).difference(historical_series_headers), key=str.casefold
        )
        current_series_labels = sorted(
            set(counterparties_in_data).union(clearing_houses_in_data), key=str.casefold
        )
        normalized_current_series_labels = set(normalized_counterparties_in_data).union(
            {normalize_counterparty(clearing_house) for clearing_house in clearing_houses_in_data}
        )
        missing_from_historical = sorted(
            set(missing_counterparty_labels).union(missing_clearing_houses), key=str.casefold
        )
        missing_from_data = sorted(
            {
                header
                for header in historical_series_headers
                if normalize_counterparty(header) not in normalized_current_series_labels
            },
            key=str.casefold,
        )
        parsed_segments = _extract_segments_from_records(parsed_sections)
        missing_expected_segments = sorted(
            expected_segments.difference(parsed_segments), key=str.casefold
        )

        if missing_from_historical:
            gap_count += len(missing_from_historical)
            missing_series.append(
                {
                    "sheet": sheet_name,
                    "missing_from_historical_headers": missing_from_historical,
                    "data_source_context": "counterparties_and_clearing_houses",
                }
            )
            warnings.append(
                "Reconciliation gap in sheet "
                f"{sheet_name!r}: series present in parsed data but missing from historical "
                f"headers ({', '.join(missing_from_historical)})"
            )

        if missing_from_data:
            gap_count += len(missing_from_data)
            warnings.append(
                "Reconciliation gap in sheet "
                f"{sheet_name!r}: series present in historical headers but missing from parsed "
                f"data ({', '.join(missing_from_data)})"
            )

        if missing_normalized_counterparties:
            sheet_exceptions: list[UnmappedCounterpartyError] = []
            raw_counterparties_for_metadata: list[str] = []
            for normalized_name in missing_normalized_counterparties:
                raw_names = sorted(
                    set(raw_counterparties_by_normalized.get(normalized_name, ())),
                    key=str.casefold,
                )
                raw_counterparties_for_metadata.extend(raw_names)
                raw_display = ", ".join(raw_names)
                warnings.append(
                    "Reconciliation unmapped counterparty in sheet "
                    f"{sheet_name!r}: raw={raw_display!r}, normalized={normalized_name!r}, "
                    "source="
                    + ",".join(
                        sorted(
                            {
                                counterparty_sources_by_raw_name.get(raw_name, "unmapped")
                                for raw_name in raw_names
                            }
                        )
                    )
                )
                for raw_name in raw_names:
                    error = UnmappedCounterpartyError(
                        normalized_counterparty=normalized_name,
                        raw_counterparty=raw_name,
                        sheet=sheet_name,
                    )
                    sheet_exceptions.append(error)
                    exceptions.append(error)

            if raw_counterparties_for_metadata:
                missing_series.append(
                    {
                        "sheet": sheet_name,
                        "error_type": "unmapped_counterparty",
                        "raw_counterparties": raw_counterparties_for_metadata,
                        "normalized_counterparties": missing_normalized_counterparties,
                    }
                )

            if fail_policy == "strict":
                if sheet_exceptions:
                    raise sheet_exceptions[0]
                raise UnmappedCounterpartyError(
                    normalized_counterparty=missing_normalized_counterparties[0],
                    raw_counterparty=missing_normalized_counterparties[0],
                    sheet=sheet_name,
                )

        if missing_expected_segments:
            gap_count += len(missing_expected_segments)
            missing_segments.append(
                {
                    "variant": variant,
                    "sheet": sheet_name,
                    "expected_segment_identifiers": missing_expected_segments,
                }
            )
            warnings.append(
                "Reconciliation gap in variant "
                f"{variant!r} sheet {sheet_name!r}: expected segments missing from parsed "
                f"results ({', '.join(missing_expected_segments)})"
            )

        canonical_key_by_series: dict[str, str] = {}
        for canonical_name, raw_name_set in normalized_counterparties_in_data.items():
            for raw in raw_name_set:
                canonical_key_by_series[raw] = canonical_name
        for ch in clearing_houses_in_data:
            canonical_key_by_series[ch] = normalize_counterparty(ch)

        by_sheet[sheet_name] = {
            "counterparties_in_data": counterparties_in_data,
            "normalized_counterparties_in_data": sorted(
                normalized_counterparties_in_data, key=str.casefold
            ),
            "clearing_houses_in_data": clearing_houses_in_data,
            "historical_series_headers": historical_series_headers,
            "normalized_historical_series_headers": sorted(
                normalized_historical_series_headers, key=str.casefold
            ),
            "current_series_labels": current_series_labels,
            "missing_from_historical_headers": missing_from_historical,
            "missing_normalized_counterparties": missing_normalized_counterparties,
            "missing_from_data": missing_from_data,
            "segments_in_data": sorted(parsed_segments, key=str.casefold),
            "missing_expected_segments": missing_expected_segments,
            "canonical_key_by_series": canonical_key_by_series,
        }
    result = {
        "by_sheet": by_sheet,
        "gap_count": gap_count,
        "warnings": warnings,
        "missing_series": missing_series,
        "missing_segments": missing_segments,
    }
    if exceptions:
        result["exceptions"] = list(exceptions)
    return result


def _validate_reconciliation_parsed_data_by_sheet(
    *, parsed_data_by_sheet: Mapping[str, Mapping[str, Any]]
) -> None:
    for raw_sheet_name, raw_sections in parsed_data_by_sheet.items():
        sheet_name = str(raw_sheet_name)
        if not isinstance(raw_sections, Mapping):
            raise ParsedDataInvalidShapeError(
                f"Invalid parsed_data shape for sheet {sheet_name!r}: expected a mapping/object"
            )

        missing_sections = [
            section for section in ("totals", "futures") if section not in raw_sections
        ]
        if missing_sections:
            raise ParsedDataMissingKeyError(
                f"Missing required parsed_data section(s) for sheet {sheet_name!r}: "
                f"{', '.join(missing_sections)}"
            )

        for section_name in ("totals", "futures"):
            section_value = raw_sections[section_name]
            if not _is_supported_table_shape(section_value):
                raise ParsedDataInvalidShapeError(
                    f"Invalid parsed_data shape for sheet {sheet_name!r} section "
                    f"{section_name!r}: expected list of mappings or tabular object with "
                    "to_dict(orient='records')/to_records()"
                )


def _is_supported_table_shape(table: Any) -> bool:
    if isinstance(table, list):
        return all(isinstance(record, Mapping) for record in table)
    return hasattr(table, "to_dict") or hasattr(table, "to_records")


def _expected_segments_for_variant(
    *,
    variant: str | None,
    expected_segments_by_variant: Mapping[str, tuple[str, ...] | list[str] | set[str]] | None,
) -> set[str]:
    if not variant or not expected_segments_by_variant:
        return set()

    for key, values in expected_segments_by_variant.items():
        if str(key).strip().casefold() != variant.strip().casefold():
            continue
        return {str(value).strip() for value in values if str(value).strip()}
    return set()


def _extract_segments_from_records(parsed_sections: Mapping[str, Any]) -> set[str]:
    segments: set[str] = set()
    for table in parsed_sections.values():
        for record in _records(table):
            raw_segment = record.get("segment", record.get("Segment", ""))
            label = str(raw_segment).strip()
            if label:
                segments.add(label)
    return segments


def _counterparty_resolution_maps_from_records(
    totals_records: list[dict[str, Any]],
) -> tuple[dict[str, set[str]], dict[str, str]]:
    normalized_to_raw: dict[str, set[str]] = {}
    sources_by_raw_name: dict[str, str] = {}
    for record in totals_records:
        raw_name = str(record.get("counterparty", "")).strip()
        if not raw_name:
            continue
        resolution = normalize_counterparty_with_source(raw_name)
        normalized_to_raw.setdefault(resolution.canonical_name, set()).add(raw_name)
        sources_by_raw_name[raw_name] = resolution.source
    return normalized_to_raw, sources_by_raw_name


def _normalized_counterparties_from_records(
    totals_records: list[dict[str, Any]],
) -> dict[str, set[str]]:
    normalized_to_raw, _ = _counterparty_resolution_maps_from_records(totals_records)
    return normalized_to_raw


def _raw_counterparties_by_normalized_from_records(
    totals_records: list[dict[str, Any]],
) -> dict[str, set[str]]:
    normalized_to_raw: dict[str, set[str]] = {}
    for record in totals_records:
        raw_value = record.get("counterparty", "")
        if raw_value is None:
            continue
        raw_name = str(raw_value)
        if not raw_name.strip():
            continue
        normalized_name = normalize_counterparty(raw_name)
        normalized_to_raw.setdefault(normalized_name, set()).add(raw_name)
    return normalized_to_raw


def _normalized_counterparties_from_parsed_data(
    parsed_sections: Mapping[str, Any],
) -> dict[str, set[str]]:
    totals_records = _records(parsed_sections.get("totals", []))
    return _normalized_counterparties_from_records(totals_records)


def _raw_counterparties_by_normalized_from_parsed_data(
    parsed_sections: Mapping[str, Any],
) -> dict[str, set[str]]:
    totals_records = _records(parsed_sections.get("totals", []))
    return _raw_counterparties_by_normalized_from_records(totals_records)


def run_pipeline(config_path: str | Path) -> Path:
    """Run the Counter Risk pipeline and return the output run directory."""

    LOGGER.info("pipeline_start config_path=%s", config_path)
    try:
        config = load_config(config_path)
    except Exception as exc:
        LOGGER.exception("pipeline_failed stage=config_load config_path=%s", config_path)
        raise RuntimeError(f"Pipeline failed during config load: {config_path}") from exc

    try:
        _validate_pipeline_config(config)
    except Exception as exc:
        LOGGER.exception("pipeline_failed stage=config_validate config_path=%s", config_path)
        raise RuntimeError(f"Pipeline failed during config validation: {config_path}") from exc

    try:
        input_paths = _resolve_input_paths(config)
        _validate_input_files(input_paths)
    except Exception as exc:
        LOGGER.exception("pipeline_failed stage=input_validation config_path=%s", config_path)
        raise RuntimeError("Pipeline failed during input validation stage") from exc

    try:
        cprs_headers = _collect_cprs_header_candidates(config=config)
        as_of_date = derive_as_of_date(config, cprs_headers)
        run_date = derive_run_date(config)
        run_date_for_directory = config.run_date
    except Exception as exc:
        LOGGER.exception("pipeline_failed stage=date_derivation config_path=%s", config_path)
        raise RuntimeError("Pipeline failed during date derivation stage") from exc

    try:
        run_dir = _create_run_directory(as_of_date=as_of_date, run_date=run_date_for_directory)
    except Exception as exc:
        LOGGER.exception(
            "pipeline_failed stage=run_dir_setup run_dir=%s",
            _resolve_repo_root() / "runs" / as_of_date.isoformat(),
        )
        raise RuntimeError("Pipeline failed during run directory setup stage") from exc

    warnings: list[str] = []
    runtime_config = config
    try:
        runtime_config = _prepare_runtime_config(
            config=config,
            run_dir=run_dir,
            as_of_date=as_of_date,
            warnings=warnings,
        )
        parsed_by_variant = _parse_inputs(_resolve_input_paths(runtime_config))
        _validate_parsed_inputs(parsed_by_variant)
        _run_reconciliation_checks(
            run_dir=run_dir,
            config=runtime_config,
            parsed_by_variant=parsed_by_variant,
            warnings=warnings,
        )
    except Exception as exc:
        LOGGER.exception("pipeline_failed stage=parse_inputs run_dir=%s", run_dir)
        raise RuntimeError("Pipeline failed during parse stage") from exc

    try:
        top_exposures, top_changes_per_variant = _compute_metrics(parsed_by_variant)
    except Exception as exc:
        LOGGER.exception("pipeline_failed stage=compute_metrics run_dir=%s", run_dir)
        raise RuntimeError("Pipeline failed during compute stage") from exc

    try:
        historical_output_paths = _update_historical_outputs(
            run_dir=run_dir,
            config=runtime_config,
            parsed_by_variant=parsed_by_variant,
            as_of_date=as_of_date,
            warnings=warnings,
        )
    except Exception as exc:
        LOGGER.exception("pipeline_failed stage=historical_update run_dir=%s", run_dir)
        raise RuntimeError("Pipeline failed during historical update stage") from exc

    try:
        output_result = _write_outputs(
            run_dir=run_dir,
            config=runtime_config,
            as_of_date=as_of_date,
            warnings=warnings,
        )
    except Exception as exc:
        LOGGER.exception("pipeline_failed stage=write_outputs run_dir=%s", run_dir)
        raise RuntimeError("Pipeline failed during output write stage") from exc
    if isinstance(output_result, tuple):
        output_paths, ppt_result = output_result
    else:
        output_paths = output_result
        ppt_result = PptProcessingResult(status=PptProcessingStatus.SUCCESS)
    output_paths = historical_output_paths + output_paths

    try:
        input_hashes = {
            name: _sha256_file(path) for name, path in _resolve_input_paths(runtime_config).items()
        }
        manifest_builder = ManifestBuilder(config=config, as_of_date=as_of_date, run_date=run_date)
        output_paths_for_manifest = [path.relative_to(run_dir) for path in output_paths]
        manifest = manifest_builder.build(
            run_dir=run_dir,
            input_hashes=input_hashes,
            output_paths=output_paths_for_manifest,
            top_exposures=top_exposures,
            top_changes_per_variant=top_changes_per_variant,
            warnings=warnings,
            ppt_status=ppt_result.status.value,
        )
        manifest_path = manifest_builder.write(run_dir=run_dir, manifest=manifest)
    except Exception as exc:
        LOGGER.exception("pipeline_failed stage=manifest_write run_dir=%s", run_dir)
        raise RuntimeError("Pipeline failed during manifest generation stage") from exc

    LOGGER.info("pipeline_complete run_dir=%s manifest=%s", run_dir, manifest_path)

    return run_dir


def _create_run_directory(*, as_of_date: date, run_date: date | None = None) -> Path:
    runs_root = _resolve_repo_root() / "runs"
    base_name = (
        as_of_date.isoformat()
        if run_date is None
        else f"{as_of_date.isoformat()}__run_{run_date.isoformat()}"
    )
    candidate_names = [base_name, *(f"{base_name}_{index}" for index in range(1, 10_000))]

    for candidate_name in candidate_names:
        run_dir = runs_root / candidate_name
        if run_dir.exists():
            continue
        run_dir.mkdir(parents=True, exist_ok=False)
        return run_dir

    raise RuntimeError(f"Unable to create unique run directory for as_of_date {base_name}")


def _resolve_repo_root() -> Path:
    """Resolve the repository root for deterministic run output layout."""

    return Path(__file__).resolve().parents[3]


def _collect_cprs_header_candidates(*, config: WorkflowConfig) -> list[str]:
    header_candidates: list[str] = []
    candidate_paths = [
        config.mosers_all_programs_xlsx,
        config.mosers_ex_trend_xlsx,
        config.mosers_trend_xlsx,
    ]
    for path in candidate_paths:
        if path is None:
            continue
        header_candidates.extend(_extract_header_text_lines(path))
    return header_candidates


def _extract_header_text_lines(
    workbook_path: Path, *, max_rows: int = 15, max_cols: int = 6
) -> list[str]:
    try:
        from openpyxl import load_workbook
    except ModuleNotFoundError:
        return []
    except Exception:
        return []

    workbook = None
    try:
        workbook = load_workbook(filename=workbook_path, read_only=True, data_only=True)
        if not workbook.sheetnames:
            return []
        worksheet = workbook[workbook.sheetnames[0]]
        lines: list[str] = []
        for row in range(1, min(int(worksheet.max_row), max_rows) + 1):
            pieces: list[str] = []
            for col in range(1, max_cols + 1):
                value = worksheet.cell(row=row, column=col).value
                if value is None:
                    continue
                text = str(value).strip()
                if text:
                    pieces.append(text)
            if pieces:
                lines.append(" ".join(pieces))
        return lines
    except Exception:
        return []
    finally:
        if workbook is not None:
            workbook.close()


def _resolve_input_paths(config: WorkflowConfig) -> dict[str, Path]:
    paths: dict[str, Path] = {
        "mosers_ex_trend_xlsx": config.mosers_ex_trend_xlsx,
        "mosers_trend_xlsx": config.mosers_trend_xlsx,
        "hist_all_programs_3yr_xlsx": config.hist_all_programs_3yr_xlsx,
        "hist_ex_llc_3yr_xlsx": config.hist_ex_llc_3yr_xlsx,
        "hist_llc_3yr_xlsx": config.hist_llc_3yr_xlsx,
        "monthly_pptx": config.monthly_pptx,
    }
    if config.mosers_all_programs_xlsx is not None:
        paths["mosers_all_programs_xlsx"] = config.mosers_all_programs_xlsx
    if config.raw_nisa_all_programs_xlsx is not None:
        paths["raw_nisa_all_programs_xlsx"] = config.raw_nisa_all_programs_xlsx
    return paths


def _validate_pipeline_config(config: WorkflowConfig) -> None:
    if config.output_root.exists() and not config.output_root.is_dir():
        raise ValueError(f"output_root must be a directory path: {config.output_root}")

    if config.mosers_all_programs_xlsx is None and config.raw_nisa_all_programs_xlsx is None:
        raise ValueError(
            "One of mosers_all_programs_xlsx or raw_nisa_all_programs_xlsx must be configured"
        )
    if config.mosers_all_programs_xlsx is not None:
        _validate_extension(
            field_name="mosers_all_programs_xlsx",
            path=config.mosers_all_programs_xlsx,
            expected_suffix=".xlsx",
        )
    if config.raw_nisa_all_programs_xlsx is not None:
        _validate_extension(
            field_name="raw_nisa_all_programs_xlsx",
            path=config.raw_nisa_all_programs_xlsx,
            expected_suffix=".xlsx",
        )
    _validate_extension(
        field_name="mosers_ex_trend_xlsx",
        path=config.mosers_ex_trend_xlsx,
        expected_suffix=".xlsx",
    )
    _validate_extension(
        field_name="mosers_trend_xlsx",
        path=config.mosers_trend_xlsx,
        expected_suffix=".xlsx",
    )
    _validate_extension(
        field_name="hist_all_programs_3yr_xlsx",
        path=config.hist_all_programs_3yr_xlsx,
        expected_suffix=".xlsx",
    )
    _validate_extension(
        field_name="hist_ex_llc_3yr_xlsx",
        path=config.hist_ex_llc_3yr_xlsx,
        expected_suffix=".xlsx",
    )
    _validate_extension(
        field_name="hist_llc_3yr_xlsx",
        path=config.hist_llc_3yr_xlsx,
        expected_suffix=".xlsx",
    )
    _validate_extension(
        field_name="monthly_pptx",
        path=config.monthly_pptx,
        expected_suffix=".pptx",
    )


def _validate_extension(*, field_name: str, path: Path, expected_suffix: str) -> None:
    if path.suffix.lower() != expected_suffix:
        raise ValueError(
            f"Invalid file type for {field_name}: expected {expected_suffix}, got '{path.suffix}'"
        )


def _validate_input_files(input_paths: dict[str, Path]) -> None:
    missing = [f"{name}: {path}" for name, path in input_paths.items() if not path.exists()]
    if missing:
        details = "; ".join(missing)
        raise FileNotFoundError(f"Missing pipeline input files: {details}")


def _prepare_runtime_config(
    *,
    config: WorkflowConfig,
    run_dir: Path,
    as_of_date: date,
    warnings: list[str],
) -> WorkflowConfig:
    raw_nisa_path = config.raw_nisa_all_programs_xlsx
    if raw_nisa_path is None:
        return config

    generated_dir = run_dir / "_generated"
    generated_dir.mkdir(parents=True, exist_ok=True)
    generated_mosers_path = generated_dir / "all_programs-generated-mosers.xlsx"
    canonical_mosers_path = run_dir / "all_programs-mosers-input.xlsx"
    generate_mosers_workbook(
        raw_nisa_path=raw_nisa_path,
        output_path=generated_mosers_path,
        as_of_date=as_of_date,
    )
    shutil.copy2(generated_mosers_path, canonical_mosers_path)
    warnings.append("Generated All Programs MOSERS workbook from raw NISA input")
    return config.model_copy(update={"mosers_all_programs_xlsx": canonical_mosers_path})


def _parse_inputs(input_paths: dict[str, Path]) -> dict[str, dict[str, Any]]:
    variants = [
        ("all_programs", input_paths["mosers_all_programs_xlsx"]),
        ("ex_trend", input_paths["mosers_ex_trend_xlsx"]),
        ("trend", input_paths["mosers_trend_xlsx"]),
    ]

    parsed: dict[str, dict[str, Any]] = {}
    for variant, workbook_path in variants:
        LOGGER.info("parse_start variant=%s file=%s", variant, workbook_path)
        totals_df = parse_fcm_totals(workbook_path)
        futures_df = parse_futures_detail(workbook_path)
        parsed[variant] = {
            "totals": totals_df,
            "futures": futures_df,
        }
        LOGGER.info(
            "parse_complete variant=%s totals_rows=%s futures_rows=%s",
            variant,
            _row_count(totals_df),
            _row_count(futures_df),
        )

    return parsed


def _validate_parsed_inputs(parsed_by_variant: dict[str, dict[str, Any]]) -> None:
    missing_variants = [
        variant for variant in _EXPECTED_VARIANTS if variant not in parsed_by_variant
    ]
    if missing_variants:
        raise ValueError(f"Missing parsed variants: {', '.join(missing_variants)}")

    for variant in _EXPECTED_VARIANTS:
        parsed_sections = parsed_by_variant[variant]
        if not isinstance(parsed_sections, Mapping):
            raise ValueError(f"Parsed payload for variant '{variant}' must be a mapping")

        missing_sections = [
            section for section in ("totals", "futures") if section not in parsed_sections
        ]
        if missing_sections:
            raise ValueError(
                f"Parsed payload for variant '{variant}' is missing sections: "
                f"{', '.join(missing_sections)}"
            )

        totals_columns = _column_names(parsed_sections["totals"])
        futures_columns = _column_names(parsed_sections["futures"])

        _require_columns(
            section_name=f"{variant}.totals",
            columns=totals_columns,
            required_columns=_REQUIRED_TOTAL_COLUMNS,
        )
        _require_columns(
            section_name=f"{variant}.futures",
            columns=futures_columns,
            required_columns=_REQUIRED_FUTURES_COLUMNS,
        )


def _require_columns(
    *, section_name: str, columns: set[str], required_columns: tuple[str, ...]
) -> None:
    missing_columns = [column for column in required_columns if column not in columns]
    if missing_columns:
        raise ValueError(
            f"Parsed section '{section_name}' is missing required columns: "
            f"{', '.join(missing_columns)}"
        )


def _compute_metrics(
    parsed_by_variant: dict[str, dict[str, Any]],
) -> tuple[dict[str, list[dict[str, Any]]], dict[str, list[dict[str, Any]]]]:
    LOGGER.info("compute_start")
    top_exposures: dict[str, list[dict[str, Any]]] = {}
    top_changes_per_variant: dict[str, list[dict[str, Any]]] = {}

    for variant, parsed in parsed_by_variant.items():
        totals_df = parsed["totals"]
        totals_records = _records(totals_df)

        if not totals_records:
            top_exposures[variant] = []
            top_changes_per_variant[variant] = []
            continue

        sorted_exposures = sorted(
            totals_records,
            key=lambda record: abs(float(record.get("Notional", 0.0) or 0.0)),
            reverse=True,
        )
        top_exposures[variant] = [
            {
                "counterparty": str(record.get("counterparty", "")),
                "notional": float(record.get("Notional", 0.0) or 0.0),
            }
            for record in sorted_exposures[:5]
        ]

        change_column = "NotionalChange"
        if not all(change_column in record for record in totals_records):
            top_changes_per_variant[variant] = []
            continue

        sorted_changes = sorted(
            totals_records,
            key=lambda record: abs(float(record.get(change_column, 0.0) or 0.0)),
            reverse=True,
        )
        top_changes_per_variant[variant] = [
            {
                "counterparty": str(record.get("counterparty", "")),
                "notional_change": float(record.get(change_column, 0.0) or 0.0),
            }
            for record in sorted_changes[:5]
        ]

    LOGGER.info("compute_complete")
    return top_exposures, top_changes_per_variant


def _write_outputs(
    *, run_dir: Path, config: WorkflowConfig, as_of_date: date, warnings: list[str]
) -> tuple[list[Path], PptProcessingResult]:
    LOGGER.info("write_outputs_start run_dir=%s", run_dir)

    variant_inputs = [
        _VariantInputs(
            name="all_programs",
            workbook_path=_require_path(
                config.mosers_all_programs_xlsx, field_name="mosers_all_programs_xlsx"
            ),
            historical_path=config.hist_all_programs_3yr_xlsx,
        ),
        _VariantInputs(
            name="ex_trend",
            workbook_path=config.mosers_ex_trend_xlsx,
            historical_path=config.hist_ex_llc_3yr_xlsx,
        ),
        _VariantInputs(
            name="trend",
            workbook_path=config.mosers_trend_xlsx,
            historical_path=config.hist_llc_3yr_xlsx,
        ),
    ]

    output_paths: list[Path] = []
    for variant_input in variant_inputs:
        source_mosers = variant_input.workbook_path
        target_monthly_book = run_dir / f"{variant_input.name}-mosers-input.xlsx"
        if source_mosers.resolve() != target_monthly_book.resolve():
            shutil.copy2(source_mosers, target_monthly_book)
        output_paths.append(target_monthly_book)

    if not config.ppt_output_enabled:
        LOGGER.info("write_outputs_skip_ppt run_dir=%s", run_dir)
        return output_paths, PptProcessingResult(status=PptProcessingStatus.SKIPPED)

    output_names = resolve_ppt_output_names(as_of_date)
    target_distribution_ppt = run_dir / output_names.distribution_filename
    output_context = OutputContext(
        config=config,
        run_dir=run_dir,
        as_of_date=as_of_date,
        run_date=config.run_date or as_of_date,
        warnings=tuple(warnings),
    )
    screenshot_generator = _build_ppt_screenshot_output_generator(
        warnings=warnings,
        ppt_output_names_resolver=resolve_ppt_output_names,
    )
    master_output_paths = list(screenshot_generator.generate(context=output_context))
    if len(master_output_paths) != 1:
        raise RuntimeError(
            "PPT screenshot output generator must produce exactly one Master PPT output"
        )
    target_master_ppt = master_output_paths[0]
    output_paths.extend(master_output_paths)
    readme_ppt_outputs = RunFolderReadmePptOutputs(
        master=target_master_ppt.relative_to(run_dir),
        distribution=target_distribution_ppt.relative_to(run_dir),
    )

    link_refresh_generator = _build_ppt_link_refresh_output_generator(warnings=warnings)
    link_refresh_generator.generate(context=output_context)
    refresh_result = _to_ppt_processing_result(link_refresh_generator.last_result)

    if refresh_result.status == PptProcessingStatus.FAILED:
        LOGGER.warning(
            "Skipping distribution PPT derivation because Master PPT refresh failed: %s",
            target_master_ppt,
        )
    else:
        chart_replaced_ppt = run_dir / f"{target_master_ppt.stem}_chart_replaced.pptx"
        chart_replacement_applied = _apply_chart_replacement(
            master_pptx_path=target_master_ppt,
            output_path=chart_replaced_ppt,
            run_dir=run_dir,
            static_mode=config.distribution_static,
            warnings=warnings,
        )
        distribution_source = chart_replaced_ppt if chart_replacement_applied else target_master_ppt
        _derive_distribution_ppt(
            master_pptx_path=distribution_source,
            distribution_pptx_path=target_distribution_ppt,
        )
        try:
            distribution_validation = validate_distribution_ppt_standalone(target_distribution_ppt)
        except RuntimeError as exc:
            warnings.append(
                "Distribution PPT standalone validation skipped; unable to parse generated deck"
            )
            LOGGER.warning(
                "Distribution PPT standalone validation skipped for %s: %s",
                target_distribution_ppt,
                exc,
            )
        else:
            if not distribution_validation.is_valid:
                rel_parts = ", ".join(distribution_validation.external_relationship_parts)
                raise RuntimeError(
                    "Distribution PPT standalone validation failed; "
                    f"found {distribution_validation.external_relationship_count} "
                    f"external relationships in: {rel_parts}"
                )
        output_paths.append(target_distribution_ppt)
        distribution_pdf_path = _export_distribution_pdf(
            source_pptx=target_distribution_ppt,
            run_dir=run_dir,
            config=config,
            as_of_date=as_of_date,
            warnings=warnings,
        )
        if distribution_pdf_path is not None:
            output_paths.append(distribution_pdf_path)
    static_output_paths = _create_static_distribution(
        source_pptx=target_master_ppt,
        run_dir=run_dir,
        config=config,
        warnings=warnings,
    )
    output_paths.extend(static_output_paths)
    if refresh_result.status == PptProcessingStatus.SUCCESS:
        readme_path = run_dir / "README.txt"
        readme_path.write_text(
            build_run_folder_readme_content(as_of_date, readme_ppt_outputs),
            encoding="utf-8",
        )
        output_paths.append(readme_path)

    LOGGER.info("write_outputs_complete output_count=%s", len(output_paths))
    return output_paths, refresh_result


def _apply_chart_replacement(
    *,
    master_pptx_path: Path,
    output_path: Path,
    run_dir: Path,
    static_mode: bool,
    warnings: list[str],
) -> bool:
    """Replace chart/OLE shapes with static images using confidence-based matching.

    Opens a COM session to export chart shapes and fallback slide images, then
    runs ``_rebuild_pptx_replacing_charts`` with confidence checks.  Returns
    ``True`` when chart replacement was applied, ``False`` when skipped.
    """

    if platform.system().lower() != "windows":
        return False

    try:
        import win32com.client
    except ImportError:
        return False

    chart_images_dir = run_dir / "_chart_images"
    chart_images_dir.mkdir(parents=True, exist_ok=True)

    app = None
    presentation = None
    chart_images: dict[tuple[int, str], Path] = {}
    fallback_slide_images: dict[int, Path] = {}
    try:
        app = win32com.client.DispatchEx("PowerPoint.Application")
        app.Visible = False
        presentation = app.Presentations.Open(str(master_pptx_path), WithWindow=False)

        chart_images = _export_chart_shapes_as_images(
            com_presentation=presentation,
            slide_images_dir=chart_images_dir,
            warnings=warnings,
        )

        if chart_images:
            slide_count = int(presentation.Slides.Count)
            for slide_idx in range(1, slide_count + 1):
                img_path = chart_images_dir / f"fallback_slide_{slide_idx:04d}.png"
                try:
                    presentation.Slides[slide_idx].Export(str(img_path), "PNG")
                    fallback_slide_images[slide_idx] = img_path
                except Exception as exc:
                    warnings.append(
                        f"chart_replacement fallback slide export failed (slide {slide_idx}): {exc}"
                    )
    except Exception as exc:
        warnings.append(f"chart_replacement COM session failed: {exc}")
        LOGGER.warning("chart_replacement_com_failed exc=%s", exc)
        return False
    finally:
        if presentation is not None:
            with contextlib.suppress(Exception):
                presentation.Close()
        if app is not None:
            with contextlib.suppress(Exception):
                app.Quit()

    if not chart_images:
        return False

    try:
        _rebuild_pptx_replacing_charts(
            source_pptx=master_pptx_path,
            output_path=output_path,
            chart_images=chart_images,
            fallback_slide_images=fallback_slide_images,
            fallback_to_full_deck_rebuild=static_mode,
        )
        LOGGER.info("chart_replacement_complete output=%s", output_path)
        return True
    except RuntimeError as exc:
        if static_mode and "full-deck static rebuild" in str(exc):
            LOGGER.info("chart_replacement_deferred_to_static_rebuild: %s", exc)
        else:
            LOGGER.warning("chart_replacement_failed exc=%s", exc)
            warnings.append(f"chart_replacement failed: {exc}")
        return False


def _derive_distribution_ppt(*, master_pptx_path: Path, distribution_pptx_path: Path) -> None:
    """Derive the distribution PPT from the generated Master PPT."""

    distribution_pptx_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        scrub_external_relationships_from_pptx(
            master_pptx_path,
            scrubbed_pptx_path=distribution_pptx_path,
        )
    except BadZipFile:
        shutil.copy2(master_pptx_path, distribution_pptx_path)
        LOGGER.debug(
            "Distribution derivation skipped link stripping for non-standard PPTX: %s",
            master_pptx_path,
        )
        return

    remaining_targets = _list_external_link_targets(distribution_pptx_path)
    if remaining_targets:
        targets_list = ", ".join(sorted(remaining_targets))
        raise RuntimeError(
            "Distribution PPT derivation did not remove all external link targets. "
            f"Remaining targets: {targets_list}"
        )


def _assert_master_preserves_external_link_targets(
    *, source_pptx_path: Path, master_pptx_path: Path
) -> None:
    source_targets = _list_external_link_targets(source_pptx_path)
    if not source_targets:
        return

    master_targets = _list_external_link_targets(master_pptx_path)
    if source_targets.issubset(master_targets):
        return

    missing_targets = sorted(source_targets.difference(master_targets))
    missing_list = ", ".join(missing_targets)
    raise RuntimeError(
        "Master PPT generation did not preserve linked chart references. "
        f"Missing external link targets: {missing_list}"
    )


def _list_external_link_targets(pptx_path: Path) -> set[str]:
    return list_external_relationship_targets(pptx_path)


def _resolve_screenshot_input_mapping(config: WorkflowConfig) -> dict[str, Path]:
    raw_inputs = config.screenshot_inputs
    normalized_pairs: list[tuple[str, Path]] = []
    seen_keys: set[str] = set()
    for raw_key, raw_path in raw_inputs.items():
        key = str(raw_key).strip()
        if not key:
            raise ValueError("Screenshot input keys must be non-empty")
        if key in seen_keys:
            raise ValueError(f"Screenshot input key is duplicated after normalization: {key!r}")
        seen_keys.add(key)

        image_path = Path(raw_path).expanduser().resolve()
        if image_path.suffix.lower() != ".png":
            raise ValueError(f"Screenshot input '{key}' must point to a PNG file: {image_path}")
        if not image_path.exists() or not image_path.is_file():
            raise FileNotFoundError(f"Screenshot input '{key}' file does not exist: {image_path}")

        normalized_pairs.append((key, image_path))

    normalized = dict(sorted(normalized_pairs, key=lambda item: item[0]))
    if config.enable_screenshot_replacement and not normalized:
        raise ValueError("Screenshot replacement is enabled but no screenshot_inputs were provided")
    return normalized


def _get_screenshot_replacer(implementation: str) -> ScreenshotReplacer:
    if implementation == "zip":
        return _replace_screenshots_with_zip_backend
    if implementation == "python-pptx":
        return _replace_screenshots_with_python_pptx_backend
    raise ValueError(f"Unsupported screenshot replacement implementation: {implementation}")


def _replace_screenshots_with_python_pptx_backend(
    source_pptx_path: Path, output_pptx_path: Path, screenshot_inputs: dict[str, Path]
) -> None:
    from counter_risk.writers.pptx_screenshots import replace_screenshot_pictures

    replace_screenshot_pictures(
        pptx_in=source_pptx_path,
        images_by_section=screenshot_inputs,
        pptx_out=output_pptx_path,
    )


def _replace_screenshots_with_zip_backend(
    source_pptx_path: Path, output_pptx_path: Path, screenshot_inputs: dict[str, Path]
) -> None:
    from counter_risk.ppt.replace_screenshots import (
        ScreenshotReplacement,
        replace_screenshots_in_pptx,
    )

    replacements = [
        ScreenshotReplacement(
            slide_number=slide_number,
            picture_index=picture_index,
            image_bytes=image_path.read_bytes(),
        )
        for screenshot_key, image_path in screenshot_inputs.items()
        for slide_number, picture_index in [_parse_zip_screenshot_key(screenshot_key)]
    ]
    replace_screenshots_in_pptx(
        source_pptx_path=source_pptx_path,
        output_pptx_path=output_pptx_path,
        replacements=replacements,
    )


def _parse_zip_screenshot_key(key: str) -> tuple[int, int]:
    slide_token, separator, picture_token = key.partition(":")
    normalized_slide = slide_token.strip().lower()
    if normalized_slide.startswith("slide"):
        normalized_slide = normalized_slide.removeprefix("slide").strip()
    if not normalized_slide.isdigit():
        raise ValueError(f"Invalid screenshot key for zip backend: {key!r}")

    picture_index = 0
    if separator:
        picture_candidate = picture_token.strip()
        if not picture_candidate.isdigit():
            raise ValueError(f"Invalid screenshot key for zip backend: {key!r}")
        picture_index = int(picture_candidate)

    return int(normalized_slide), picture_index


def _refresh_ppt_links(pptx_path: Path) -> PptProcessingResult:
    """Best-effort refresh of linked PowerPoint content via COM automation."""

    if platform.system().lower() != "windows":
        LOGGER.info("ppt_link_refresh_skipped file=%s reason=unsupported_platform", pptx_path)
        return PptProcessingResult(
            status=PptProcessingStatus.SKIPPED,
            error_detail="unsupported platform",
        )

    try:
        import win32com.client
    except ImportError:
        LOGGER.info("ppt_link_refresh_skipped file=%s reason=win32com_unavailable", pptx_path)
        return PptProcessingResult(
            status=PptProcessingStatus.SKIPPED,
            error_detail="win32com unavailable",
        )

    app = None
    presentation = None
    try:
        app = win32com.client.DispatchEx("PowerPoint.Application")
        app.Visible = False
        presentation = app.Presentations.Open(str(pptx_path), WithWindow=False)
        presentation.UpdateLinks()
        presentation.Save()
        LOGGER.info("ppt_link_refresh_complete file=%s", pptx_path)
        return PptProcessingResult(status=PptProcessingStatus.SUCCESS)
    except Exception as exc:
        LOGGER.exception("ppt_link_refresh_failed file=%s", pptx_path)
        raise RuntimeError(f"PPT link refresh failed for '{pptx_path}': {exc}") from exc
    finally:
        if presentation is not None:
            presentation.Close()
        if app is not None:
            app.Quit()


def _create_static_distribution(
    *,
    source_pptx: Path,
    run_dir: Path,
    config: WorkflowConfig,
    warnings: list[str],
) -> list[Path]:
    """Produce a static distribution copy of the presentation.

    On Windows with COM available: exports each slide as a PNG image and rebuilds
    the deck as a flat image-only PPTX (no live Excel links).

    When COM is unavailable (non-Windows or missing pywin32): appends a human-readable
    warning to *warnings* and returns an empty list so the caller's regular
    (non-static) PPT remains the sole deliverable.
    """

    if not config.distribution_static:
        return []

    if platform.system().lower() != "windows":
        warnings.append(
            "distribution_static requested but PowerPoint COM is only available on Windows; "
            "no static distribution produced"
        )
        LOGGER.info("distribution_static_skipped reason=non_windows")
        return []

    try:
        import win32com.client
    except ImportError:
        warnings.append(
            "distribution_static requested but win32com is not installed; "
            "no static distribution produced"
        )
        LOGGER.info("distribution_static_skipped reason=win32com_unavailable")
        return []

    output_paths: list[Path] = []
    slide_images_dir = run_dir / "_distribution_slides"
    static_pptx_path = run_dir / f"{source_pptx.stem}_distribution_static.pptx"
    app = None
    presentation = None
    try:
        app = win32com.client.DispatchEx("PowerPoint.Application")
        app.Visible = False
        presentation = app.Presentations.Open(str(source_pptx), WithWindow=False)

        # Preferred: export each slide as a PNG and rebuild the entire deck as
        # one picture per slide, removing all live chart/OLE objects.
        slide_images_dir.mkdir(parents=True, exist_ok=True)
        slide_images = _export_slides_as_images(
            com_presentation=presentation,
            slide_images_dir=slide_images_dir,
            warnings=warnings,
        )
        _rebuild_pptx_from_slide_images(
            source_pptx=source_pptx,
            output_path=static_pptx_path,
            slide_images=slide_images,
        )
        output_paths.append(static_pptx_path)
        LOGGER.info("distribution_static_pptx_complete path=%s", static_pptx_path)

    except Exception as exc:
        warnings.append(f"distribution_static generation failed: {exc}")
        LOGGER.exception("distribution_static_failed")
    finally:
        if presentation is not None:
            with contextlib.suppress(Exception):
                presentation.Close()
        if app is not None:
            with contextlib.suppress(Exception):
                app.Quit()

    return output_paths


def _export_distribution_pdf(
    *,
    source_pptx: Path,
    run_dir: Path,
    config: WorkflowConfig,
    as_of_date: date,
    warnings: list[str],
) -> Path | None:
    output_context = OutputContext(
        config=config,
        run_dir=run_dir,
        as_of_date=as_of_date,
        run_date=config.run_date or as_of_date,
        warnings=tuple(warnings),
    )
    pdf_output_generator = _build_pdf_export_output_generator(
        source_pptx=source_pptx,
        warnings=warnings,
    )
    generated_paths = pdf_output_generator.generate(context=output_context)
    if not generated_paths:
        return None

    if len(generated_paths) != 1:
        raise RuntimeError("PDF export output generator must produce at most one PDF output")

    return generated_paths[0]


def _export_pptx_to_pdf(*, source_pptx: Path, pdf_path: Path) -> None:
    from counter_risk.outputs.pdf_export import export_pptx_to_pdf_via_com

    export_pptx_to_pdf_via_com(source_pptx=source_pptx, pdf_path=pdf_path)


def _export_slides_as_images(
    *,
    com_presentation: Any,
    slide_images_dir: Path,
    warnings: list[str],
) -> list[Path]:
    """Export each slide in the COM presentation as a PNG image."""
    slide_images: list[Path] = []
    slide_count = int(com_presentation.Slides.Count)
    for slide_idx in range(1, slide_count + 1):
        image_path = slide_images_dir / f"slide_{slide_idx:04d}.png"
        try:
            com_presentation.Slides[slide_idx].Export(str(image_path), "PNG")
            slide_images.append(image_path)
        except Exception as exc:
            warnings.append(f"distribution_static slide export failed (slide {slide_idx}): {exc}")
            LOGGER.warning(
                "distribution_static_slide_export_failed slide=%d exc=%s", slide_idx, exc
            )
            raise
    return slide_images


def _export_chart_shapes_as_images(
    *,
    com_presentation: Any,
    slide_images_dir: Path,
    warnings: list[str],
) -> dict[tuple[int, str], Path]:
    """Export OLE/chart shapes from a COM presentation to individual PNG files.

    Iterates every slide in *com_presentation* and exports shapes whose COM
    type is ``msoEmbeddedOLEObject`` (7), ``msoLinkedOLEObject`` (10), or
    whose ``HasChart`` property is ``True``.  The resulting PNG files are
    written into *slide_images_dir*.

    Returns a mapping of ``(1-based slide index, shape name)`` to the exported
    image ``Path``.  Shapes that fail to export are recorded in *warnings* and
    omitted from the result so the caller can decide how to handle them.
    """
    chart_images: dict[tuple[int, str], Path] = {}
    slide_count: int = com_presentation.Slides.Count
    for slide_idx in range(1, slide_count + 1):
        com_slide = com_presentation.Slides[slide_idx]
        for com_shape in com_slide.Shapes:
            shape_type: int = com_shape.Type
            has_chart = False
            with contextlib.suppress(Exception):
                has_chart = bool(com_shape.HasChart)
            if (
                shape_type not in (_COM_SHAPE_TYPE_EMBEDDED_OLE, _COM_SHAPE_TYPE_LINKED_OLE)
                and not has_chart
            ):
                continue
            img_path = slide_images_dir / f"chart_{slide_idx:04d}_{com_shape.Id}.png"
            try:
                com_shape.Export(str(img_path), _COM_SHAPE_FORMAT_PNG)
                chart_images[(slide_idx, com_shape.Name)] = img_path
                LOGGER.info(
                    "chart_shape_exported slide=%d shape=%s path=%s",
                    slide_idx,
                    com_shape.Name,
                    img_path,
                )
            except Exception as exc:
                warnings.append(
                    f"distribution_static chart export failed "
                    f"(slide {slide_idx}, shape '{com_shape.Name}'): {exc}"
                )
                LOGGER.warning(
                    "chart_shape_export_failed slide=%d shape=%s exc=%s",
                    slide_idx,
                    com_shape.Name,
                    exc,
                )
    return chart_images


def _rebuild_pptx_replacing_charts(
    *,
    source_pptx: Path,
    output_path: Path,
    chart_images: dict[tuple[int, str], Path],
    fallback_slide_images: dict[int, Path] | None = None,
    fallback_to_full_deck_rebuild: bool = False,
) -> None:
    """Create a PPTX replacing OLE/chart shapes with static PNG images.

    For each entry in *chart_images* the corresponding shape is removed from
    its slide and replaced with a ``Picture`` shape at the exact same position
    and size.  All other shapes (titles, text boxes, decorative elements) are
    left untouched so the deck remains editable and titles stay selectable.

    When *chart_images* is empty (no chart shapes found or all exports failed)
    the source presentation is saved unchanged to *output_path*, ensuring a
    deliverable is always produced.
    """
    from pptx import Presentation

    fallback_slide_images = fallback_slide_images or {}
    prs = Presentation(str(source_pptx))
    low_confidence_slides: set[int] = set()
    unresolved_low_confidence_slides: set[int] = set()

    for slide_idx, slide in enumerate(prs.slides, start=1):
        replacements: list[tuple[Path, int, int, int, int]] = []
        shapes_to_remove: list[Any] = []
        replacement_requests = {
            shape_name: image_path
            for (shape_slide_idx, shape_name), image_path in chart_images.items()
            if shape_slide_idx == slide_idx
        }

        for shape_name, img_path in replacement_requests.items():
            candidate_shapes = [
                shape
                for shape in list(slide.shapes)
                if str(getattr(shape, "name", "")) == shape_name
            ]
            confidence = _shape_match_confidence(
                target_name=shape_name,
                candidate_shapes=candidate_shapes,
            )
            if confidence < _SHAPE_MATCH_CONFIDENCE_THRESHOLD:
                low_confidence_slides.add(slide_idx)
                LOGGER.warning(
                    "chart_replace_low_confidence slide=%d shape=%s confidence=%.2f",
                    slide_idx,
                    shape_name,
                    confidence,
                )
                continue

            matched_shape = candidate_shapes[0]
            left, top, width, height = (
                matched_shape.left,
                matched_shape.top,
                matched_shape.width,
                matched_shape.height,
            )
            shapes_to_remove.append(matched_shape)
            replacements.append((img_path, left, top, width, height))

        for shape in shapes_to_remove:
            sp_tree = shape.element.getparent()
            sp_tree.remove(shape.element)

        for img_path, left, top, width, height in replacements:
            slide.shapes.add_picture(str(img_path), left, top, width, height)

        if fallback_to_full_deck_rebuild:
            continue

        if slide_idx in low_confidence_slides:
            fallback_image = fallback_slide_images.get(slide_idx)
            if fallback_image is None:
                unresolved_low_confidence_slides.add(slide_idx)
            else:
                _replace_slide_with_image(slide=slide, slide_image=fallback_image)

    if unresolved_low_confidence_slides:
        unresolved_list = ", ".join(
            str(index) for index in sorted(unresolved_low_confidence_slides)
        )
        raise RuntimeError(
            "Chart replacement confidence check failed and no slide-image fallback was provided "
            f"for slide(s): {unresolved_list}"
        )

    if low_confidence_slides and fallback_to_full_deck_rebuild:
        low_confidence_list = ", ".join(str(index) for index in sorted(low_confidence_slides))
        raise RuntimeError(
            "Chart replacement confidence check failed; full-deck static rebuild required "
            f"for slide(s): {low_confidence_list}"
        )

    prs.save(str(output_path))


def _shape_match_confidence(*, target_name: str, candidate_shapes: list[Any]) -> float:
    """Score chart shape matching confidence using name uniqueness and optional geometry.

    Name uniqueness is weighted heavily because duplicate names on a slide make
    deterministic replacement unreliable. When one candidate is present and no
    geometry information is available, confidence remains high.
    """

    if not candidate_shapes:
        return 0.0

    unique_name_score = 1.0 if len(candidate_shapes) == 1 else 0.2
    position_scores = [
        _shape_position_confidence(shape)
        for shape in candidate_shapes
        if str(getattr(shape, "name", "")) == target_name
    ]
    position_score = max(position_scores, default=0.0)
    return min(1.0, (0.75 * unique_name_score) + (0.25 * position_score))


def _shape_position_confidence(shape: Any) -> float:
    """Return a bounded confidence based on shape geometry validity."""

    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    if width <= 0 or height <= 0:
        return 0.0

    left = int(getattr(shape, "left", 0) or 0)
    top = int(getattr(shape, "top", 0) or 0)
    if left < 0 or top < 0:
        return 0.5
    return 1.0


def _replace_slide_with_image(*, slide: Any, slide_image: Path) -> None:
    """Fallback for low-confidence matching: flatten the full slide to one image."""

    width = slide.part.slide_layout.part.package.presentation_part.presentation.slide_width
    height = slide.part.slide_layout.part.package.presentation_part.presentation.slide_height

    for shape in list(slide.shapes):
        sp_tree = shape.element.getparent()
        sp_tree.remove(shape.element)
    slide.shapes.add_picture(str(slide_image), left=0, top=0, width=width, height=height)


def _update_historical_outputs(
    *,
    run_dir: Path,
    config: WorkflowConfig,
    parsed_by_variant: dict[str, dict[str, Any]],
    as_of_date: date,
    warnings: list[str],
) -> list[Path]:
    LOGGER.info("historical_update_start run_dir=%s as_of_date=%s", run_dir, as_of_date.isoformat())
    output_generator = _build_historical_workbook_output_generator(
        parsed_by_variant=parsed_by_variant,
        warnings=warnings,
    )
    output_context = OutputContext(
        config=config,
        run_dir=run_dir,
        as_of_date=as_of_date,
        run_date=config.run_date or as_of_date,
        warnings=tuple(warnings),
    )
    output_paths = list(output_generator.generate(context=output_context))

    LOGGER.info("historical_update_complete workbook_count=%s", len(output_paths))
    return output_paths


def _build_historical_workbook_output_generator(
    *,
    parsed_by_variant: dict[str, dict[str, Any]],
    warnings: list[str],
) -> OutputGenerator:
    from counter_risk.outputs.historical_workbook import HistoricalWorkbookOutputGenerator

    return HistoricalWorkbookOutputGenerator(
        parsed_by_variant=parsed_by_variant,
        warnings=warnings,
        workbook_merger=_merge_historical_workbook,
        records_extractor=_records,
    )


def _build_ppt_screenshot_output_generator(
    *,
    warnings: list[str],
    ppt_output_names_resolver: Callable[[date], PptOutputNames],
) -> OutputGenerator:
    from counter_risk.outputs.ppt_screenshot import PptScreenshotOutputGenerator

    return PptScreenshotOutputGenerator(
        warnings=warnings,
        screenshot_input_mapping_resolver=_resolve_screenshot_input_mapping,
        screenshot_replacer_resolver=_get_screenshot_replacer,
        master_link_target_validator=_assert_master_preserves_external_link_targets,
        ppt_output_names_resolver=ppt_output_names_resolver,
    )


def _build_ppt_link_refresh_output_generator(
    *, warnings: list[str]
) -> PptLinkRefreshOutputGenerator:
    from counter_risk.outputs.ppt_link_refresh import PptLinkRefreshOutputGenerator

    return PptLinkRefreshOutputGenerator(
        warnings=warnings,
        ppt_link_refresher=_refresh_ppt_links,
    )


def _build_pdf_export_output_generator(
    *, source_pptx: Path, warnings: list[str]
) -> OutputGenerator:
    from counter_risk.outputs.pdf_export import PDFExportGenerator

    return PDFExportGenerator(
        source_pptx=source_pptx,
        warnings=warnings,
        pptx_to_pdf_exporter=lambda src, dst: _export_pptx_to_pdf(source_pptx=src, pdf_path=dst),
    )


def _to_ppt_processing_result(refresh_result: object | None) -> PptProcessingResult:
    if refresh_result is None:
        raise RuntimeError("PPT link refresh output generator did not record a refresh result")

    status_value = getattr(refresh_result, "status", None)
    if status_value is None:
        raise TypeError("PPT link refresh result is missing a status")

    status_text = str(getattr(status_value, "value", status_value)).strip().lower()
    try:
        status = PptProcessingStatus(status_text)
    except ValueError as exc:
        raise ValueError(f"Unsupported PPT processing status: {status_text!r}") from exc

    error_detail = getattr(refresh_result, "error_detail", None)
    normalized_error = None if error_detail is None else str(error_detail)
    return PptProcessingResult(status=status, error_detail=normalized_error)


def _merge_historical_workbook(
    *,
    workbook_path: Path,
    variant: str,
    as_of_date: date,
    totals_records: list[dict[str, Any]],
    warnings: list[str],
) -> None:
    try:
        from openpyxl import load_workbook
    except ImportError:
        message = (
            f"Historical workbook update skipped for variant '{variant}'; openpyxl unavailable"
        )
        LOGGER.warning("historical_update_skipped variant=%s reason=openpyxl_unavailable", variant)
        warnings.append(message)
        return

    workbook = None
    try:
        workbook = load_workbook(filename=workbook_path)
        preferred_sheet = _PREFERRED_HISTORICAL_SHEET_BY_VARIANT.get(variant)
        worksheet = _select_historical_worksheet(
            workbook=workbook,
            preferred_sheet_name=preferred_sheet,
        )
        _validate_historical_headers(worksheet=worksheet)
        append_row = int(getattr(worksheet, "max_row", 0)) + 1

        total_notional = sum(float(record.get("Notional", 0.0) or 0.0) for record in totals_records)
        counterparties = len(
            {
                str(record.get("counterparty", "")).strip()
                for record in totals_records
                if str(record.get("counterparty", "")).strip()
            }
        )

        worksheet.cell(row=append_row, column=1).value = as_of_date
        worksheet.cell(row=append_row, column=2).value = total_notional
        worksheet.cell(row=append_row, column=3).value = counterparties
        workbook.save(workbook_path)
        LOGGER.info(
            "historical_update_variant_complete variant=%s row=%s notional=%s counterparties=%s",
            variant,
            append_row,
            total_notional,
            counterparties,
        )
    except Exception as exc:
        raise RuntimeError(
            f"Failed to update historical workbook for variant '{variant}': {workbook_path}"
        ) from exc
    finally:
        if workbook is not None:
            workbook.close()


def _normalize_header(value: Any) -> str:
    if value is None:
        return ""
    return " ".join(str(value).split()).casefold()


def _select_historical_worksheet(*, workbook: Any, preferred_sheet_name: str | None) -> Any:
    sheet_names = list(getattr(workbook, "sheetnames", []))
    if not sheet_names:
        raise ValueError("Historical workbook has no worksheets")

    if preferred_sheet_name and preferred_sheet_name in sheet_names:
        return workbook[preferred_sheet_name]

    fallback_sheet_name = sorted(sheet_names, key=str.casefold)[0]
    if preferred_sheet_name:
        LOGGER.warning(
            "historical_sheet_preferred_missing preferred=%s fallback=%s",
            preferred_sheet_name,
            fallback_sheet_name,
        )
    return workbook[fallback_sheet_name]


def _validate_historical_headers(*, worksheet: Any) -> None:
    worksheet_title = str(getattr(worksheet, "title", "<unknown>"))
    header_row = _find_historical_header_row(worksheet=worksheet)
    date_value = _normalize_header(worksheet.cell(row=header_row, column=1).value)
    first_series_value = _normalize_header(worksheet.cell(row=header_row, column=2).value)
    second_series_value = _normalize_header(worksheet.cell(row=header_row, column=3).value)

    missing: list[str] = []
    if date_value not in _DATE_HEADER_CANDIDATES:
        missing.append(_REQUIRED_HISTORICAL_APPEND_HEADERS[0])
    if not first_series_value:
        missing.append(_REQUIRED_HISTORICAL_APPEND_HEADERS[1])
    if not second_series_value:
        missing.append(_REQUIRED_HISTORICAL_APPEND_HEADERS[2])

    if missing:
        raise ValueError(
            "Historical workbook sheet "
            f"{worksheet_title!r} is missing required columns for append: {', '.join(missing)}"
        )


def _find_historical_header_row(*, worksheet: Any, max_scan_rows: int = 25) -> int:
    max_row = int(getattr(worksheet, "max_row", max_scan_rows))
    for row in range(1, min(max_row, max_scan_rows) + 1):
        if _normalize_header(worksheet.cell(row=row, column=1).value) in _DATE_HEADER_CANDIDATES:
            return row
    raise ValueError(
        "Historical workbook sheet "
        f"{getattr(worksheet, 'title', '<unknown>')!r} is missing a date header in column A"
    )


def _sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(64 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _require_path(path: Path | None, *, field_name: str) -> Path:
    if path is None:
        raise ValueError(f"{field_name} is required for pipeline execution")
    return path


def _records(table: Any) -> list[dict[str, Any]]:
    if isinstance(table, list):
        return [dict(record) for record in table if isinstance(record, Mapping)]
    if hasattr(table, "to_dict"):
        as_records = table.to_dict(orient="records")
        return [dict(record) for record in as_records if isinstance(record, Mapping)]
    if hasattr(table, "to_records"):
        as_records = table.to_records()
        return [dict(record) for record in as_records if isinstance(record, Mapping)]
    return []


def _column_names(table: Any) -> set[str]:
    if hasattr(table, "columns"):
        raw_columns = table.columns
        try:
            return {str(column) for column in raw_columns}
        except TypeError:
            return set()

    records = _records(table)
    if not records:
        return set()
    return {str(column) for column in records[0]}


def _row_count(table: Any) -> int:
    return len(_records(table))


def _has_reconciliation_rows(parsed_sections: Any) -> bool:
    if not isinstance(parsed_sections, Mapping):
        return False
    return any(_row_count(table) > 0 for table in parsed_sections.values())


def _run_reconciliation_checks(
    *,
    run_dir: Path,
    config: WorkflowConfig,
    parsed_by_variant: dict[str, dict[str, Any]],
    warnings: list[str],
) -> None:
    variant_historical_paths: dict[str, Path] = {
        "all_programs": config.hist_all_programs_3yr_xlsx,
        "ex_trend": config.hist_ex_llc_3yr_xlsx,
        "trend": config.hist_llc_3yr_xlsx,
    }

    total_gap_count = 0
    impacted_series_count = 0
    impacted_rows_count = 0
    first_unmapped_counterparty_error: UnmappedCounterpartyError | None = None
    reconciliation_by_variant: dict[str, dict[str, Any]] = {}
    for variant, historical_path in variant_historical_paths.items():
        parsed_sections = parsed_by_variant.get(variant, {})
        if not _has_reconciliation_rows(parsed_sections):
            continue
        historical_headers_by_sheet = _extract_historical_series_headers_by_sheet(historical_path)
        parsed_data_by_sheet = _build_parsed_data_by_sheet(
            parsed_sections=parsed_sections,
            historical_series_headers_by_sheet=historical_headers_by_sheet,
        )
        result = reconcile_series_coverage(
            parsed_data_by_sheet=parsed_data_by_sheet,
            historical_series_headers_by_sheet=historical_headers_by_sheet,
            variant=variant,
            expected_segments_by_variant=config.reconciliation.expected_segments_by_variant,
            fail_policy=config.reconciliation.fail_policy,
        )
        reconciliation_by_variant[variant] = result
        total_gap_count += int(result.get("gap_count", 0))
        result_exceptions = result.get("exceptions")
        if first_unmapped_counterparty_error is None and isinstance(result_exceptions, list):
            for exception in result_exceptions:
                if isinstance(exception, UnmappedCounterpartyError):
                    first_unmapped_counterparty_error = exception
                    break
        by_sheet_result = result.get("by_sheet", {})
        if isinstance(by_sheet_result, Mapping):
            for sheet_name, sheet_result in by_sheet_result.items():
                if not isinstance(sheet_result, Mapping):
                    continue
                series_delta, rows_delta = _calculate_impacted_scope_for_sheet(
                    parsed_sections=parsed_data_by_sheet.get(str(sheet_name), {}),
                    reconciliation_sheet_result=sheet_result,
                )
                impacted_series_count += series_delta
                impacted_rows_count += rows_delta
        for warning in result.get("warnings", []):
            warnings.append(f"Reconciliation ({variant}): {warning}")

    if total_gap_count == 0:
        return

    warnings.append(
        "Reconciliation summary: "
        f"gaps={total_gap_count}, impacted_series={impacted_series_count}, "
        f"impacted_rows={impacted_rows_count}"
    )
    _write_needs_mapping_updates(
        run_dir=run_dir,
        fail_policy=config.reconciliation.fail_policy,
        reconciliation_by_variant=reconciliation_by_variant,
        total_gap_count=total_gap_count,
        impacted_series_count=impacted_series_count,
        impacted_rows_count=impacted_rows_count,
    )
    if config.reconciliation.fail_policy == "strict":
        if first_unmapped_counterparty_error is not None:
            raise first_unmapped_counterparty_error
        raise ValueError(
            "Reconciliation strict mode failed due to missing/unmapped series; "
            f"gap_count={total_gap_count}"
        )


def _build_parsed_data_by_sheet(
    *,
    parsed_sections: Mapping[str, Any],
    historical_series_headers_by_sheet: Mapping[str, tuple[str, ...] | list[str] | set[str]],
) -> dict[str, dict[str, list[dict[str, Any]]]]:
    sheet_names = [str(name) for name in historical_series_headers_by_sheet if str(name)]
    if not sheet_names:
        return {}

    parsed_data_by_sheet: dict[str, dict[str, list[dict[str, Any]]]] = {
        sheet_name: {"totals": [], "futures": []} for sheet_name in sheet_names
    }
    normalized_headers_by_sheet: dict[str, set[str]] = {
        sheet_name: {
            normalize_counterparty(header)
            for header in (
                str(value).strip()
                for value in historical_series_headers_by_sheet.get(sheet_name, ())
            )
            if header
        }
        for sheet_name in sheet_names
    }

    totals_fallback_sheet = _select_fallback_sheet_name(
        sheet_names=sheet_names,
        preferred_tokens=("total", "counterparty", "fcm"),
    )
    futures_fallback_sheet = _select_fallback_sheet_name(
        sheet_names=sheet_names,
        preferred_tokens=("future", "fcm", "ch", "clearing"),
        default=totals_fallback_sheet,
    )

    for record in _records(parsed_sections.get("totals", [])):
        raw_counterparty = str(record.get("counterparty", "")).strip()
        normalized_counterparty = (
            normalize_counterparty(raw_counterparty) if raw_counterparty else ""
        )
        target_sheet = _sheet_for_series_label(
            normalized_label=normalized_counterparty,
            sheet_names=sheet_names,
            normalized_headers_by_sheet=normalized_headers_by_sheet,
            fallback_sheet=totals_fallback_sheet,
        )
        if target_sheet:
            parsed_data_by_sheet[target_sheet]["totals"].append(record)

    for record in _records(parsed_sections.get("futures", [])):
        raw_clearing_house = str(record.get("clearing_house", "")).strip()
        normalized_clearing_house = (
            normalize_counterparty(raw_clearing_house) if raw_clearing_house else ""
        )
        target_sheet = _sheet_for_series_label(
            normalized_label=normalized_clearing_house,
            sheet_names=sheet_names,
            normalized_headers_by_sheet=normalized_headers_by_sheet,
            fallback_sheet=futures_fallback_sheet,
        )
        if target_sheet:
            parsed_data_by_sheet[target_sheet]["futures"].append(record)

    return parsed_data_by_sheet


def _calculate_impacted_scope_for_sheet(
    *,
    parsed_sections: Mapping[str, Any],
    reconciliation_sheet_result: Mapping[str, Any],
) -> tuple[int, int]:
    normalized_impacted_series = _identify_impacted_entities_for_sheet(
        reconciliation_sheet_result=reconciliation_sheet_result
    )
    impacted_rows = _count_rows_for_impacted_entities(
        parsed_sections=parsed_sections,
        normalized_impacted_series=normalized_impacted_series,
    )
    impacted_series = len(normalized_impacted_series)
    return impacted_series, impacted_rows


def _identify_impacted_entities_for_sheet(
    *,
    reconciliation_sheet_result: Mapping[str, Any],
) -> set[str]:
    missing_series_labels = {
        str(label).strip()
        for key in (
            "missing_from_historical_headers",
            "missing_from_data",
            "missing_normalized_counterparties",
        )
        for label in reconciliation_sheet_result.get(key, [])
        if str(label).strip()
    }
    normalized_impacted_series = {
        normalize_counterparty(label) for label in missing_series_labels if label
    }
    return normalized_impacted_series


def _count_rows_for_impacted_entities(
    *,
    parsed_sections: Mapping[str, Any],
    normalized_impacted_series: set[str],
) -> int:
    impacted_rows = 0
    for impacted_label in normalized_impacted_series:
        impacted_rows += _count_rows_for_normalized_label(
            parsed_sections=parsed_sections,
            normalized_label=impacted_label,
        )
    return impacted_rows


def _count_rows_for_normalized_label(
    *,
    parsed_sections: Mapping[str, Any],
    normalized_label: str,
) -> int:
    impacted_rows = 0
    for record in _records(parsed_sections.get("totals", [])):
        if (
            _record_normalized_label(record=record, raw_label_key="counterparty")
            == normalized_label
        ):
            impacted_rows += 1

    for record in _records(parsed_sections.get("futures", [])):
        if (
            _record_normalized_label(record=record, raw_label_key="clearing_house")
            == normalized_label
        ):
            impacted_rows += 1
    return impacted_rows


def _record_normalized_label(*, record: Mapping[str, Any], raw_label_key: str) -> str:
    normalized_field = str(record.get("normalized_label", "")).strip()
    if normalized_field:
        return normalize_counterparty(normalized_field)

    raw_label = str(record.get(raw_label_key, "")).strip()
    if not raw_label:
        return ""
    return normalize_counterparty(raw_label)


def _select_fallback_sheet_name(
    *,
    sheet_names: list[str],
    preferred_tokens: tuple[str, ...],
    default: str | None = None,
) -> str | None:
    for sheet_name in sheet_names:
        normalized_sheet_name = sheet_name.casefold()
        if any(token in normalized_sheet_name for token in preferred_tokens):
            return sheet_name
    if default:
        return default
    return sheet_names[0] if sheet_names else None


def _sheet_for_series_label(
    *,
    normalized_label: str,
    sheet_names: list[str],
    normalized_headers_by_sheet: Mapping[str, set[str]],
    fallback_sheet: str | None,
) -> str | None:
    if normalized_label:
        for sheet_name in sheet_names:
            if normalized_label in normalized_headers_by_sheet.get(sheet_name, set()):
                return sheet_name
    return fallback_sheet


def _extract_historical_series_headers_by_sheet(workbook_path: Path) -> dict[str, tuple[str, ...]]:
    def _raise_with_context(*, exc: Exception, context: str) -> NoReturn:
        message = str(exc)
        suffix = f": {message}" if message else ""
        raise type(exc)(f"{context}{suffix}") from exc

    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError(
            "Reconciliation requires openpyxl to read historical workbook headers"
        ) from exc

    try:
        workbook = load_workbook(filename=workbook_path, read_only=True, data_only=True)
    except (KeyError, ValueError, OSError) as exc:
        raise RuntimeError(
            f"Unable to extract historical series headers from workbook: {workbook_path}"
        ) from exc
    except Exception as exc:
        _raise_with_context(
            exc=exc,
            context=f"Unexpected error while loading historical workbook {workbook_path!s}",
        )
    workbook_obj: Any = workbook
    try:
        headers_by_sheet: dict[str, tuple[str, ...]] = {}
        for sheet_name in getattr(workbook_obj, "sheetnames", []):
            try:
                worksheet = workbook_obj[sheet_name]
            except KeyError:
                continue
            except Exception as exc:
                _raise_with_context(
                    exc=exc,
                    context=(
                        "Unexpected error while loading historical worksheet in "
                        f"workbook {workbook_path!s}, sheet {sheet_name!r}"
                    ),
                )

            try:
                header_row = _find_historical_header_row(worksheet=worksheet)
            except ValueError:
                continue
            except Exception as exc:
                _raise_with_context(
                    exc=exc,
                    context=(
                        "Unexpected error while scanning historical header row in "
                        f"workbook {workbook_path!s}, sheet {sheet_name!r}"
                    ),
                )

            try:
                max_column = int(getattr(worksheet, "max_column", 0))
                headers = [
                    label
                    for label in (
                        str(worksheet.cell(row=header_row, column=column_index).value or "").strip()
                        for column_index in range(2, max_column + 1)
                    )
                    if label
                ]
            except (KeyError, ValueError):
                continue
            except Exception as exc:
                _raise_with_context(
                    exc=exc,
                    context=(
                        "Unexpected error while reading historical series headers in "
                        f"workbook {workbook_path!s}, sheet {sheet_name!r}"
                    ),
                )

            headers_by_sheet[sheet_name] = tuple(headers)
        return headers_by_sheet
    finally:
        workbook_obj.close()


def _write_needs_mapping_updates(
    *,
    run_dir: Path,
    fail_policy: str,
    reconciliation_by_variant: Mapping[str, Mapping[str, Any]],
    total_gap_count: int,
    impacted_series_count: int,
    impacted_rows_count: int,
) -> Path:
    timestamp = utc_now_isoformat()
    lines: list[str] = [
        "Counter Risk Reconciliation Gaps",
        f"timestamp_utc: {timestamp}",
        f"run_identifier: {run_dir.name}",
        f"fail_policy: {fail_policy}",
        f"total_gap_count: {total_gap_count}",
        f"impacted_series_count: {impacted_series_count}",
        f"impacted_rows_count: {impacted_rows_count}",
        "",
    ]

    for variant in sorted(reconciliation_by_variant, key=str.casefold):
        payload = reconciliation_by_variant[variant]
        lines.append(f"[variant: {variant}]")

        missing_series = payload.get("missing_series", [])
        if isinstance(missing_series, list):
            for item in missing_series:
                if not isinstance(item, Mapping):
                    continue
                sheet_name = str(item.get("sheet", "unknown"))
                missing_labels = item.get("missing_from_historical_headers", [])
                if isinstance(missing_labels, list) and missing_labels:
                    lines.append(
                        f"- sheet {sheet_name}: missing_from_historical_headers="
                        f"{', '.join(str(label) for label in missing_labels)}"
                    )

        missing_segments = payload.get("missing_segments", [])
        if isinstance(missing_segments, list):
            for item in missing_segments:
                if not isinstance(item, Mapping):
                    continue
                sheet_name = str(item.get("sheet", "unknown"))
                expected_labels = item.get("expected_segment_identifiers", [])
                if isinstance(expected_labels, list) and expected_labels:
                    lines.append(
                        f"- sheet {sheet_name}: missing_expected_segments="
                        f"{', '.join(str(label) for label in expected_labels)}"
                    )

        lines.append("")

    output_path = run_dir / "NEEDS_MAPPING_UPDATES.txt"
    output_path.write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")
    return output_path
